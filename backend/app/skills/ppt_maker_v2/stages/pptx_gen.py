"""Stage: PPTX generation — merge single pages into a .pptx file."""

import os
import re
import uuid

from app.skills.base import SkillContext, SkillResult
from app.services._paths import PUBLIC_DIR
from ..constants import SLIDE_WIDTH_INCHES, SLIDE_HEIGHT_INCHES


async def handle(context: SkillContext, session: dict, sessions: dict, session_id: str) -> SkillResult:
    """Build PPTX from selected page images."""
    page_images = session.get("single_pages", [])
    if not page_images:
        return SkillResult(
            success=True,
            message="暂无可生成 PPTX 的页面。请先完成第 3 步逐页生成。",
            data={"skill": "ppt_maker", "stage": "pptx_scope"},
        )

    msg = context.user_message.strip()
    selected = _select_pages(msg, page_images)

    try:
        result = _build_pptx(session_id, selected)
    except Exception as exc:
        return SkillResult(success=False, message=f"生成 PPTX 失败：{exc}")

    sessions[session_id]["stage"] = "completed"
    return SkillResult(
        success=True,
        message=(
            f"✅ PPTX 已生成！共 {len(selected)} 页。\n\n"
            f"[下载 PPTX](/api/skills/download/{result['filename']})\n\n"
            f"如需继续制作新 PPT，请重新触发技能。"
        ),
        data={"skill": "ppt_maker", "download_url": f"/api/skills/download/{result['filename']}", "stage": "completed"},
    )


def _select_pages(request: str, page_images: list[dict]) -> list[dict]:
    """Parse user's page selection request."""
    r = request.strip().lower().replace(" ", "")

    # "全部" / "all"
    if r in {"全部", "all", "整套", "所有", "都有"}:
        return page_images

    # Range: "1-3" / "1~3" / "1至3"
    range_match = re.match(r"(\d+)[-~至](\d+)", r)
    if range_match:
        start, end = int(range_match.group(1)), int(range_match.group(2))
        return [p for p in page_images if start <= p["index"] <= end]

    # Individual numbers
    nums = set()
    for m in re.finditer(r"\d+", r):
        nums.add(int(m.group()))
    if nums:
        return [p for p in page_images if p["index"] in nums]

    # Default: all
    return page_images


def _build_pptx(session_id: str, page_images: list[dict]) -> dict:
    """Build a 16:9 PPTX from page images."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_INCHES)
    prs.slide_height = Inches(SLIDE_HEIGHT_INCHES)

    for img in page_images:
        path = img.get("path", "")
        if not path or not os.path.exists(path):
            continue
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        slide.shapes.add_picture(
            path, 0, 0,
            width=prs.slide_width, height=prs.slide_height,
        )

    output_dir = str(PUBLIC_DIR)
    os.makedirs(output_dir, exist_ok=True)
    filename = f"ppt_maker_{session_id[:8]}_{uuid.uuid4().hex[:10]}.pptx"
    filepath = os.path.join(output_dir, filename)
    prs.save(filepath)

    return {
        "filename": filename,
        "path": filepath,
        "download_url": f"/api/skills/download/{filename}",
    }
