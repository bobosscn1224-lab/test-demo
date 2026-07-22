from __future__ import annotations

import asyncio
import logging
import os
import re
import uuid

from app.config import settings
from app.services._paths import PUBLIC_DIR
from app.services.llm_service import llm_service
from app.skills.base import BaseSkill, SkillContext, SkillResult
from app.utils.file_parser import parse_file_sync
from app.core.skill_session import SkillSessionHelper
from . import visual_systems

logger = logging.getLogger(__name__)

SKILL_NAME = "ppt_maker"
_sessions: dict[str, dict] = {}
_helper = SkillSessionHelper(SKILL_NAME, _sessions)

EXIT_WORDS = {
    "退出",
    "返回",
    "不做了",
    "停止",
    "结束",
    "取消",
    "先不做",
    "退出skill",
    "退出 skill",
}

CONFIRM_WORDS = {"确认", "可以", "没问题", "通过", "就这样", "ok", "okay", "yes", "确定", "开始", "继续"}


class PPTMakerSkill(BaseSkill):
    name = "ppt_maker"
    description = "根据用户输入、上传文档、PPT大纲或视觉稿制作PPT，支持大纲、缩略图、逐页高清图和PPTX输出。"
    triggers = ["做PPT", "制作PPT", "生成PPT", "帮我做PPT", "PPT", "ppt", "幻灯片", "演示文稿"]
    keywords = ["ppt", "powerpoint", "slides", "deck", "提案", "路演", "汇报材料", "课件"]

    async def execute(self, context: SkillContext) -> SkillResult:
        msg = context.user_message.strip()
        session_id = context.session_id or "default"

        # Restore persisted sessions on first access
        await _helper.restore()

        if self._is_exit(msg):
            await _helper.delete(session_id)
            return SkillResult(success=True, message="好的，已退出当前 PPT 制作流程。")

        session = _sessions.get(session_id)
        if not session:
            return self._show_entry_menu(session_id)

        stage = session.get("stage")
        result = None
        if stage == "awaiting_entry_choice":
            result = await self._handle_entry_choice(context, session_id)
        elif stage == "awaiting_content":
            result = await self._handle_content(context, session_id)
        elif stage == "awaiting_outline_confirm":
            result = await self._handle_outline_confirm(context, session_id)
        elif stage == "awaiting_outline_for_visual":
            result = await self._handle_outline_for_visual(context, session_id)
        elif stage == "awaiting_collage_for_pages":
            result = await self._handle_collage_for_pages(context, session_id)
        elif stage == "awaiting_page_info_for_step3":
            result = await self._handle_page_info_for_step3(context, session_id)
        elif stage == "awaiting_page_image_for_editable_ppt":
            result = await self._handle_page_image_for_editable_ppt(context, session_id)
        elif stage == "visual_direction":
            result = await self._handle_visual_direction_confirm(context, session_id)
        elif stage == "entry3_ready":
            result = await self._handle_entry3_ready(context, session_id)
        elif stage == "awaiting_visual_choice":
            result = self._handle_visual_choice(context, session_id)
        elif stage == "awaiting_step3_start":
            result = self._ack_step3_start(session_id)
        elif stage in {"generating_visual_direction", "generating_single_pages"}:
            result = self._return_generation_status(session_id)
        elif stage == "awaiting_editable_ppt_scope":
            result = self._handle_editable_ppt_scope(context, session_id)
        else:
            result = self._show_entry_menu(session_id)

        # Persist session to DB after every state change
        current = _sessions.get(session_id)
        if current:
            await _helper.save(session_id, current)

        return result

    async def execute_stream(self, context: SkillContext):
        session_id = context.session_id or "default"
        session = _sessions.get(session_id)

        if session and session.get("stage") == "awaiting_outline_for_visual":
            source_text = await self._collect_source_text(context)
            if not self._has_enough_content(context.user_message, source_text):
                yield await self._handle_outline_for_visual(context, session_id)
                return
            _sessions[session_id] = {**session, "stage": "visual_direction", "outline": source_text}
            yield "收到 PPT 大纲，开始从第 2 步生成三版视觉缩略图。\n\n"
            async for item in self._run_step2_visual_collages(session_id, source_text):
                yield item
            return

        if session and session.get("stage") == "visual_direction" and self._is_confirm(context.user_message):
            outline = str(session.get("outline") or "").strip()
            if not outline:
                yield SkillResult(success=False, message="第二步出错了：没有找到已确认的大纲内容。请重新提供大纲后再继续。")
                return
            yield "收到，开始从第 2 步生成三版 PPT 视觉缩略图。\n\n"
            async for item in self._run_step2_visual_collages(session_id, outline):
                yield item
            return

        if session and session.get("stage") == "entry3_ready" and self._is_confirm(context.user_message):
            outline = str(session.get("outline") or "").strip()
            if not outline:
                yield SkillResult(success=False, message="第 3 步出错了：没有找到页面信息。请重新提供页数和标题。")
                return
            yield "收到，跳过第 2 步，直接进入第 3 步：基于上传的缩略图逐页生成高清单页视觉稿。\n\n"
            await asyncio.sleep(0.1)
            async for item in self._start_step3_single_pages(session_id, "REF"):
                yield item
            return

        if session and session.get("stage") == "awaiting_outline_confirm" and self._is_confirm(context.user_message):
            outline = str(session.get("outline") or "").strip()
            if not outline:
                yield SkillResult(
                    success=False,
                    message="第二步出错了：没有找到已确认的大纲内容。请重新上传资料或重新生成大纲后再确认。",
                    data={"skill": self.name, "stage": "error"},
                )
                return
            yield "收到，大纲已确认。\n\n"
            async for item in self._run_step2_visual_collages(session_id, outline):
                yield item
            return

        if session and session.get("stage") == "awaiting_visual_choice":
            choice = self._parse_visual_choice(context.user_message)
            if choice:
                yield f"已接收到你的选择：方案 {choice}。\n"
                yield "正在进入第 3 步：逐页生成高清 16:9 单页 PPT 视觉稿。\n\n"
                await asyncio.sleep(0.1)
                async for item in self._start_step3_single_pages(session_id, choice):
                    yield item
                return
            yield self._handle_visual_choice(context, session_id)
            return

        if session and session.get("stage") == "awaiting_step3_start":
            choice = self._selected_visual_choice(session)
            if choice:
                yield f"已接收到第 3 步启动指令，继续使用方案 {choice}。\n"
                yield "正在进入第 3 步：逐页生成高清 16:9 单页 PPT 视觉稿。\n\n"
                await asyncio.sleep(0.1)
                async for item in self._start_step3_single_pages(session_id, choice):
                    yield item
                return
            yield SkillResult(
                success=False,
                message="第 3 步无法启动：没有找到已选择的视觉方案。请重新选择方案 A、B 或 C。",
                data={"skill": self.name, "stage": "awaiting_visual_choice"},
            )
            return

        result = await self.execute(context)
        yield result

    def _is_exit(self, msg: str) -> bool:
        normalized = msg.strip().lower().replace(" ", "")
        return len(normalized) <= 12 and normalized in {w.lower().replace(" ", "") for w in EXIT_WORDS}

    def _is_confirm(self, msg: str) -> bool:
        normalized = msg.strip().lower()
        return len(normalized) <= 12 and any(word in normalized for word in CONFIRM_WORDS)

    def _show_entry_menu(self, session_id: str) -> SkillResult:
        _sessions[session_id] = {"stage": "awaiting_entry_choice"}
        return SkillResult(
            success=True,
            message=(
                "请选择本次 PPT 制作的入口：\n\n"
                "1. 上传文档或直接输入内容，生成大纲制作 PPT（走第 1-2-3-4 步）\n"
                "2. 直接提供 PPT 大纲，开始制作 PPT 缩略图（从第 2 步开始）\n"
                "3. 直接提供 PPT 整体详细缩略图，开始输出分页高清图片风格图（从第 3 步开始）\n"
                "4. 提供某一页高清 PPT 风格图，制作可编辑 PPT（直接第 4 步）\n\n"
                "请回复数字 1、2、3 或 4。"
            ),
            data={"skill": self.name, "stage": "awaiting_entry_choice"},
        )

    def _parse_entry_choice(self, msg: str) -> str | None:
        normalized = msg.strip().lower()
        compact = normalized.replace(" ", "")
        mapping = {
            "1": "1", "选1": "1", "选择1": "1", "第一项": "1", "第1项": "1",
            "2": "2", "选2": "2", "选择2": "2", "第二项": "2", "第2项": "2",
            "3": "3", "选3": "3", "选择3": "3", "第三项": "3", "第3项": "3",
            "4": "4", "选4": "4", "选择4": "4", "第四项": "4", "第4项": "4",
        }
        if compact in mapping:
            return mapping[compact]
        if "生成大纲" in normalized or "上传文档" in normalized or "直接输入内容" in normalized:
            return "1"
        if "ppt大纲" in compact or "ppt的大纲" in compact:
            return "2"
        if "整体详细缩略图" in normalized or "分页高清" in normalized or "高清图片风格图" in normalized:
            return "3"
        if "可编辑ppt" in compact or "可编辑pptx" in compact or "某一页高清" in normalized:
            return "4"
        return None

    async def _handle_entry_choice(self, context: SkillContext, session_id: str) -> SkillResult:
        choice = self._parse_entry_choice(context.user_message)
        if not choice:
            return SkillResult(success=True, message="请回复数字 1、2、3 或 4，选择本次 PPT 制作入口。")

        if choice == "1":
            _sessions[session_id] = {"stage": "awaiting_content", "entry_choice": "1"}
            return SkillResult(
                success=True,
                message="已选择入口 1。请上传用于制作 PPT 的文档，或直接输入内容、主题、受众和制作要求。",
                data={"skill": self.name, "stage": "awaiting_content"},
            )
        if choice == "2":
            _sessions[session_id] = {"stage": "awaiting_outline_for_visual", "entry_choice": "2"}
            return SkillResult(
                success=True,
                message="已选择入口 2。请粘贴已确认的 PPT 大纲和逐页内容，或上传包含大纲的文档。",
                data={"skill": self.name, "stage": "awaiting_outline_for_visual"},
            )
        if choice == "3":
            _sessions[session_id] = {"stage": "awaiting_collage_for_pages", "entry_choice": "3"}
            return SkillResult(
                success=True,
                message="已选择入口 3。请上传 PPT 整体详细缩略图，或输入每页缩略图的详细描述、页数和页序。",
                data={"skill": self.name, "stage": "awaiting_collage_for_pages"},
            )

        _sessions[session_id] = {"stage": "awaiting_page_image_for_editable_ppt", "entry_choice": "4"}
        return SkillResult(
            success=True,
            message="已选择入口 4。请上传这一页高清 PPT 风格图，并说明需要还原成单页 PPTX 还是加入现有 PPT。",
            data={"skill": self.name, "stage": "awaiting_page_image_for_editable_ppt"},
        )

    async def _handle_content(self, context: SkillContext, session_id: str) -> SkillResult:
        source_text = await self._collect_source_text(context)
        if not self._has_enough_content(context.user_message, source_text):
            return SkillResult(
                success=True,
                message="我还没有收到足够的 PPT 资料。请上传文档，或直接输入主题、内容、受众和制作要求。",
                data={"skill": self.name, "stage": "awaiting_content"},
            )
        return await self._generate_outline(context, session_id, source_text)

    async def _handle_outline_confirm(self, context: SkillContext, session_id: str) -> SkillResult:
        if self._is_confirm(context.user_message):
            return SkillResult(
                success=True,
                message="收到确认。请回复“开始”，我将进入第 2 步生成三版 PPT 缩略图。",
                data={"skill": self.name, "stage": "awaiting_outline_confirm"},
            )
        source_text = await self._collect_source_text(context)
        if source_text or len(context.user_message.strip()) > 10:
            return await self._generate_outline(context, session_id, source_text, revision=context.user_message)
        return SkillResult(success=True, message="请确认大纲是否可以进入下一步；如果需要调整，请直接告诉我。")

    async def _handle_outline_for_visual(self, context: SkillContext, session_id: str) -> SkillResult:
        source_text = await self._collect_source_text(context)
        if not self._has_enough_content(context.user_message, source_text):
            return SkillResult(
                success=True,
                message="请提供已确认的 PPT 大纲和逐页内容，或上传包含大纲的文档。",
                data={"skill": self.name, "stage": "awaiting_outline_for_visual"},
            )
        _sessions[session_id] = {"stage": "visual_direction", "entry_choice": "2", "outline": source_text}
        return SkillResult(success=True, message="已收到 PPT 大纲。请回复“开始”，我将从第 2 步生成三版 PPT 缩略图。")

    async def _handle_collage_for_pages(self, context: SkillContext, session_id: str) -> SkillResult:
        if not context.uploaded_files and len(context.user_message.strip()) <= 20:
            return SkillResult(
                success=True,
                message=(
                    "请上传 PPT 整体详细缩略图（图片文件），并同时用文字补充以下信息：\n"
                    "1. 总页数\n"
                    "2. 每页标题（按页序排列）\n\n"
                    "这些文字信息是为了确保生成时页码和内容准确对应，缩略图本身会作为视觉风格参考。"
                ),
            )
        _sessions[session_id] = {
            "stage": "awaiting_page_info_for_step3",
            "entry_choice": "3",
            "collage_files": context.uploaded_files,
        }
        return SkillResult(
            success=True,
            message=(
                "已收到缩略图。请补充以下文字信息，我会进入第 3 步逐页生成高清风格图：\n"
                "1. 总页数（如：12 页）\n"
                "2. 每页标题（按第 1 页到最后一页的顺序列出）\n\n"
                "例如：\n"
                "第 1 页：封面 - Q2 业务复盘\n"
                "第 2 页：核心指标总览\n"
                "第 3 页：MO 漏斗转化分析\n"
                "...\n\n"
                "这些信息用于确保逐页生成时页码和标题准确对位。"
            ),
            data={"skill": self.name, "stage": "awaiting_page_info_for_step3"},
        )

    async def _handle_page_info_for_step3(self, context: SkillContext, session_id: str) -> SkillResult:
        """Handle user providing page count and titles for step 3."""
        msg = context.user_message.strip()
        if len(msg) <= 10:
            return SkillResult(
                success=True,
                message="请提供总页数和每页标题。例如：\n总 10 页\n第 1 页：封面\n第 2 页：目录\n...",
                data={"skill": self.name, "stage": "awaiting_page_info_for_step3"},
            )

        # Extract slide info from user input
        page_count = None
        m = re.search(r'(\d+)\s*页', msg)
        if m:
            page_count = int(m.group(1))

        # Extract per-page titles
        titles = []
        for m in re.finditer(r'第\s*(\d+)\s*页[：:]\s*(.+?)(?=第\s*\d+\s*页|$)', msg):
            titles.append({"page": int(m.group(1)), "title": m.group(2).strip()})

        if not titles:
            # Try line-by-line parsing
            lines = [l.strip() for l in msg.split('\n') if l.strip()]
            for i, line in enumerate(lines, 1):
                clean = re.sub(r'^\d+[\.\)、]?\s*', '', line).strip()
                if clean:
                    titles.append({"page": i, "title": clean})

        if not titles or len(titles) < 2:
            return SkillResult(
                success=True,
                message="未能识别出足够的页面信息。请按格式列出：\n第 1 页：[标题]\n第 2 页：[标题]\n...",
                data={"skill": self.name, "stage": "awaiting_page_info_for_step3"},
            )

        actual_count = page_count or len(titles)
        if actual_count != len(titles) and page_count:
            actual_count = page_count

        # Build outline from the provided info
        outline = f"PPT 共 {actual_count} 页\n\n"
        for t in titles:
            outline += f"第 {t['page']} 页：{t['title']}\n"

        session = _sessions.get(session_id, {})
        collage_files = session.get("collage_files") or []

        # For entry 3: the uploaded collage IS the visual master — skip Step 2, go direct to Step 3
        _sessions[session_id] = {
            **session,
            "stage": "entry3_ready",
            "entry_choice": "3",
            "outline": outline,
            "selected_visual": {"label": "REF", "filename": collage_files[0].get("filename", "") if collage_files else "uploaded_collage", "path": collage_files[0].get("path", "") if collage_files else ""},
        }
        collages: list[dict] = [{"label": "REF", "filename": collage_files[0].get("filename", "") if collage_files else "", "path": collage_files[0].get("path", "") if collage_files else ""}]
        _sessions[session_id]["visual_collages"] = collages

        return SkillResult(
            success=True,
            message=(
                f"已识别 {len(titles)} 页标题，上传的缩略图将直接作为视觉母版。\n\n"
                f"请回复「开始」，我将跳过第 2 步（三版缩略图），直接进入第 3 步逐页生成高清单页视觉稿。"
            ),
            data={"skill": self.name, "stage": "entry3_ready"},
        )

    async def _handle_page_image_for_editable_ppt(self, context: SkillContext, session_id: str) -> SkillResult:
        session = _sessions.get(session_id, {})
        msg = context.user_message.strip().lower()

        # 1. User uploads images → store and show conversion options
        if context.uploaded_files:
            stored = list(session.get("uploaded_page_images", []))
            for f in context.uploaded_files:
                path = f.get("path", "")
                if path and path not in stored:
                    stored.append(path)
            _sessions[session_id] = {**session, "uploaded_page_images": stored, "stage": "awaiting_page_image_for_editable_ppt"}
            return self._show_conversion_options(session_id, stored)

        # 2. Get stored images
        stored = session.get("uploaded_page_images", [])

        # 3. User picks conversion mode
        # PRECISE / VERIFIED / OBJECT — local object-first reconstruction.
        if msg in ("precise", "verified", "verify", "object", "objects", "本地对象", "对象化", "精确重建", "高仿", "高仿真", "验证模式"):
            if not stored:
                return self._ask_upload_or_format()
            return await self._generate_precise_pptx(session_id, stored)

        # LAYOUT — fast layout-aware
        if msg in ("layout", "布局"):
            if not stored:
                return self._ask_upload_or_format()
            return await self._generate_layout_pptx(session_id, stored)

        # CODEX
        if msg in ("codex", "ai生成", "ai视觉"):
            if not stored:
                return self._ask_upload_or_format()
            return await self._generate_codex_pptx(session_id, stored)

        # AGNES — AI layout understanding
        if msg in ("agnes", "ai布局", "智能布局"):
            if not stored:
                return self._ask_upload_or_format()
            return await self._generate_agnes_pptx(session_id, stored)

        # DECKWEAVER — full deckweaver pipeline (ENABLE_NATIVE_OUTLINE_SHAPES=True)
        if msg in ("deckweaver", "dw", "dw布局", "原生布局"):
            if not stored:
                return self._ask_upload_or_format()
            return await self._generate_deckweaver_pptx(session_id, stored)

        # BATCH — process all uploaded images into a single combined PPTX
        if msg in ("batch", "批量", "批量转换", "合成", "合并"):
            if not stored:
                return self._ask_upload_or_format()
            if len(stored) < 2:
                return SkillResult(success=False, message="批量转换需要至少上传 2 张图片。请继续上传。")
            return await self._generate_batch_pptx(session_id, stored)

        # VBA
        if msg in ("vba", "宏代码"):
            if not stored:
                return self._ask_upload_or_format()
            return await self._generate_vba(session_id, stored)

        # SVG / DRAWIO
        if msg in ("svg",):
            return await self._generate_svg_from_stored(session_id, session)
        if msg in ("drawio", "draw.io"):
            return await self._generate_drawio_from_stored(session_id, session)

        # Generic confirm
        if self._is_confirm(context.user_message):
            if not stored:
                return self._ask_upload_or_format()
            # Default to PRECISE for best quality
            return await self._generate_precise_pptx(session_id, stored)

        # Exit
        if msg in ("退出", "exit", "quit"):
            _sessions[session_id] = {**session, "stage": None}
            return SkillResult(success=True, message="已退出图片转 PPTX。你可以继续其他操作。")

        # Fallback: ask to upload or pick a mode
        if not stored:
            return self._ask_upload_or_format()
        return self._show_conversion_options(session_id, stored)

    # ── Helper methods ──

    def _show_conversion_options(self, session_id: str, stored: list[str]) -> SkillResult:
        lines = [
            f"已收到 {len(stored)} 张图片。请选择转换方式：",
            "",
            "**🥇 OBJECT / VERIFIED** — 本地对象化精确重建（推荐，较慢）",
            "- OCR + 文字擦除 + 对象提取 + 原生形状 + 可编辑文本",
            "- 启用 PPT 渲染预览和候选字号校准，不依赖大模型",
            "",
            "**🥈 LAYOUT** — 快速布局（~30s）",
            "- DeckWeaver 架构，文字可编辑",
            "",
            "**🥉 CODEX** — AI 高仿真（需 GPT-5.5）",
            "",
            "**其他：VBA / SVG / DRAWIO / DECKWEAVER**",
            "",
            "回复 **object** / **verified** / **precise** 使用高质量本地管线；回复 **layout** / **deckweaver** 使用不同架构",
            "也支持 **codex** / **agnes** / **vba** / **svg** / **drawio**",
            "也可继续上传图片。输入 **退出** 结束。",
        ]
        return SkillResult(success=True, message="\n".join(lines),
                          data={"skill": self.name, "stage": "awaiting_page_image_for_editable_ppt"})

    def _ask_upload_or_format(self) -> SkillResult:
        return SkillResult(
            success=True,
            message="请上传高清 PPT 视觉稿图片，然后回复 **object** / **verified** / **precise** / **layout** / **codex** / **vba** 选择转换方式。")

    async def _generate_precise_pptx(self, session_id: str, image_paths: list[str]) -> SkillResult:
        """PRECISE pipeline: full reconstruction with COM calibration."""
        if not image_paths:
            return SkillResult(success=False, message="请先上传图片。")
        from app.services.local_object_reconstruction import PIPELINE_VERSION, reconstruct

        output_dir = os.path.normpath(os.path.join(
            os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))),
            "outputs"))
        results, errors = [], []
        for path in image_paths:
            if not os.path.exists(path):
                errors.append(f"{os.path.basename(path)}：文件不存在"); continue
            try:
                r = await reconstruct(
                    path,
                    session_id=session_id,
                    output_dir=output_dir,
                    render_preview=True,
                    max_calibration_passes=3,
                )
            except Exception as e:
                logger.exception("PRECISE failed"); r = {"error": str(e)}
            (results if not r.get("error") else errors).append(r if not r.get("error") else None)
            if r.get("error"):
                errors.append(f"{os.path.basename(path)}：{r['error']}")

        if not results:
            return SkillResult(success=False, message="PRECISE 重建失败：\n- " + "\n- ".join(filter(None, errors)))

        base_url = f"http://localhost:{settings.port}"
        lines = [f"## PRECISE 精确重建完成 ({PIPELINE_VERSION})", ""]
        for idx, r in enumerate(results):
            dl = f"{base_url}/api/skills/download/{r['filename']}"
            lines.append(f"### 第 {idx+1} 页" if len(results) > 1 else "### 结果")
            lines.append(f"📥 [点击下载 PPTX]({dl})")
            lines.append(f"📂 `{r.get('path', '')}`")
            rep = r.get("report", {})
            ed = rep.get("editable", {})
            lines.append(f"- 文本框 {ed.get('text_boxes',0)} 个 | 形状 {ed.get('native_shapes',0)} 个 | 图片素材 {rep.get('image_assets',{}).get('visual_tiles',0)+rep.get('image_assets',{}).get('movable_objects',0)} 个")
            qa = rep.get("visual_qa", {})
            lines.append(f"- SSIM {qa.get('ssim','N/A')} | 文字SSIM {qa.get('text_ssim','N/A')} | 字号缩放 {qa.get('selected_font_scale',1.0)}")
            if r.get("preview_url"):
                lines.append(f"🖼 [预览图]({base_url}{r['preview_url']})")
            if r.get("comparison_url"):
                lines.append(f"🔍 [对照图]({base_url}{r['comparison_url']})")
            lines.append("")
        if errors:
            lines.append("### 失败\n" + "\n".join(f"- {e}" for e in errors if e))
        return SkillResult(success=True, message="\n".join(lines),
                          data={"skill": self.name, "stage": "completed", "download_url": f"{base_url}/api/skills/download/{results[0]['filename']}"})

    async def _generate_layout_pptx(self, session_id: str, image_paths: list[str]) -> SkillResult:
        """LAYOUT pipeline: DeckWeaver v5 fast reconstruction."""
        if not image_paths: return SkillResult(success=False, message="请先上传图片。")
        from app.services.layout_reconstructor import reconstruct as run_layout
        output_dir = os.path.normpath(os.path.join(
            os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))),
            "outputs"))
        r = await run_layout(image_paths[0], output_dir=output_dir, session_id=session_id, enable_shapes=True)
        if not r: return SkillResult(success=False, message="Layout 重建失败。")
        base_url = f"http://localhost:{settings.port}"
        dl = f"{base_url}/api/skills/download/{r['filename']}"
        msg = (f"## Layout 重建完成\n\n📥 [下载 PPTX]({dl})\n📂 `{r.get('path','')}`\n"
               f"- 文字 {r.get('text_items',0)} | 形状 {r.get('native_shapes',0)} | 图片 {r.get('image_items',0)}\n"
               f"- 布局 {r.get('layout_type','N/A')} | 配色 {', '.join(r.get('color_scheme',{}).get('palette',['N/A'])[:3])}")
        return SkillResult(success=True, message=msg, data={"skill": self.name, "download_url": dl, "stage": "completed"})

    async def _generate_codex_pptx(self, session_id: str, image_paths: list[str]) -> SkillResult:
        if not image_paths: return SkillResult(success=False, message="请先上传图片。")
        from app.services.codex_pptx_service import generate_pptx
        r = await generate_pptx(image_paths[0], session_id=session_id, output_dir=os.path.join("data","outputs"), timeout=300)
        if not r or "error" in r:
            return SkillResult(success=False, message=f"Codex 生成失败：{r.get('error','未知错误') if r else '生成失败'}")
        base_url = f"http://localhost:{settings.port}"
        dl = f"{base_url}/api/skills/download/{r.get('filename','output.pptx')}"
        return SkillResult(success=True, message=f"## Codex AI PPTX\n\n📥 [下载]({dl})\n- {r.get('size',0)//1024} KB", data={"skill": self.name, "download_url": dl, "stage": "completed"})

    async def _generate_agnes_pptx(self, session_id: str, image_paths: list[str]) -> SkillResult:
        if not image_paths: return SkillResult(success=False, message="请先上传图片。")
        if not settings.agnes_api_key:
            return SkillResult(success=False, message="未配置 Agnes API Key。")
        from app.services.agnes_pipeline import PIPELINE_VERSION, reconstruct
        output_dir = os.path.normpath(os.path.join(
            os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))),
            "outputs"))
        r = await reconstruct(image_paths[0], session_id=session_id, output_dir=output_dir,
                              agnes_api_key=s.agnes_api_key, render_preview=True, max_calibration_passes=2)
        if r.get("error"): return SkillResult(success=False, message=f"Agnes 失败：{r['error']}")
        base_url = f"http://localhost:{settings.port}"
        dl = f"{base_url}/api/skills/download/{r['filename']}"
        rep = r.get("report", {})
        ag = rep.get("agnes", {})
        msg = (f"## Agnes 智能布局 ({PIPELINE_VERSION})\n\n📥 [下载 PPTX]({dl})\n📂 `{r.get('path','')}`\n"
               f"- 布局 {ag.get('layout_type','?')} | 风格 {ag.get('visual_style','?')} | {ag.get('element_count',0)} 元素\n"
               f"- 文本框 {rep.get('editable',{}).get('text_boxes',0)} | 形状 {rep.get('editable',{}).get('native_shapes',0)}")
        qa = rep.get("visual_qa", {})
        if qa.get("ssim"):
            msg += f"\n- SSIM {qa['ssim']} | 文字SSIM {qa.get('text_ssim','N/A')}"
        return SkillResult(success=True, message=msg, data={"skill": self.name, "download_url": dl, "stage": "completed"})

    async def _generate_deckweaver_pptx(self, session_id: str, image_paths: list[str]) -> SkillResult:
        """DeckWeaver pipeline with ENABLE_NATIVE_OUTLINE_SHAPES=True."""
        if not image_paths: return SkillResult(success=False, message="请先上传图片。")
        from app.services.deckweaver_service import PIPELINE_VERSION, convert_image_to_pptx
        output_dir = os.path.normpath(os.path.join(
            os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))),
            "outputs"))
        r = await convert_image_to_pptx(image_paths[0], session_id=session_id, output_dir=output_dir)
        if r.get("error"): return SkillResult(success=False, message=f"DeckWeaver 失败：{r['error']}")
        base_url = f"http://localhost:{settings.port}"
        dl = f"{base_url}/api/skills/download/{r['filename']}"
        rep = r.get("report", {})
        msg = (f"## DeckWeaver 原生布局 ({PIPELINE_VERSION})\n\n"
               f"📥 [下载 PPTX]({dl})\n📂 `{r.get('path','')}`\n"
               f"- 文字 {rep.get('text_items', '?')} | 原生形状已启用")
        return SkillResult(success=True, message=msg, data={"skill": self.name, "download_url": dl, "stage": "completed"})

    async def _generate_batch_pptx(self, session_id: str, image_paths: list[str]) -> SkillResult:
        """Batch convert all uploaded images into one combined PPTX."""
        if not image_paths: return SkillResult(success=False, message="请先上传图片。")
        from app.services.batch_pptx_service import PIPELINE_VERSION, batch_convert
        output_dir = os.path.normpath(os.path.join(
            os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))),
            "outputs"))

        # Sort by filename for consistent ordering
        sorted_paths = sorted(image_paths, key=lambda p: os.path.basename(p))
        r = await batch_convert(sorted_paths, session_id=session_id, output_dir=output_dir)

        if r.get("error") and not r.get("page_count"):
            errs = "\n".join(f"- {e['page']}: {e['error']}" for e in r.get("errors", [])[:5])
            return SkillResult(success=False, message=f"批量转换失败：\n{errs}")

        base_url = f"http://localhost:{settings.port}"
        dl = f"{base_url}/api/skills/download/{r['filename']}"
        pages_done = r.get("page_count", 0)
        pages_total = r.get("total_pages", 0)
        errs = r.get("errors", [])

        lines = [
            f"## 批量转换完成 ({PIPELINE_VERSION})",
            f"📥 [下载完整 PPTX ({pages_done}/{pages_total} 页)]({dl})",
            f"📂 `{r.get('path','')}`",
            "",
        ]
        for res in r.get("results", []):
            lines.append(f"- 第 {res['page']} 页 ✅ `{res['file']}` ({res.get('text_items',0)} 个文字)")
        if errs:
            lines.append("")
            lines.append("### 失败页面")
            for e in errs:
                lines.append(f"- 第 {e['page']} 页 ❌ `{e['file']}`：{e['error']}")
        return SkillResult(success=True, message="\n".join(lines),
                          data={"skill": self.name, "download_url": dl, "stage": "completed"})

    async def _generate_vba(self, session_id: str, image_paths: list[str]) -> SkillResult:
        if not image_paths: return SkillResult(success=False, message="请先上传图片。")
        from app.services.vba_pptx_service import analyze_and_generate_vba
        r = await analyze_and_generate_vba(image_paths[0], session_id=session_id, output_dir=os.path.join("data","outputs"))
        if not r: return SkillResult(success=False, message="VBA 生成失败。")
        base_url = f"http://localhost:{settings.port}"
        dl = f"{base_url}/api/skills/download/{r.get('filename','output.bas')}"
        return SkillResult(success=True, message=f"## VBA 宏代码\n\n📥 [下载]({dl})\n在 PowerPoint 中 Alt+F11 导入并运行。", data={"skill": self.name, "download_url": dl, "stage": "completed"})

    async def _generate_svg_from_stored(self, session_id: str, session: dict) -> SkillResult:
        stored = session.get("uploaded_page_images", [])
        slides = [{"index": i+1, "title": f"Slide {i+1}"} for i in range(len(stored))]
        return await self._generate_svg_slides(session_id, slides)

    async def _generate_drawio_from_stored(self, session_id: str, session: dict) -> SkillResult:
        stored = session.get("uploaded_page_images", [])
        slides = [{"index": i+1, "title": f"Slide {i+1}"} for i in range(len(stored))]
        return await self._generate_drawio(session_id, "", slides)

    async def _collect_source_text(self, context: SkillContext) -> str:
        parts: list[str] = []
        if context.user_message.strip():
            parts.append(f"用户输入：\n{context.user_message.strip()}")
        for uploaded in context.uploaded_files:
            path = uploaded.get("path")
            filename = uploaded.get("filename") or (os.path.basename(path) if path else "upload")
            if not path:
                continue
            try:
                parsed = await asyncio.wait_for(asyncio.to_thread(parse_file_sync, path), timeout=30)
            except Exception:
                parsed = None
            if parsed:
                parts.append(f"上传文件：{filename}\n{parsed[:30000]}")
            else:
                parts.append(f"上传文件：{filename}\n系统暂时无法从该文件中提取文本。")
        return "\n\n---\n\n".join(parts).strip()

    def _has_enough_content(self, message: str, source_text: str) -> bool:
        stripped = message.strip()
        trigger_only = any(stripped.lower() == trigger.lower() for trigger in self.triggers)
        return bool(source_text) and (len(source_text) > 80 or not trigger_only)

    async def _generate_outline(
        self,
        context: SkillContext,
        session_id: str,
        source_text: str,
        revision: str | None = None,
    ) -> SkillResult:
        system_prompt = (
            "你是专业商业PPT策划顾问，拥有10年管理咨询经验。擅长将复杂信息提炼为结构清晰、"
            "观点鲜明的汇报材料。只做PPT大纲和逐页内容，不生成PPT文件。"
            "必须基于用户材料，不新增未经确认的数据、品牌、人物、产品或来源。"
            "每页内容要具体充实，避免空洞的泛词，正文要点至少3-5条。"
        )
        revision_text = f"\n\n用户修改要求：\n{revision}" if revision else ""
        prompt = f"""
请阅读以下资料，生成一份详细专业的PPT大纲和逐页内容。

要求：
1. 第一页必须是封面页（含主标题+副标题+日期/场合）。
2. 每页必须包含：
   - 页码和页面类型（封面/目录/内容/总结等）
   - 结论式标题（一句话点明本页核心观点）
   - 核心观点（2-3句展开说明）
   - 正文要点（至少3-5条具体内容，每条10-30字，避免空洞泛词）
   - 视觉建议（图表类型、配图方向、排版建议）
3. 不要生成pptx，不要生成图片。
4. 结构适合正式商业汇报，逻辑清晰、层层递进。
5. 严格基于资料，不编造未确认信息。
6. 内容要足够详细，每页正文至少150字，给后续视觉设计提供充分素材。
7. 最后询问用户是否确认大纲，确认后再进入下一步。

资料：
{source_text}
{revision_text}
""".strip()
        try:
            response = await llm_service.chat(
                interaction_name="outline_generation",
                system_prompt=system_prompt,
                messages=[{"role": "user", "content": prompt}],
                max_tokens=8192,
                temperature=0.3,
                timeout=180,
                thinking={"type": "disabled"},
            )
            outline = ""
            if response.content:
                for block in response.content:
                    if hasattr(block, "text"):
                        outline += block.text
        except Exception as exc:
            return SkillResult(success=False, message=f"生成 PPT 大纲时出错：{exc}")
        outline = outline.strip()
        if "确认" not in outline:
            outline += "\n\n请确认以上 PPT 大纲和逐页内容是否可以进入下一步。"
        _sessions[session_id] = {"stage": "awaiting_outline_confirm", "outline": outline}
        return SkillResult(success=True, message=outline, data={"skill": self.name, "stage": "awaiting_outline_confirm"})

    def _imagegen_exe(self) -> str | None:
        exe = settings.ruizhi_imagegen_exe or os.environ.get("RUIZHI_IMAGEGEN_EXE")
        if exe and os.path.exists(exe):
            return exe
        default = os.path.join(
            os.environ.get("LOCALAPPDATA", ""),
            "Programs",
            "Codex",
            "resources",
            "bin",
            "ruizhi-imagegen.exe",
        )
        return default if default and os.path.exists(default) else None

    def _imagegen_env(self) -> dict[str, str]:
        env = os.environ.copy()
        if settings.ruizhi_home:
            env["RUIZHI_HOME"] = settings.ruizhi_home
        if settings.codex_home:
            env["CODEX_HOME"] = settings.codex_home
        if settings.ruizhi_api_key:
            env["RUIZHI_API_KEY"] = settings.ruizhi_api_key
        if settings.openai_api_key:
            env["OPENAI_API_KEY"] = settings.openai_api_key
        return env

    def _visual_prompts(self, outline: str) -> list[tuple[str, str]]:
        from app.services.collage_prompt_spec import build_collage_prompt, strip_visual_suggestions
        cleaned = strip_visual_suggestions(outline)
        page_count = len(dict.fromkeys(re.findall(r'第\s*(\d+)\s*页', cleaned))) or 8
        return [
            (label, build_collage_prompt(
                total_pages=page_count,
                cleaned_outline=cleaned,
                variant_label=label,
                project_context="应用场景：正式商业汇报\n目标受众：管理层\n视觉风格：专业严谨",
            ))
            for label in ("A", "B", "C")
        ]

    async def _run_step2_visual_collages(self, session_id: str, outline: str):
        yield "进度 1/4：正在读取已确认的大纲和逐页内容...\n"
        yield "进度 2/4：正在准备三版不同视觉方向的 PPT 拼图提示词...\n"
        yield "进度 3/4：正在检查 imagegen 图片生成能力...\n"
        async for item in self._generate_visual_collages(session_id, outline):
            yield item

    async def _generate_visual_collages(self, session_id: str, outline: str):
        exe = None  # retained only for the legacy method signature
        output_dir = str(PUBLIC_DIR)
        os.makedirs(output_dir, exist_ok=True)
        run_id = uuid.uuid4().hex[:10]
        generated: list[dict] = []
        _sessions[session_id] = {**_sessions.get(session_id, {}), "stage": "generating_visual_direction", "outline": outline}

        for index, (label, prompt) in enumerate(self._visual_prompts(outline), 1):
            filename = f"ppt_maker_{session_id[:8]}_{run_id}_{label.lower()}.png"
            out_path = os.path.join(output_dir, filename)
            yield f"进度 4/4：正在生成方案 {label} 的完整 PPT 拼图（{index}/3）...\n"
            page_count = len(dict.fromkeys(re.findall(r'第\s*(\d+)\s*页', outline))) or 8
            result = await self._run_imagegen(
                exe, prompt, out_path, timeout=420,
                interaction_name="ppt_collage",
                validation_context={"expected_pages": page_count, "columns": 3, "outline": outline},
            )
            if result:
                yield SkillResult(success=False, message=f"第二步出错了：方案 {label} 未能正常生成。\n\n{result}")
                return
            item = {"label": label, "filename": filename, "path": out_path}
            generated.append(item)
            _sessions[session_id] = {**_sessions.get(session_id, {}), "visual_collages": generated}
            yield SkillResult(
                success=True,
                message=(
                    f"方案 {label} 已生成：\n\n"
                    f"![方案 {label}](/api/skills/download/{filename})\n\n"
                    f"[下载方案 {label} 拼图](/api/skills/download/{filename})\n\n"
                ),
                data={"skill": self.name, "stage": "generating_visual_direction", "visual_collages": generated},
            )

        _sessions[session_id] = {
            **_sessions.get(session_id, {}),
            "stage": "awaiting_visual_choice",
            "outline": outline,
            "visual_collages": generated,
        }
        yield SkillResult(
            success=True,
            message="第二步已完成。请在方案 A、方案 B、方案 C 中选择一个视觉方向；如果都不满意，请回复“重新生成”。",
            data={"skill": self.name, "stage": "awaiting_visual_choice", "visual_collages": generated},
        )

    async def _run_imagegen(
        self, exe: str, prompt: str, out_path: str, timeout: int,
        interaction_name: str = "ppt_slide", validation_context: dict | None = None,
    ) -> str | None:
        """Generate only through the unified paid-image quality gate."""
        try:
            from app.services.image_gen_service import generate_image

            os.makedirs(os.path.dirname(out_path) or ".", exist_ok=True)
            result = await generate_image(
                prompt, out_path, interaction_name=interaction_name,
                validation_context=validation_context,
                size="1792x1024",
            )
            if result.success:
                return None
            return result.error
        except Exception as exc:
            return f"图片生成失败：{exc}"

    def _parse_visual_choice(self, msg: str) -> str | None:
        compact = msg.strip().upper().replace(" ", "")
        compact = compact.replace("方案", "").replace("选择", "").replace("我选", "").replace("选", "")
        compact = compact.strip("。.!！,，：:")
        if compact in {"A", "B", "C"}:
            return compact
        match = re.search(r"方案\s*([ABC])", msg, re.IGNORECASE)
        return match.group(1).upper() if match else None

    def _selected_visual_choice(self, session: dict) -> str | None:
        selected = session.get("selected_visual") or {}
        label = str(selected.get("label", "")).upper()
        return label if label in {"A", "B", "C"} else None

    def _handle_visual_choice(self, context: SkillContext, session_id: str) -> SkillResult:
        msg = context.user_message.strip()
        if any(word in msg for word in ["不满意", "都不满意", "重新生成", "重做", "再生成"]):
            session = _sessions.get(session_id, {})
            session["stage"] = "awaiting_outline_confirm"
            _sessions[session_id] = session
            return SkillResult(success=True, message="好的。请告诉我重新生成前是否有特殊视觉要求；如果没有，直接回复“确认”。")
        choice = self._parse_visual_choice(msg)
        if not choice:
            return SkillResult(success=True, message="请回复方案 A、方案 B 或方案 C。若三版都不满意，也可以回复“重新生成”。")
        session = _sessions.get(session_id, {})
        generated = session.get("visual_collages") or []
        selected = next((item for item in generated if str(item.get("label", "")).upper() == choice), None)
        if not selected:
            return SkillResult(success=True, message=f"没有找到方案 {choice} 的缩略图记录。请从已生成的方案中选择 A、B 或 C。")
        session["selected_visual"] = selected
        session["stage"] = "awaiting_step3_start"
        _sessions[session_id] = session
        return SkillResult(success=True, message=f"已选择方案 {choice}。请回复“开始”或“继续”，我会进入第 3 步逐页生成高清单页图。")

    async def _handle_entry3_ready(self, context: SkillContext, session_id: str) -> SkillResult:
        """Handle user message when entry 3 is ready to go — confirm starts step 3 directly."""
        if self._is_confirm(context.user_message):
            session = _sessions.get(session_id, {})
            outline = str(session.get("outline") or "").strip()
            if not outline:
                return SkillResult(success=False, message="第 3 步出错了：没有找到页面信息。")
            _sessions[session_id] = {**session, "stage": "generating_single_pages"}
            return SkillResult(
                success=True,
                message="正在开始第 3 步逐页生成高清单页视觉稿...请稍候。",
                data={"skill": self.name, "stage": "generating_single_pages"},
            )
        if self._is_exit(context.user_message):
            _sessions.pop(session_id, None)
            return SkillResult(success=True, message="好的，已退出 PPT 制作流程。")
        return SkillResult(
            success=True,
            message="请回复「开始」或「继续」，我将直接进入第 3 步逐页生成高清单页视觉稿。如需退出请回复「退出」。",
            data={"skill": self.name, "stage": "entry3_ready"},
        )

    async def _handle_visual_direction_confirm(self, context: SkillContext, session_id: str) -> SkillResult:
        """Handle user message when outline is ready and stage is visual_direction."""
        if self._is_confirm(context.user_message):
            session = _sessions.get(session_id, {})
            outline = str(session.get("outline") or "").strip()
            if not outline:
                return SkillResult(success=False, message="第二步出错了：没有找到已确认的大纲内容。")
            _sessions[session_id] = {**session, "stage": "generating_visual_direction"}
            return SkillResult(
                success=True,
                message="正在开始第 2 步生成三版 PPT 缩略图...请稍候。",
                data={"skill": self.name, "stage": "generating_visual_direction"},
            )
        if self._is_exit(context.user_message):
            _sessions.pop(session_id, None)
            return SkillResult(success=True, message="好的，已退出 PPT 制作流程。")
        return SkillResult(
            success=True,
            message="请回复「开始」或「继续」，我将从第 2 步生成三版 PPT 缩略图。如需退出请回复「退出」。",
            data={"skill": self.name, "stage": "visual_direction"},
        )

    def _ack_step3_start(self, session_id: str) -> SkillResult:
        session = _sessions.get(session_id, {})
        choice = self._selected_visual_choice(session)
        if not choice:
            return SkillResult(success=True, message="第 3 步还不能开始：没有找到已选择的视觉方案。请重新选择方案 A、B 或 C。")
        return SkillResult(success=True, message=f"已接收到你的选择：方案 {choice}。请回复“开始”或“继续”，我会进入第 3 步。")

    def _return_generation_status(self, session_id: str) -> SkillResult:
        session = _sessions.get(session_id, {})
        page_images = session.get("single_pages") or []
        if page_images:
            message = "第 3 步正在执行，当前已生成的单页图如下：\n\n"
            for item in page_images:
                message += f"第 {item['page']} 页\n\n![第 {item['page']} 页](/api/skills/download/{item['filename']})\n\n"
        else:
            message = "第 3 步已接收并开始执行，目前正在准备或生成第 1 页。请稍候。"
        return SkillResult(success=True, message=message, data={"skill": self.name, "stage": session.get("stage", "generating_single_pages")})

    def _return_visual_collages(self, session_id: str) -> SkillResult:
        session = _sessions.get(session_id, {})
        generated = session.get("visual_collages") or []
        if not generated:
            return SkillResult(success=True, message="第二步还没有生成可查看的拼图。请回复“确认”重新开始生成三版视觉方案。")
        message = "当前已生成的 PPT 拼图方案如下：\n\n"
        for item in generated:
            message += f"方案 {item['label']}：完整 PPT 拼图\n\n![方案 {item['label']}](/api/skills/download/{item['filename']})\n\n"
        message += "请在方案 A、方案 B、方案 C 中选择一个视觉方向。"
        return SkillResult(success=True, message=message, data={"skill": self.name, "stage": session.get("stage", "awaiting_visual_choice")})

    def _extract_slide_sections(self, outline: str) -> list[dict]:
        matches = list(re.finditer(r"(第\s*\d+\s*页[:：]?.*?)(?=第\s*\d+\s*页[:：]?|\Z)", outline, re.S))
        slides: list[dict] = []
        for idx, match in enumerate(matches, 1):
            text = match.group(1).strip()
            title_line = text.splitlines()[0].strip() if text else f"第 {idx} 页"
            title = re.sub(r"^第\s*\d+\s*页[:：]?", "", title_line).strip() or title_line
            slides.append({"index": idx, "title": title, "content": text[:3000]})
        if slides:
            return slides
        chunks = [c.strip() for c in re.split(r"\n\s*\n", outline) if c.strip()]
        for idx, chunk in enumerate(chunks[:20], 1):
            title = chunk.splitlines()[0].strip()[:80] or f"第 {idx} 页"
            slides.append({"index": idx, "title": title, "content": chunk[:3000]})
        return slides

    # _style_text_for_choice, _build_visual_system_description, _detect_page_layout
    # moved to visual_systems.py — imported via `from . import visual_systems`

    async def _start_step3_single_pages(self, session_id: str, choice: str):
        session = _sessions.get(session_id, {})
        generated = session.get("visual_collages") or []
        selected = next((item for item in generated if str(item.get("label", "")).upper() == choice), None)
        outline = str(session.get("outline") or "").strip()
        if not selected:
            yield SkillResult(success=False, message=f"第 3 步出错了：没有找到方案 {choice} 的缩略图记录。请重新选择 A/B/C。")
            return
        if not outline:
            yield SkillResult(success=False, message="第 3 步出错了：没有找到已确认的大纲和逐页内容，无法生成单页高清图。")
            return
        slides = self._extract_slide_sections(outline)
        if not slides:
            yield SkillResult(success=False, message="第 3 步出错了：无法从大纲中识别页码和逐页内容。请补充每页标题和正文要点后重试。")
            return
        session["selected_visual"] = selected
        session["stage"] = "generating_single_pages"
        _sessions[session_id] = session
        total = len(slides)
        yield f"已选择方案 {choice}。\n\n第 3 步：基于拼图逐页还原高清单页视觉稿（共 {total} 页）。\n\n"
        await asyncio.sleep(0.1)
        yield "页面清单：\n"
        for s in slides:
            yield f"  第 {s['index']} 页 — {s.get('title', '')}\n"
        yield "\n按页序依次生成...\n"
        await asyncio.sleep(0.1)
        async for item in self._generate_single_page_images(session_id, choice, slides):
            yield item

    async def _generate_single_page_images(self, session_id: str, choice: str, slides: list[dict]):
        exe = self._imagegen_exe()
        if not exe:
            yield SkillResult(success=False, message="第 3 步出错了：未找到 imagegen 图片生成工具，无法生成逐页高清图。")
            return
        output_dir = str(PUBLIC_DIR)
        os.makedirs(output_dir, exist_ok=True)
        run_id = uuid.uuid4().hex[:10]
        page_images: list[dict] = []
        style = visual_systems.style_text_for_choice(choice)
        visual_system = visual_systems.build_visual_system(choice)
        total = len(slides)
        session = _sessions.get(session_id, {})
        selected_visual = session.get("selected_visual") or {}
        collage_filename = selected_visual.get("filename", "")

        for slide in slides:
            idx = slide["index"]
            filename = f"ppt_maker_{session_id[:8]}_{run_id}_page_{idx:02d}.png"
            out_path = os.path.join(output_dir, filename)

            # Build page-specific layout hint from content structure
            content = slide.get("content", "")
            page_type = slide.get("title", f"Page {idx}")
            # Detect page type for layout hint
            layout_hint = visual_systems.detect_page_layout(content)

            prompt = f"""Design task: Reproduce page {idx} of {total} from a confirmed PPT collage master as a standalone high-resolution 16:9 slide.

This is NOT a redesign. This is faithful expansion — taking one page from the collage and rendering it as a clean, sharp, full-resolution single slide. The collage is the master reference for visual style, layout, and content structure.

COLLAGE CONTEXT:
- Selected visual direction: Plan {choice} — {style}
- Collage filename: {collage_filename}
- Page {idx} of {total}

VISUAL SYSTEM (must replicate exactly from the collage):
{visual_system}

LAYOUT TYPE FOR THIS PAGE:
{layout_hint}

CONFIRMED CONTENT (from the approved outline):
{content[:2000]}

CRITICAL RULES:
1. Output page {idx} ONLY. One single 16:9 slide image. No collage, no multi-page, no PPTX file.
2. Match the collage page layout faithfully — same title position, content blocks, chart areas, module boundaries. Just clearer and higher resolution.
3. {idx} must be clearly visible as the page number.
4. All Chinese text must be real, correct, readable Chinese characters. NO gibberish or pseudo-Chinese.
5. Use ONLY confirmed content above. Never invent brands, logos, people, products, data, or sources.
6. If collage details (small text, chart numbers) are unclear, supplement from the confirmed outline above while preserving the original visual layout and proportions.
7. Redraw complex charts at high resolution matching the visual proportions and style — do not fabricate precise data values.
8. The result must look like a finished, formal business PPT slide — NOT a wireframe, sketch, or design note.
9. Maintain the exact same design system as the collage across all pages. Allow for the layout variations that already exist in the collage.
10. Do NOT simply crop-zoom the collage. Re-render each page at full resolution with clearer text, icons, charts, and details.
11. This slide will be used as a blueprint for an editable PPTX — composition, text hierarchy, module boundaries, chart relationships, page number, and visual rhythm must be clear, stable, and reproducible."""

            yield f"正在生成第 {idx}/{total} 页单页 PPT 视觉稿...\n"
            await asyncio.sleep(0.1)
            error = await self._run_imagegen(
                exe, prompt, out_path, timeout=420,
                interaction_name="ppt_slide",
                validation_context={"page": idx, "expected_aspect": 16 / 9},
            )
            if error:
                yield SkillResult(
                    success=False,
                    message=f"第 3 步出错了：第 {idx} 页未能正常生成。\n\n{error}",
                    data={"skill": self.name, "stage": "generating_single_pages", "single_pages": page_images},
                )
                return
            page_images.append({"page": idx, "title": slide.get("title", ""), "filename": filename, "path": out_path})
            _sessions[session_id] = {**_sessions.get(session_id, {}), "stage": "generating_single_pages", "single_pages": page_images}
            yield SkillResult(
                success=True,
                message=(
                    f"第 {idx} 页高清 PPT 风格图已生成：\n\n"
                    f"![第 {idx} 页高清风格图](/api/skills/download/{filename})\n\n"
                    f"[下载第 {idx} 页高清图](/api/skills/download/{filename})\n\n"
                    "我会继续生成下一页。"
                ),
                data={"skill": self.name, "stage": "generating_single_pages", "single_pages": page_images, "latest_page": page_images[-1]},
            )
            await asyncio.sleep(0.1)
        _sessions[session_id] = {**_sessions.get(session_id, {}), "stage": "awaiting_editable_ppt_scope", "single_pages": page_images}
        yield SkillResult(
            success=True,
            message="第 3 步已完成。你希望把哪一部分做成可编辑 PPT？可以回复某一页、几页，或回复“全部”。",
            data={"skill": self.name, "stage": "awaiting_editable_ppt_scope", "single_pages": page_images},
        )

    def _select_page_images(self, request: str, page_images: list[dict]) -> list[dict]:
        if not page_images:
            return []
        req = (request or "").strip().lower()
        if not req or "全部" in req or "all" in req or "整套" in req:
            return page_images
        selected_pages: set[int] = set()
        for start, end in re.findall(r"(\d+)\s*[-~至到]\s*(\d+)", req):
            a, b = int(start), int(end)
            if a > b:
                a, b = b, a
            selected_pages.update(range(a, b + 1))
        for num in re.findall(r"\d+", req):
            selected_pages.add(int(num))
        if not selected_pages:
            return page_images
        return [item for item in page_images if int(item.get("page", 0)) in selected_pages]

    def _build_pptx_from_page_images(self, session_id: str, page_images: list[dict]) -> dict:
        from pptx import Presentation
        from pptx.util import Inches

        output_dir = str(PUBLIC_DIR)
        os.makedirs(output_dir, exist_ok=True)
        filename = f"ppt_maker_{session_id[:8]}_{uuid.uuid4().hex[:10]}.pptx"
        out_path = os.path.join(output_dir, filename)
        prs = Presentation()
        prs.slide_width = Inches(13.333333)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]
        for item in page_images:
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(item["path"], 0, 0, width=prs.slide_width, height=prs.slide_height)
        prs.save(out_path)
        return {"filename": filename, "path": out_path, "download_url": f"/api/skills/download/{filename}"}

    def _handle_editable_ppt_scope(self, context: SkillContext, session_id: str) -> SkillResult:
        session = _sessions.get(session_id, {})
        page_images = session.get("single_pages") or []
        if not page_images:
            return SkillResult(success=False, message="第 4 步出错了：没有找到第 3 步生成的高清页面图，无法生成 PPTX。")
        selected = self._select_page_images(context.user_message.strip(), page_images)
        if not selected:
            return SkillResult(success=False, message="没有匹配到你指定的页码。请回复“全部”或类似“第 1 页”“1-3 页”的范围。")
        try:
            pptx = self._build_pptx_from_page_images(session_id, selected)
        except Exception as exc:
            return SkillResult(success=False, message=f"第 4 步出错了：生成 PPTX 失败。\n\n错误类型：{exc.__class__.__name__}\n错误原因：{str(exc) or '无详细错误文本'}")
        session["editable_scope_request"] = context.user_message.strip()
        session["stage"] = "completed"
        session["pptx"] = pptx
        _sessions[session_id] = session
        return SkillResult(
            success=True,
            message=(
                "第 4 步已完成，PPTX 文件已生成。\n\n"
                f"[下载 PPTX 文件]({pptx['download_url']})\n\n"
                "说明：当前版本先将第 3 步生成的高清视觉稿按页放入 16:9 PPTX，保证可下载、可演示。"
                "后续如果要做到主要文字完全可编辑，还需要继续做文字识别、遮罩和文本框重建。"
            ),
            data={"skill": self.name, "stage": "completed", "download_url": pptx["download_url"], "filename": pptx["filename"], "path": pptx["path"]},
            follow_up_action="download",
        )
