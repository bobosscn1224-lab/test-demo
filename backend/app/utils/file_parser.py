import os
import io
import logging

logger = logging.getLogger(__name__)

TEXT_EXTENSIONS = {
    ".txt", ".md", ".py", ".js", ".ts", ".jsx", ".tsx", ".go", ".java",
    ".yaml", ".yml", ".json", ".csv", ".xml", ".html", ".css", ".sql",
    ".sh", ".bat", ".ps1", ".ini", ".cfg", ".toml", ".rst",
}

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tiff", ".tif"}

SUPPORTED_EXTENSIONS = TEXT_EXTENSIONS | IMAGE_EXTENSIONS | {".pdf", ".docx", ".xlsx", ".pptx"}


def parse_file_sync(file_path: str) -> str | None:
    """Parse a single file and return its text content. Returns None if unsupported or unreadable."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext not in SUPPORTED_EXTENSIONS:
        return None

    try:
        if ext in TEXT_EXTENSIONS:
            return _read_text(file_path)
        elif ext in IMAGE_EXTENSIONS:
            return None  # images handled by async path only
        elif ext == ".pdf":
            return _read_pdf(file_path)
        elif ext == ".docx":
            return _read_docx(file_path)
        elif ext == ".xlsx":
            return _read_xlsx(file_path)
        elif ext == ".pptx":
            return _read_pptx(file_path)
    except Exception:
        logger.warning("Failed to parse %s", file_path, exc_info=True)
        return None


async def parse_file(file_path: str) -> str | None:
    """Async wrapper — tries text extraction first, then async image-aware parsing."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext in IMAGE_EXTENSIONS:
        return await _read_image_file(file_path)
    if ext in TEXT_EXTENSIONS:
        return _read_text(file_path)

    # For rich formats, use async path that includes image descriptions
    try:
        if ext == ".pdf":
            return await _read_pdf_async(file_path)
        elif ext == ".docx":
            return await _read_docx_async(file_path)
        elif ext == ".xlsx":
            return _read_xlsx(file_path)
        elif ext == ".pptx":
            return await _read_pptx_async(file_path)
    except Exception:
        logger.warning("Failed to parse %s", file_path, exc_info=True)
        return None


def _read_text(file_path: str) -> str:
    encodings = ["utf-8", "gbk", "gb2312", "latin-1"]
    for enc in encodings:
        try:
            with open(file_path, "r", encoding=enc) as f:
                return f.read()
        except UnicodeDecodeError:
            continue
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def _read_pdf(file_path: str) -> str:
    import fitz
    doc = fitz.open(file_path)
    try:
        parts = []
        for page in doc:
            # Try text extraction with word-level sorting first (fixes single-char lines)
            text = _extract_pdf_page_text(page)
            if text:
                parts.append(text)
        return "\n".join(parts)
    finally:
        doc.close()


def _extract_pdf_page_text(page) -> str:
    """Extract text from a PDF page, handling single-character layout issues.

    Strategy: always use word-level extraction with line grouping (Method 1).
    It produces better results than raw text even for single-char-heavy pages.
    """
    words = page.get_text("words")
    if words:
        words.sort(key=lambda w: (round(w[1], 0), w[0]))
        lines = []
        current_line = []
        current_y = round(words[0][1], 0) if words else 0
        for w in words:
            wy = round(w[1], 0)
            if abs(wy - current_y) > 3:
                if current_line:
                    lines.append(" ".join(current_line))
                current_line = [w[4]]
                current_y = wy
            else:
                current_line.append(w[4])
        if current_line:
            lines.append(" ".join(current_line))
        result = "\n".join(lines)
        normal_lines = [l for l in lines if len(l) > 5]
        if len(normal_lines) > len(lines) * 0.5:
            return result
        # Low-quality but still potentially useful — clean it further
        cleaned = _merge_single_chars(result)
        if cleaned:
            return cleaned

    # Method 2: No words at all (e.g., image-only page)
    raw = page.get_text("text")
    return _merge_single_chars(raw)


def _merge_single_chars(text: str) -> str:
    """Merge isolated single characters into words. Handles both:
    - Single char per line (PDF extraction artifact)
    - Single chars separated by spaces on the same line
    """
    lines = text.split("\n")
    out = []
    for line in lines:
        stripped = line.strip()
        if not stripped:
            if out:
                out.append("")
            continue
        # Split line into tokens
        tokens = stripped.split()
        # Merge single-char tokens with neighbors
        merged = []
        buf = ""
        for tok in tokens:
            if len(tok) == 1:
                buf += tok
            else:
                if buf:
                    merged.append(buf)
                    buf = ""
                merged.append(tok)
        if buf:
            merged.append(buf)
        out.append(" ".join(merged) if merged else stripped)
    return "\n".join(out)


def _read_docx(file_path: str) -> str:
    from docx import Document
    doc = Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs if p.text)


def _read_xlsx(file_path: str) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(file_path, data_only=True)
    try:
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"--- Sheet: {sheet_name} ---")
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(c) if c is not None else "" for c in row)
                if row_text.strip():
                    parts.append(row_text)
        return "\n".join(parts)
    finally:
        wb.close()


def _read_pptx(file_path: str) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx not installed, cannot parse %s", file_path)
        return None

    prs = Presentation(file_path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        parts.append(text)
    return "\n".join(parts)


# ── Async image-aware parsers ──────────────────────────────────────────

async def _read_image_file(file_path: str) -> str | None:
    """Describe a standalone image file using the Vision API."""
    try:
        with open(file_path, "rb") as f:
            image_bytes = f.read()
    except OSError:
        return None

    from app.services.image_service import image_service
    fname = os.path.basename(file_path)
    description = await image_service.describe_image(image_bytes, source_hint=fname)
    if description:
        return f"[图片: {fname}]\n{description}"
    return None


async def _read_pdf_async(file_path: str) -> str:
    """Parse PDF text AND describe embedded images."""
    import fitz
    from app.services.image_service import image_service

    doc = fitz.open(file_path)
    try:
        parts = []
        for page_idx, page in enumerate(doc):
            page_num = page_idx + 1

            # Extract text
            text = _extract_pdf_page_text(page)
            if text:
                parts.append(text)

            # Extract and describe images on this page
            image_list = page.get_images()
            if image_list:
                for img_idx, img_info in enumerate(image_list):
                    try:
                        xref = img_info[0]
                        base_image = doc.extract_image(xref)
                        img_bytes = base_image["image"]
                        img_ext = base_image.get("ext", "png")
                        hint = f"{os.path.basename(file_path)} 第{page_num}页 图片{img_idx + 1}"
                        description = await image_service.describe_image(img_bytes, source_hint=hint)
                        if description:
                            parts.append(f"[文档图片 — 第{page_num}页 图{img_idx + 1}]\n{description}")
                    except Exception:
                        logger.debug("Failed to describe image %d on page %d of %s", img_idx, page_num, file_path)

        return "\n".join(parts)
    finally:
        doc.close()


async def _read_docx_async(file_path: str) -> str:
    """Parse DOCX text AND describe embedded images."""
    from docx import Document
    from app.services.image_service import image_service

    doc = Document(file_path)
    parts = []

    # Extract images from the DOCX zip
    image_map: dict[str, bytes] = {}
    try:
        import zipfile
        with zipfile.ZipFile(file_path, "r") as zf:
            for name in zf.namelist():
                if name.startswith("word/media/") and any(
                    name.lower().endswith(ext) for ext in (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp")
                ):
                    image_map[name] = zf.read(name)
    except Exception:
        logger.debug("Could not extract images from DOCX zip", exc_info=True)

    img_idx = 0
    for para in doc.paragraphs:
        if para.text:
            parts.append(para.text)

        # Check for inline images in this paragraph
        for run in para.runs:
            if run._element.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}drawing'):
                # Found an image — try to find it in the document relationships
                for rel in doc.part.rels.values():
                    if "image" in rel.reltype and hasattr(rel, "target_part"):
                        try:
                            img_bytes = rel.target_part.blob
                        except Exception:
                            continue
                        if img_bytes:
                            img_idx += 1
                            hint = f"{os.path.basename(file_path)} 图片{img_idx}"
                            description = await image_service.describe_image(img_bytes, source_hint=hint)
                            if description:
                                parts.append(f"[文档图片 — 图{img_idx}]\n{description}")

    # Also describe any orphan images not referenced in paragraphs
    for rel in doc.part.rels.values():
        if "image" in rel.reltype and hasattr(rel, "target_part"):
            try:
                img_bytes = rel.target_part.blob
            except Exception:
                continue
            if img_bytes:
                img_idx += 1
                hint = f"{os.path.basename(file_path)} 图片{img_idx}"
                description = await image_service.describe_image(img_bytes, source_hint=hint)
                if description:
                    parts.append(f"[文档图片 — 图{img_idx}]\n{description}")

    return "\n".join(parts) if parts else _read_docx(file_path)


async def _read_pptx_async(file_path: str) -> str:
    """Parse PPTX text AND describe embedded images."""
    try:
        from pptx import Presentation
        from pptx.shapes.picture import Picture
    except ImportError:
        return _read_pptx(file_path)

    from app.services.image_service import image_service

    prs = Presentation(file_path)
    parts = []
    img_idx = 0

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        has_text = False

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        if not has_text:
                            parts.append(f"--- Slide {slide_num} ---")
                            has_text = True
                        parts.append(text)

            if isinstance(shape, Picture):
                try:
                    img_bytes = shape.image.blob
                    img_idx += 1
                    hint = f"{os.path.basename(file_path)} 幻灯片{slide_num} 图片{img_idx}"
                    description = await image_service.describe_image(img_bytes, source_hint=hint)
                    if description:
                        if not has_text:
                            parts.append(f"--- Slide {slide_num} ---")
                            has_text = True
                        parts.append(f"[幻灯片图片 — 第{slide_num}页 图{img_idx}]\n{description}")
                except Exception:
                    logger.debug("Failed to describe image on slide %d of %s", slide_num, file_path)

    return "\n".join(parts) if parts else _read_pptx(file_path)
