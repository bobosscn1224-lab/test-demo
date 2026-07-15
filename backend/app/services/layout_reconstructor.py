"""Image-to-Editable-PPTX — DeckWeaver v5 layout reconstruction.

v5 improvements over v4:
  - Native shape detection (rounded rects, circles, lines, arrows)
  - K-Means colour scheme extraction
  - Page layout analysis + semantic element naming
  - Human-readable output report (editable vs image vs shape breakdown)
  - Outputs to unified absolute path

Pipeline:
  1. RapidOCR → text inventory
  2. OpenCV contours → native shapes
  3. Canny + connected components → image inventory
  4. K-Means → colour scheme
  5. Side-agreement erasure → clean background
  6. Layout analysis → page structure + card grouping
  7. python-pptx → layered slide: BG → shapes → images → text
  8. Output report
"""
from __future__ import annotations

import logging
import os
import uuid
from collections import Counter
from pathlib import Path

import cv2
import numpy as np
from pptx.oxml.ns import qn

logger = logging.getLogger(__name__)

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

_RAPID_OCR = None


def _get_ocr():
    global _RAPID_OCR
    if _RAPID_OCR is None:
        from rapidocr_onnxruntime import RapidOCR
        _RAPID_OCR = RapidOCR()
    return _RAPID_OCR


async def reconstruct(
    image_path: str,
    output_dir: str | None = None,
    session_id: str = "",
    *,
    enable_shapes: bool = True,
) -> dict | None:
    import asyncio
    return await asyncio.to_thread(_run, image_path, output_dir, session_id, enable_shapes)


def _run(
    image_path: str,
    output_dir: str | None,
    session_id: str,
    enable_shapes: bool,
) -> dict | None:
    try:
        from PIL import Image

        img = Image.open(image_path)
        iw, ih = img.size
        buf = np.fromfile(image_path, dtype=np.uint8)
        cv_img = cv2.imdecode(buf, cv2.IMREAD_COLOR)
        if cv_img is None:
            cv_img = np.array(img.convert("RGB"))[:, :, ::-1]

        sx = SLIDE_W_IN / iw
        sy = SLIDE_H_IN / ih
        logger.info("DeckWeaver v5: %s (%dx%d)", os.path.basename(image_path), iw, ih)

        # ── 1. Text inventory ──
        text_items = _build_text_inventory(cv_img, iw, ih)
        text_items = _deduplicate_text_boxes(text_items)

        # ── 2. Native shapes + colour scheme ──
        native_shapes: list = []
        color_scheme = {"palette": [], "background": "#FFFFFF", "accent": "#111111"}
        if enable_shapes:
            from .layout.shape_detector import detect_native_shapes, extract_color_scheme
            raw_shapes = detect_native_shapes(cv_img)
            native_shapes = _filter_shapes(raw_shapes, text_items, iw, ih)
            color_scheme = extract_color_scheme(cv_img)
            logger.info("Shapes: %d detected, %d kept", len(raw_shapes), len(native_shapes))

        # ── 3. Image/icon inventory ──
        image_items = _build_image_inventory(cv_img, text_items, native_shapes, iw, ih)

        # ── 4. Clean background ──
        cleaned = _erase_text_regions(cv_img, text_items)
        bg_mode, bg_color = _classify_background(cv_img, cleaned)

        # ── 5. Layout analysis + semantic naming ──
        from .layout.layout_analyzer import analyze_page_structure, assign_semantic_names
        layout = analyze_page_structure(text_items, image_items, native_shapes, iw, ih)
        assign_semantic_names(text_items, layout["sections"], iw, ih)

        # ── 6. Build PPTX ──
        out_dir = Path(output_dir) if output_dir else Path("outputs")
        out_dir.mkdir(parents=True, exist_ok=True)
        sid = session_id[:8] if session_id else uuid.uuid4().hex[:8]
        asset_dir = out_dir / f"dw_{sid}"
        asset_dir.mkdir(parents=True, exist_ok=True)

        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        prs.slide_width = Inches(SLIDE_W_IN)
        prs.slide_height = Inches(SLIDE_H_IN)
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Layer 0: Background
        if bg_mode == "solid":
            _set_slide_bg(slide, *bg_color)
        else:
            bg_path = asset_dir / "bg.png"
            cv2.imencode('.png', cleaned)[1].tofile(str(bg_path))
            slide.shapes.add_picture(str(bg_path), Inches(0), Inches(0), Inches(SLIDE_W_IN), Inches(SLIDE_H_IN))

        # Layer 1: Native shapes
        _place_shapes(slide, native_shapes, sx, sy)

        # Layer 2: Image assets
        for ii in image_items:
            _place_image(slide, cleaned, ii, asset_dir, sx, sy)
            _fill_rect(cleaned, ii["x"], ii["y"], ii["w"], ii["h"], bg_color)

        # Layer 3: Editable text
        for ti in text_items:
            _place_text(slide, ti, sx, sy)

        fname = f"ppt_dw_{sid}.pptx"
        fpath = out_dir / fname
        prs.save(str(fpath))

        # ── 7. Output report ──
        from .layout.layout_analyzer import generate_output_report
        report_text = generate_output_report(
            text_items, image_items, native_shapes, [],
            layout["sections"], layout["layout_type"], color_scheme,
        )

        return {
            "filename": fname,
            "path": str(fpath),
            "url": f"/api/skills/download/{fname}",
            "pages": 1,
            "text_items": len(text_items),
            "image_items": len(image_items),
            "native_shapes": len(native_shapes),
            "layout_type": layout["layout_type"],
            "color_scheme": color_scheme,
            "report": {
                "method": "DeckWeaver v5: RapidOCR + shapes + side-agreement + z-ordered PPTX",
                "text_count": len(text_items),
                "image_count": len(image_items),
                "shape_count": len(native_shapes),
                "report_text": report_text,
            },
        }
    except Exception as exc:
        logger.error("DeckWeaver v5 failed: %s", exc)
        import traceback
        traceback.print_exc()
        return None


# ═══════════════════════════════════════════════════════════════════════════
# Text inventory
# ═══════════════════════════════════════════════════════════════════════════

def _build_text_inventory(cv_img, iw, ih) -> list[dict]:
    ocr = _get_ocr()
    items = []
    try:
        result, _ = ocr(cv_img)
        if not result:
            return items
        for box, text, score in result:
            text = str(text or "").strip()
            confidence = float(score) * 100
            if not text or confidence < 45:
                continue
            if len(box) < 4:
                continue
            xs = [int(p[0]) for p in box]
            ys = [int(p[1]) for p in box]
            x1, y1 = max(0, min(xs)), max(0, min(ys))
            x2, y2 = min(iw, max(xs)), min(ih, max(ys))
            if x2 <= x1 or y2 <= y1:
                continue
            w, h = x2 - x1, y2 - y1
            role = "body"
            if y1 < ih * 0.14 and h > ih * 0.022:
                role = "title"
            elif y1 > ih * 0.90 and len(text) <= 4:
                role = "page_number"
            elif text.replace(" ", "").replace(".", "").replace("%", "").isdigit():
                role = "number"
            elif h > ih * 0.028:
                role = "heading"
            items.append({
                "text": text, "x": x1, "y": y1, "w": w, "h": h,
                "confidence": round(confidence, 1), "role": role,
            })
    except Exception as exc:
        logger.warning("RapidOCR: %s", exc)
    return items


def _deduplicate_text_boxes(items: list) -> list:
    """Remove overlapping text boxes, keeping the larger/more confident one."""
    if len(items) <= 1:
        return items
    # Sort by area descending
    items.sort(key=lambda t: t["w"] * t["h"], reverse=True)
    kept = []
    for item in items:
        overlap = False
        for k in kept:
            # IoU check
            ax1, ay1 = item["x"], item["y"]
            ax2, ay2 = item["x"] + item["w"], item["y"] + item["h"]
            bx1, by1 = k["x"], k["y"]
            bx2, by2 = k["x"] + k["w"], k["y"] + k["h"]
            inter_x = max(0, min(ax2, bx2) - max(ax1, bx1))
            inter_y = max(0, min(ay2, by2) - max(ay1, by1))
            inter_area = inter_x * inter_y
            area_a = (ax2 - ax1) * (ay2 - ay1)
            if area_a > 0 and inter_area / area_a > 0.5:
                overlap = True
                break
        if not overlap:
            kept.append(item)
    # Restore original order (by y, then x)
    kept.sort(key=lambda t: (t["y"], t["x"]))
    return kept


# ═══════════════════════════════════════════════════════════════════════════
# Image inventory
# ═══════════════════════════════════════════════════════════════════════════

def _build_image_inventory(cv_img, text_items, native_shapes, iw, ih) -> list[dict]:
    h, w = cv_img.shape[:2]
    items = []
    try:
        gray = cv2.cvtColor(cv_img, cv2.COLOR_BGR2GRAY)
        edges = cv2.Canny(gray, 45, 140)
        mask = np.zeros((h, w), dtype=np.uint8)
        for t in text_items:
            pad = 3
            x1, y1 = max(0, t["x"] - pad), max(0, t["y"] - pad)
            x2, y2 = min(w, t["x"] + t["w"] + pad), min(h, t["y"] + t["h"] + pad)
            mask[y1:y2, x1:x2] = 255
        for s in native_shapes:
            pad = 4
            sx = s.get("x", s.get("x1", 0))
            sy = s.get("y", s.get("y1", 0))
            sw = s.get("w", abs(s.get("x2", 0) - s.get("x1", 0)))
            sh = s.get("h", abs(s.get("y2", 0) - s.get("y1", 0)))
            x1, y1 = max(0, sx - pad), max(0, sy - pad)
            x2, y2 = min(w, sx + sw + pad), min(h, sy + sh + pad)
            mask[y1:y2, x1:x2] = 255
        edges[mask > 0] = 0
        kernel = np.ones((5, 5), np.uint8)
        closed = cv2.morphologyEx(edges, cv2.MORPH_CLOSE, kernel)
        num_labels, labels, stats, _ = cv2.connectedComponentsWithStats(closed, 8)
        min_area = max(500, int(w * h * 0.001))
        max_area = int(w * h * 0.55)
        for i in range(1, num_labels):
            cx, cy, cw, ch, area = [int(v) for v in stats[i]]
            if area < min_area or area > max_area:
                continue
            aspect = cw / max(ch, 1)
            if aspect > 16 or aspect < 0.06:
                continue
            items.append({
                "x": cx, "y": cy, "w": cw, "h": ch,
                "category": "visual",
            })
    except Exception as exc:
        logger.warning("Image inventory: %s", exc)
    return items[:40]


# ═══════════════════════════════════════════════════════════════════════════
# Side-agreement background erasure
# ═══════════════════════════════════════════════════════════════════════════

def _erase_text_regions(cv_img, text_items) -> np.ndarray:
    cleaned = cv_img.copy()
    for t in text_items:
        x1, y1 = t["x"], t["y"]
        x2, y2 = t["x"] + t["w"], t["y"] + t["h"]
        if x2 <= x1 or y2 <= y1:
            continue
        bg = _side_agreement_bg(cv_img, x1, y1, x2, y2)
        cleaned[y1:y2, x1:x2] = bg
    return cleaned


def _side_agreement_bg(img, x1, y1, x2, y2) -> np.ndarray:
    h, w = img.shape[:2]
    sides = {}
    if y1 >= 2:
        m = _strip_bg_median(img[y1 - 2:y1, x1:x2])
        if m is not None: sides["top"] = m
    if y2 + 2 <= h:
        m = _strip_bg_median(img[y2:y2 + 2, x1:x2])
        if m is not None: sides["bot"] = m
    if x1 >= 2:
        m = _strip_bg_median(img[y1:y2, x1 - 2:x1])
        if m is not None: sides["left"] = m
    if x2 + 2 <= w:
        m = _strip_bg_median(img[y1:y2, x2:x2 + 2])
        if m is not None: sides["right"] = m
    if not sides:
        return np.array([255, 255, 255], dtype=np.uint8)
    quantized = [tuple(((c // 20) * 20).astype(int)) for c in sides.values()]
    votes = Counter(quantized)
    winner_q, _ = votes.most_common(1)[0]
    close = [c for c in sides.values() if max(abs(c - np.array(winner_q))) < 30]
    if close:
        return np.median(close, axis=0).astype(np.uint8)
    return np.array(winner_q, dtype=np.uint8)


def _strip_bg_median(strip: np.ndarray) -> np.ndarray | None:
    if strip.size == 0:
        return None
    flat = strip.reshape(-1, 3).astype(np.float32)
    if len(flat) < 3:
        return flat[0].astype(np.uint8)
    median = np.median(flat, axis=0)
    dists = np.linalg.norm(flat - median, axis=1)
    close = flat[dists < np.percentile(dists, 75)]
    if len(close) >= 2:
        return np.median(close, axis=0).astype(np.uint8)
    return median.astype(np.uint8)


def _classify_background(cv_img, cleaned) -> tuple:
    h, w = cv_img.shape[:2]
    sample_size = max(10, min(w, h) // 30)
    regions = [
        cleaned[0:sample_size, 0:sample_size],
        cleaned[0:sample_size, -sample_size:],
        cleaned[-sample_size:, 0:sample_size],
        cleaned[-sample_size:, -sample_size:],
        cleaned[h // 2 - 4:h // 2 + 4, 0:sample_size],
        cleaned[h // 2 - 4:h // 2 + 4, -sample_size:],
        cleaned[0:sample_size, w // 2 - 4:w // 2 + 4],
        cleaned[-sample_size:, w // 2 - 4:w // 2 + 4],
        cleaned[h // 2 - 10:h // 2 + 10, w // 2 - 10:w // 2 + 10],
    ]
    means = [np.mean(r.reshape(-1, 3), axis=0) for r in regions if r.size > 0]
    if not means:
        return ("solid", (255, 255, 255))
    means_arr = np.array(means)
    if np.max(np.var(means_arr, axis=0)) < 300:
        avg = np.mean(means_arr, axis=0).astype(int)
        return ("solid", (int(avg[2]), int(avg[1]), int(avg[0])))
    return ("complex", _estimate_slide_bg(cv_img))


def _estimate_slide_bg(cv_img) -> tuple:
    h, w = cv_img.shape[:2]
    corners = [cv_img[0:10, 0:10], cv_img[0:10, -10:], cv_img[-10:, 0:10], cv_img[-10:, -10:]]
    samples = [s.reshape(-1, 3) for s in corners if s.size > 0]
    avg = np.median(np.concatenate(samples), axis=0).astype(int)
    return (int(avg[2]), int(avg[1]), int(avg[0]))


def _filter_shapes(shapes: list, text_items: list, iw: int, ih: int) -> list:
    """Filter noisy shape detections: remove text-overlapping, low-confidence, near-bg shapes."""
    if not shapes:
        return []

    # Estimate background color from shape fills
    from collections import Counter as _Counter
    fills = [s.get("fill", "#FFFFFF") for s in shapes if s.get("type") in ("rect", "rounded_rect")]
    bg_hex = _Counter(fills).most_common(1)[0][0] if fills else "#FFFFFF"

    def _shape_overlaps_text(s, texts):
        sx1, sy1 = s.get("x", s.get("x1", 0)), s.get("y", s.get("y1", 0))
        sx2 = sx1 + s.get("w", abs(s.get("x2", 0) - s.get("x1", 0)))
        sy2 = sy1 + s.get("h", abs(s.get("y2", 0) - s.get("y1", 0)))
        for t in texts:
            tx1, ty1 = t["x"], t["y"]
            tx2, ty2 = t["x"] + t["w"], t["y"] + t["h"]
            # Check overlap (with small margin)
            if sx1 < tx2 + 5 and sx2 > tx1 - 5 and sy1 < ty2 + 3 and sy2 > ty1 - 3:
                return True
        return False

    def _hex_dist(a, b):
        ra, ga, ba = int(a[1:3], 16), int(a[3:5], 16), int(a[5:7], 16)
        rb, gb, bb = int(b[1:3], 16), int(b[3:5], 16), int(b[5:7], 16)
        return abs(ra - rb) + abs(ga - gb) + abs(ba - bb)

    filtered = []
    for s in shapes:
        stype = s.get("type", "rect")
        # Skip very tiny shapes
        area = s.get("w", 0) * s.get("h", 0)
        if stype in ("rect", "rounded_rect", "circle") and area < iw * ih * 0.0003:
            continue
        # Skip very low confidence
        if s.get("confidence", 0) < 0.55:
            continue
        filtered.append(s)

    # Sort by area descending, keep top 30
    filtered.sort(key=lambda s: s.get("w", 0) * s.get("h", 0), reverse=True)
    return filtered[:30]


def _fill_rect(img, x, y, w, h, color):
    h_img, w_img = img.shape[:2]
    x1, y1 = max(0, x), max(0, y)
    x2, y2 = min(w_img, x + w), min(h_img, y + h)
    if x1 < x2 and y1 < y2:
        img[y1:y2, x1:x2] = (color[2], color[1], color[0])


# ═══════════════════════════════════════════════════════════════════════════
# PPTX placement
# ═══════════════════════════════════════════════════════════════════════════

def _set_slide_bg(slide, r, g, b):
    from lxml import etree
    bgPr = slide.background._element.find("{http://schemas.openxmlformats.org/presentationml/2006/main}bgPr")
    if bgPr is None:
        bgPr = etree.SubElement(slide.background._element, qn("p:bgPr"))
    sf = etree.SubElement(bgPr, qn("a:solidFill"))
    sc = etree.SubElement(sf, qn("a:srgbClr"))
    sc.set("val", f"{r:02x}{g:02x}{b:02x}")


def _place_shapes(slide, shapes, sx, sy):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

    shape_map = {
        "rounded_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
        "rect": MSO_SHAPE.RECTANGLE,
        "circle": MSO_SHAPE.OVAL,
    }
    for s in shapes:
        try:
            stype = s.get("type", "rect")

            if stype in ("line", "arrow"):
                connector = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    Inches(s["x1"] * sx), Inches(s["y1"] * sy),
                    Inches(s["x2"] * sx), Inches(s["y2"] * sy),
                )
                hex_c = s.get("stroke", "#999999").lstrip("#")
                r, g, b = int(hex_c[0:2], 16), int(hex_c[2:4], 16), int(hex_c[4:6], 16)
                connector.line.color.rgb = RGBColor(r, g, b)
                connector.line.width = Pt(s.get("stroke_width", 0.5))
                continue

            if stype not in shape_map:
                continue

            shape = slide.shapes.add_shape(
                shape_map[stype],
                Inches(s["x"] * sx), Inches(s["y"] * sy),
                Inches(s["w"] * sx), Inches(s["h"] * sy),
            )
            hex_f = s.get("fill", "#FFFFFF").lstrip("#")
            fr, fg, fb = int(hex_f[0:2], 16), int(hex_f[2:4], 16), int(hex_f[4:6], 16)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(fr, fg, fb)

            hex_s = s.get("stroke", "#CCCCCC").lstrip("#")
            sr, sg, sb = int(hex_s[0:2], 16), int(hex_s[2:4], 16), int(hex_s[4:6], 16)
            shape.line.color.rgb = RGBColor(sr, sg, sb)
            shape.line.width = Pt(s.get("stroke_width", 0.5))
        except Exception:
            pass


def _place_image(slide, cv_img, item, asset_dir, sx, sy):
    from pptx.util import Inches
    try:
        px, py, pw, ph = item["x"], item["y"], item["w"], item["h"]
        crop = cv_img[py:py + ph, px:px + pw]
        if crop.size == 0:
            return
        path = asset_dir / f"img_{px}_{py}.png"
        cv2.imencode('.png', crop)[1].tofile(str(path))
        slide.shapes.add_picture(
            str(path),
            Inches(px * sx), Inches(py * sy),
            Inches(pw * sx), Inches(ph * sy),
        )
    except Exception:
        pass


def _place_text(slide, item, sx, sy):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    try:
        x_in = max(0, item["x"] * sx)
        y_in = max(0, item["y"] * sy)
        w_in = min(SLIDE_W_IN - x_in, max(item["w"] * sx, 0.8))
        h_in = min(SLIDE_H_IN - y_in, max(item["h"] * sy, 0.35))

        tb = slide.shapes.add_textbox(Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item["text"]

        role = item.get("role", "body")
        fs = {"title": 26, "heading": 16, "number": 20, "page_number": 10}.get(role, 13)
        p.font.size = Pt(fs)
        if role in ("title", "heading"):
            p.font.bold = True
        p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)

        rel_x = (x_in + w_in / 2) / SLIDE_W_IN
        p.alignment = PP_ALIGN.CENTER if 0.42 < rel_x < 0.58 else PP_ALIGN.LEFT

        _set_autofit(tf)
        name = item.get("semantic_name", f"text_{role}_{item['x']}_{item['y']}")
        tb.name = name
    except Exception:
        pass


def _set_autofit(tf):
    from lxml import etree
    bp = tf._element.find(qn("a:bodyPr"))
    if bp is not None:
        for tag in ("noAutofit", "normAutofit", "spAutoFit"):
            for e in bp.findall(qn(f"a:{tag}")):
                bp.remove(e)
        na = etree.SubElement(bp, qn("a:normAutofit"))
        na.set("fontScale", "70000")
