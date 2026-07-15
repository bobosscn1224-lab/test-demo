"""Background Handler — manages the background layer for PPTX reconstruction.

Strategy:
  - Solid background → native PPTX solidFill (pure, no image)
  - Gradient background → native OOXML gradientFill (pure, no image)
  - Complex background (photo/texture) → inpaint text regions, use as image
  - Complex (inpainting failed) → fallback to original image

Returns a uniform interface: apply_to_slide(slide) callable.
"""
from __future__ import annotations

import logging
import os
from pathlib import Path

import cv2
import numpy as np
from pptx.oxml.ns import qn

logger = logging.getLogger(__name__)

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5


def handle_background(
    cv_img: np.ndarray,
    pil_img,
    background_analysis: dict,
    text_regions: list[dict],
    work_dir: Path,
    normalized_path: Path,
) -> dict:
    """Determine background strategy and prepare background layer.

    Args:
        cv_img: BGR image as numpy array
        pil_img: PIL Image
        background_analysis: {"type": "solid"|"gradient"|"complex", "solid_color": "#hex", "gradient": {...}}
        text_regions: list of text region dicts with x,y,w,h
        work_dir: working directory for saving intermediate files
        normalized_path: path to normalized source image

    Returns:
        {
            "type": "native_solid" | "native_gradient" | "inpainted_image" | "original_image",
            "solid_color": "#hex" | None,
            "gradient_info": dict | None,
            "bg_image_path": str | None,
            "apply_to_slide": callable(slide) -> None,
        }
    """
    bg_type = background_analysis.get("type", "complex")

    if bg_type == "solid":
        solid_hex = background_analysis.get("solid_color", "#FFFFFF")
        return _prepare_solid_bg(solid_hex)

    if bg_type == "gradient":
        gradient_info = background_analysis.get("gradient")
        if gradient_info and gradient_info.get("stops"):
            return _prepare_gradient_bg(gradient_info)

    # Complex: try inpainting
    if bg_type == "complex" and text_regions:
        return _prepare_inpainted_bg(cv_img, text_regions, work_dir)

    # Last resort: original image
    return _prepare_original_bg(normalized_path)


# ── Solid background ──────────────────────────────────────────────────────

def _prepare_solid_bg(hex_color: str) -> dict:
    """Prepare a native solid-fill background."""
    hex_color = hex_color.lstrip("#")
    r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)

    def apply_to_slide(slide):
        _set_slide_bg(slide, r, g, b)

    return {
        "type": "native_solid",
        "solid_color": f"#{hex_color}",
        "gradient_info": None,
        "bg_image_path": None,
        "apply_to_slide": apply_to_slide,
    }


def _set_slide_bg(slide, r: int, g: int, b: int):
    """Set slide solid background color via OOXML."""
    from lxml import etree
    bg_xml = slide.background._element
    nsmap = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main",
             "a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    bgPr = bg_xml.find(".//p:bgPr", nsmap)
    if bgPr is None:
        bgPr = etree.SubElement(bg_xml, qn("p:bgPr"))
    # Remove existing fills
    for child in list(bgPr):
        if child.tag in (qn("a:solidFill"), qn("a:gradFill"), qn("a:noFill")):
            bgPr.remove(child)
    solidFill = etree.SubElement(bgPr, qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", f"{r:02x}{g:02x}{b:02x}")


# ── Gradient background ───────────────────────────────────────────────────

def _prepare_gradient_bg(gradient_info: dict) -> dict:
    """Prepare a native OOXML gradient-fill background.

    gradient_info: {"type": "linear", "angle": 90, "stops": [{"pos": "0%", "color": "#hex"}, ...]}
    """
    return {
        "type": "native_gradient",
        "solid_color": None,
        "gradient_info": gradient_info,
        "bg_image_path": None,
        "apply_to_slide": lambda slide: _apply_gradient_bg(slide, gradient_info),
    }


def _apply_gradient_bg(slide, gradient_info: dict):
    """Apply a native gradient fill to the slide background via OOXML.

    Since python-pptx has no gradient API, we construct the XML directly.
    """
    from lxml import etree

    bg_xml = slide.background._element
    bgPr = bg_xml.find(qn("p:bgPr"))
    if bgPr is None:
        bgPr = etree.SubElement(bg_xml, qn("p:bgPr"))

    # Remove existing fills
    for child in list(bgPr):
        if child.tag in (qn("a:solidFill"), qn("a:gradFill"), qn("a:noFill")):
            bgPr.remove(child)

    gradFill = etree.SubElement(bgPr, qn("a:gradFill"))
    gradFill.set("rotWithShape", "1")

    # Build gradient stop list
    gsLst = etree.SubElement(gradFill, qn("a:gsLst"))
    stops = gradient_info.get("stops", [])
    for stop in stops:
        pos_str = stop.get("pos", "0%").replace("%", "")
        pos_val = int(float(pos_str) * 1000)  # OOXML: 0..100000
        gs = etree.SubElement(gsLst, qn("a:gs"))
        gs.set("pos", str(pos_val))
        srgbClr = etree.SubElement(gs, qn("a:srgbClr"))
        color_hex = stop.get("color", "#FFFFFF").lstrip("#")
        srgbClr.set("val", color_hex)

    # Linear gradient direction
    angle = gradient_info.get("angle", 90)
    lin = etree.SubElement(gradFill, qn("a:lin"))
    # OOXML angle: 0=left-to-right, 9000000=bottom-to-top, etc.
    # Convert degrees to 1/60000 degree units
    lin.set("ang", str(int(angle * 60000)))
    lin.set("scaled", "1")


# ── Inpainted background ──────────────────────────────────────────────────

def _prepare_inpainted_bg(cv_img: np.ndarray, text_regions: list[dict],
                           work_dir: Path) -> dict:
    """Inpaint text regions and save as background image."""
    from app.services.ppt_reconstruction.analysis_engine import inpaint_text_regions

    try:
        inpainted = inpaint_text_regions(cv_img, text_regions)
        bg_path = work_dir / "bg_inpainted.png"
        cv2.imwrite(str(bg_path), inpainted)
        logger.info("Inpainted background saved: %s", bg_path)

        def apply_to_slide(slide):
            from pptx.util import Inches
            slide.shapes.add_picture(
                str(bg_path), Inches(0), Inches(0),
                Inches(SLIDE_W_IN), Inches(SLIDE_H_IN),
            )

        return {
            "type": "inpainted_image",
            "solid_color": None,
            "gradient_info": None,
            "bg_image_path": str(bg_path),
            "apply_to_slide": apply_to_slide,
        }
    except Exception as exc:
        logger.warning("Inpainting failed: %s, falling back to original image", exc)
        return _prepare_original_bg(work_dir.parent)  # fallback


# ── Original image background (fallback) ──────────────────────────────────

def _prepare_original_bg(image_path: Path) -> dict:
    """Use the original image as background (last resort fallback)."""
    def apply_to_slide(slide):
        from pptx.util import Inches
        slide.shapes.add_picture(
            str(image_path), Inches(0), Inches(0),
            Inches(SLIDE_W_IN), Inches(SLIDE_H_IN),
        )

    return {
        "type": "original_image",
        "solid_color": None,
        "gradient_info": None,
        "bg_image_path": str(image_path),
        "apply_to_slide": apply_to_slide,
    }
