"""PPT Reconstruction Pipeline — CV/OCR-driven element-by-element PPTX.

PRIMARY: CV + OCR hybrid (image_to_pptx) -> individual editable PPTX elements.
  - Text: OCR-detected -> individual editable text boxes
  - Shapes: CV-detected lines/rectangles -> native PPTX shapes
  - Images: entropy-detected logos/photos -> cropped image assets
  - Background: solid fill detected from edges

FALLBACK: visual_locked (full image bg + text overlay).
"""
from __future__ import annotations

import asyncio
import json
import logging
import os
import subprocess
import sys
import uuid
from pathlib import Path

import cv2
import numpy as np
from pptx.oxml.ns import qn

logger = logging.getLogger(__name__)

SCRIPTS_DIR = Path(__file__).parent / "scripts"


async def run_pipeline(
    image_path: str,
    session_id: str = "",
    quality_mode: str = "balanced",
    output_dir: str | None = None,
) -> dict | None:
    """Run the reconstruction pipeline.

    PRIMARY: CV/OCR hybrid — every detected element as independent editable PPTX object.
    Uses Tesseract + PaddleOCR for text, OpenCV for shapes, entropy for visual assets.
    """
    out_dir = Path(output_dir) if output_dir else Path("data/outputs")
    out_dir.mkdir(parents=True, exist_ok=True)
    sid = session_id[:8] if session_id else uuid.uuid4().hex[:8]

    # PRIMARY: CV/OCR hybrid conversion (element-by-element, no LLM needed)
    try:
        logger.info("Running CV/OCR hybrid conversion (PRIMARY)...")
        result = await _cv_ocr_hybrid_pptx(image_path, sid, out_dir)
        if result and result.get("editable_elements", 0) > 2:
            logger.info("CV/OCR hybrid complete: %d editable elements", result.get("editable_elements", 0))
            return result
        logger.warning("CV/OCR produced too few elements, trying visual_locked fallback")
    except Exception as exc:
        logger.warning("CV/OCR hybrid failed: %s, falling back", exc)

    # FALLBACK: visual_locked (image bg + text overlay)
    try:
        result = await _build_visual_locked_pptx(image_path, sid, out_dir)
        if result:
            return result
    except Exception as exc:
        logger.warning("Visual_locked fallback failed: %s", exc)

    # LAST-RESORT
    return await _fallback_basic(image_path, session_id, str(out_dir))


async def _cv_ocr_hybrid_pptx(image_path: str, sid: str, out_dir: Path) -> dict | None:
    """Primary: CV + OCR hybrid conversion -> element-by-element editable PPTX.

    Uses image_to_pptx._convert_sync which creates:
      - Layer 0: Solid background color (detected from edges)
      - Layer 1: Cropped visual assets (logos, photos, complex graphics)
      - Layer 2: Native shapes (lines, rectangles)
      - Layer 3: Individual editable text boxes with matched fonts/colors

    Each element is an independent, selectable, movable PPTX object.
    """
    import asyncio as _asyncio
    from app.services.image_to_pptx import _convert_sync

    # Run synchronous CV/OCR conversion in thread pool with timeout
    try:
        result = await _asyncio.wait_for(
            _asyncio.to_thread(_convert_sync, image_path, str(out_dir), sid),
            timeout=120.0,
        )
        if result:
            result["report"]["method"] = "CV/OCR hybrid: individual editable elements (text boxes + shapes + image crops)"
            result["report"]["quality_mode"] = "cv_ocr_hybrid"
        return result
    except _asyncio.TimeoutError:
        logger.warning("CV/OCR conversion timed out")
        return None


# Legacy helpers (kept for fallback paths)

def _preprocess_image(image_path: str, work_dir: Path) -> tuple[Path, dict]:
    from PIL import Image, ImageOps, ImageFilter
    normalized = work_dir / "normalized.png"
    img = Image.open(image_path)
    img = ImageOps.exif_transpose(img).convert("RGB")
    if img.width > 2400:
        scale = 2400 / img.width
        img = img.resize((2400, round(img.height * scale)), Image.Resampling.LANCZOS)
    img = img.filter(ImageFilter.UnsharpMask(radius=1.2, percent=120, threshold=3))
    normalized.parent.mkdir(parents=True, exist_ok=True)
    img.save(str(normalized))
    return normalized, {
        "source": image_path, "normalized": str(normalized),
        "original_width": img.width, "original_height": img.height,
        "width": img.width, "height": img.height,
        "aspect_ratio": round(img.width / img.height, 6),
    }


async def _build_visual_locked_pptx(image_path: str, sid: str, out_dir: Path) -> dict | None:
    """Fallback: 2-page PPTX with full image bg + editable text overlay."""
    import pytesseract
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    tesseract_path = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
    if os.path.exists(tesseract_path):
        pytesseract.pytesseract.tesseract_cmd = tesseract_path

    SLIDE_W = 13.333
    SLIDE_H = 7.5

    img = Image.open(image_path)
    iw, ih = img.size
    buf = np.fromfile(image_path, dtype=np.uint8)
    cv_img = cv2.imdecode(buf, cv2.IMREAD_COLOR)
    if cv_img is None:
        cv_img = np.array(img.convert("RGB"))[:, :, ::-1]

    sx = SLIDE_W / iw
    sy = SLIDE_H / ih

    from app.services.image_to_pptx import _ocr_text_regions, _detect_lines_and_rects
    text_regions = _ocr_text_regions(img, iw, ih)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]

    s1 = prs.slides.add_slide(blank)
    s1.shapes.add_picture(image_path, Inches(0), Inches(0), Inches(SLIDE_W), Inches(SLIDE_H))

    s2 = prs.slides.add_slide(blank)
    s2.shapes.add_picture(image_path, Inches(0), Inches(0), Inches(SLIDE_W), Inches(SLIDE_H))

    for tr in text_regions:
        try:
            x_in = max(0, tr["x"] * sx - 0.02)
            y_in = max(0, tr["y"] * sy - 0.02)
            w_in = min(SLIDE_W - x_in, tr["w"] * sx + 0.06)
            h_in = min(SLIDE_H - y_in, tr["h"] * sy + 0.06)
            if w_in < 0.3: w_in = 0.5
            if h_in < 0.15: h_in = 0.22
            px, py = int(tr["x"]), int(tr["y"])
            pw, ph = int(tr["w"]), int(tr["h"])
            local_bg = _sample_bg_around(cv_img, px, py, pw, ph)
            txBox = s2.shapes.add_textbox(Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in))
            tf = txBox.text_frame; tf.word_wrap = True
            _set_shape_fill(txBox, *local_bg, alpha=0.92)
            p = tf.paragraphs[0]; p.text = tr["text"]
            fs = tr.get("font_size_est", 14)
            p.font.size = Pt(max(8, min(48, fs)))
            p.font.color.rgb = RGBColor(*_sample_text_fg(cv_img, px, py, pw, ph))
            if tr.get("is_title") or fs > 22: p.font.bold = True
            rel_x = (x_in + w_in / 2) / SLIDE_W
            if rel_x < 0.3: p.alignment = PP_ALIGN.LEFT
            elif rel_x > 0.7: p.alignment = PP_ALIGN.RIGHT
            elif 0.4 < rel_x < 0.6: p.alignment = PP_ALIGN.CENTER
            _set_text_autofit(tf)
        except Exception: pass

    shapes = _detect_lines_and_rects(cv_img, iw, ih, sx, sy)
    _add_shapes_to_slide(s2, shapes[:15])

    fname = f"ppt_recon_{sid}.pptx"
    fpath = out_dir / fname
    prs.save(str(fpath))
    total_chars = sum(len(t["text"]) for t in text_regions)
    return {
        "filename": fname, "path": str(fpath),
        "url": f"http://127.0.0.1:8001/api/skills/download/{fname}",
        "pages": 2, "chars_extracted": total_chars,
        "text_regions": len(text_regions), "shapes_detected": len(shapes),
        "editable_elements": len(text_regions) + len(shapes),
        "report": {
            "editable_text_boxes": len(text_regions),
            "editable_shapes": len(shapes),
            "method": "visual_locked fallback: image bg + editable text overlay",
        },
    }


def _sample_bg_around(cv_img, x, y, w, h):
    h_img, w_img = cv_img.shape[:2]; m = 4; pixels = []
    if y > m: pixels.append(cv_img[max(0,y-m):y, x:x+w].reshape(-1, 3))
    if y+h+m < h_img: pixels.append(cv_img[y+h:min(h_img,y+h+m), x:x+w].reshape(-1, 3))
    if pixels:
        avg = np.mean(np.concatenate(pixels), axis=0).astype(int)
    else:
        roi = cv_img[y:y+h, x:x+w]
        avg = np.mean(roi.reshape(-1,3), axis=0).astype(int) if roi.size>0 else np.array([255,255,255])
    return (int(avg[2]), int(avg[1]), int(avg[0]))


def _sample_text_fg(cv_img, x, y, w, h):
    roi = cv_img[y:y+h, x:x+w]
    if roi.size == 0: return (0,0,0)
    gray = cv2.cvtColor(roi, cv2.COLOR_BGR2GRAY)
    dark = gray <= np.percentile(gray, 35)
    if dark.sum() > 3: avg = np.mean(roi[dark], axis=0).astype(int)
    else: avg = np.array([0,0,0])
    return (int(avg[2]), int(avg[1]), int(avg[0]))


def _set_shape_fill(shape, r, g, b, alpha=0.90):
    from lxml import etree
    spPr = shape._element.find(qn("p:spPr"))
    if spPr is None: spPr = etree.SubElement(shape._element, qn("p:spPr"))
    for child in list(spPr):
        if child.tag in (qn("a:solidFill"), qn("a:gradFill"), qn("a:noFill")): spPr.remove(child)
    solidFill = etree.Element(qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr")); srgbClr.set("val", f"{r:02x}{g:02x}{b:02x}")
    alpha_elem = etree.SubElement(srgbClr, qn("a:alpha")); alpha_elem.set("val", str(int(alpha*100000)))
    geom = spPr.find(qn("a:prstGeom"))
    if geom is not None: spPr.insert(list(spPr).index(geom), solidFill)
    else: spPr.append(solidFill)


def _set_text_autofit(tf):
    from lxml import etree
    bodyPr = tf._element.find(qn("a:bodyPr"))
    if bodyPr is not None:
        for tag_suffix in ("noAutofit", "normAutofit", "spAutoFit"):
            for existing in bodyPr.findall(qn(f"a:{tag_suffix}")): bodyPr.remove(existing)
        etree.SubElement(bodyPr, qn("a:normAutofit")).set("fontScale", "70000")


def _add_shapes_to_slide(slide, shapes):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    for s in shapes[:15]:
        try:
            if s["type"] == "line":
                conn = slide.shapes.add_connector(1, Inches(s["x1"]), Inches(s["y1"]), Inches(s["x2"]), Inches(s["y2"]))
                conn.line.width = Pt(s.get("width_pt", 1.0)); conn.line.color.rgb = RGBColor(0x99, 0x99, 0x99)
            elif s["type"] == "rectangle":
                shape = slide.shapes.add_shape(1, Inches(s["x"]), Inches(s["y"]), Inches(s["w"]), Inches(s["h"]))
                shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(*s.get("fill", (240,240,240)))
                shape.line.color.rgb = RGBColor(*s.get("stroke", (200,200,200))); shape.line.width = Pt(0.5)
        except Exception: pass


async def _fallback_basic(image_path, session_id, output_dir):
    from app.services.image_to_pptx import convert_image_to_pptx
    return await convert_image_to_pptx(image_path, str(output_dir) if output_dir else None, session_id)
