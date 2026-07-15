#!/usr/bin/env python3
"""Fallback converter: SVG subset → editable PPTX.

This parses basic SVG elements (rect, circle, line, text, image) and maps them to
native PowerPoint elements. Complex paths/filters/masks are not faithfully
converted; use plan_to_pptx.py for production.
"""
from __future__ import annotations

import argparse
import base64
import re
import tempfile
from pathlib import Path
from typing import Dict, Tuple

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

SVG_NS = "{http://www.w3.org/2000/svg}"
XLINK = "{http://www.w3.org/1999/xlink}href"


def strip_ns(tag: str) -> str:
    return tag.split("}", 1)[-1] if "}" in tag else tag


def parse_num(value, default=0.0) -> float:
    if value is None:
        return default
    m = re.match(r"[-+]?\d*\.?\d+", str(value))
    return float(m.group(0)) if m else default


def parse_color(value: str | None, default=(0, 0, 0)) -> RGBColor:
    if not value or value == "none":
        return RGBColor(*default)
    value = value.strip()
    if value.startswith("#"):
        value = value[1:]
    if len(value) == 3:
        value = "".join(c * 2 for c in value)
    try:
        return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))
    except Exception:
        return RGBColor(*default)


def get_style(el) -> Dict[str, str]:
    style: Dict[str, str] = {}
    raw = el.get("style")
    if raw:
        for part in raw.split(";"):
            if ":" in part:
                k, v = part.split(":", 1)
                style[k.strip()] = v.strip()
    for k, v in el.attrib.items():
        if strip_ns(k) in ["fill", "stroke", "stroke-width", "font-size", "font-family", "font-weight", "text-anchor"]:
            style[strip_ns(k)] = v
    return style


class Map:
    def __init__(self, svg_w: float, svg_h: float, slide_w_in: float = 13.333333):
        self.svg_w = svg_w
        self.svg_h = svg_h
        self.slide_w = Inches(slide_w_in)
        self.slide_h = int(self.slide_w * svg_h / svg_w)
        self.pt_per_px = (slide_w_in * 72) / svg_w

    def x(self, v): return int(self.slide_w * parse_num(v) / self.svg_w)
    def y(self, v): return int(self.slide_h * parse_num(v) / self.svg_h)
    def w(self, v): return int(self.slide_w * parse_num(v) / self.svg_w)
    def h(self, v): return int(self.slide_h * parse_num(v) / self.svg_h)
    def pt(self, v): return Pt(max(1, parse_num(v, 24) * self.pt_per_px))


def decode_data_uri(uri: str, tmp: Path) -> Path | None:
    if not uri.startswith("data:") or ";base64," not in uri:
        return None
    header, b64 = uri.split(",", 1)
    ext = ".png"
    if "jpeg" in header or "jpg" in header:
        ext = ".jpg"
    elif "webp" in header:
        ext = ".webp"
    path = tmp / f"embedded_{abs(hash(uri))}{ext}"
    path.write_bytes(base64.b64decode(b64))
    return path


def text_content(el) -> str:
    parts = []
    if el.text:
        parts.append(el.text)
    for child in el:
        if child.text:
            parts.append(child.text)
        if child.tail:
            parts.append(child.tail)
    return "".join(parts).strip()


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("svg", type=Path)
    parser.add_argument("--out", type=Path, default=Path("work/reconstructed.pptx"))
    parser.add_argument("--slide-width", type=float, default=13.333333)
    args = parser.parse_args()

    root = etree.parse(str(args.svg)).getroot()
    viewbox = root.get("viewBox")
    if viewbox:
        nums = [float(x) for x in viewbox.replace(",", " ").split()]
        svg_w, svg_h = nums[2], nums[3]
    else:
        svg_w = parse_num(root.get("width"), 1920)
        svg_h = parse_num(root.get("height"), 1080)
    m = Map(svg_w, svg_h, args.slide_width)

    prs = Presentation()
    prs.slide_width = m.slide_w
    prs.slide_height = m.slide_h
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tmp = Path(tempfile.mkdtemp(prefix="svg_pptx_"))
    for el in root.iter():
        tag = strip_ns(el.tag)
        style = get_style(el)
        if tag == "rect":
            rx = parse_num(el.get("rx"), 0)
            shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rx > 0 else MSO_SHAPE.RECTANGLE
            shape = slide.shapes.add_shape(shape_type, m.x(el.get("x")), m.y(el.get("y")), m.w(el.get("width")), m.h(el.get("height")))
            fill = style.get("fill", "#FFFFFF")
            if fill == "none":
                shape.fill.background()
            else:
                shape.fill.solid(); shape.fill.fore_color.rgb = parse_color(fill, (255,255,255))
            stroke = style.get("stroke", "none")
            if stroke == "none":
                shape.line.fill.background()
            else:
                shape.line.color.rgb = parse_color(stroke)
                shape.line.width = m.pt(style.get("stroke-width", 1))
        elif tag == "circle":
            cx, cy, r = parse_num(el.get("cx")), parse_num(el.get("cy")), parse_num(el.get("r"))
            shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, m.x(cx-r), m.y(cy-r), m.w(2*r), m.h(2*r))
            fill = style.get("fill", "#FFFFFF")
            if fill == "none": shape.fill.background()
            else: shape.fill.solid(); shape.fill.fore_color.rgb = parse_color(fill, (255,255,255))
            stroke = style.get("stroke", "none")
            if stroke == "none": shape.line.fill.background()
            else: shape.line.color.rgb = parse_color(stroke); shape.line.width = m.pt(style.get("stroke-width", 1))
        elif tag == "line":
            shape = slide.shapes.add_connector(1, m.x(el.get("x1")), m.y(el.get("y1")), m.x(el.get("x2")), m.y(el.get("y2")))
            shape.line.color.rgb = parse_color(style.get("stroke", "#000000"))
            shape.line.width = m.pt(style.get("stroke-width", 1))
        elif tag == "text":
            txt = text_content(el)
            if not txt:
                continue
            x = parse_num(el.get("x")); y = parse_num(el.get("y"))
            font_size = parse_num(style.get("font-size"), 24)
            box = slide.shapes.add_textbox(m.x(x), m.y(y-font_size), m.w(svg_w - x), m.h(font_size * 2.4))
            tf = box.text_frame; tf.clear(); tf.text = txt
            anchor = style.get("text-anchor", "start")
            align = {"start": PP_ALIGN.LEFT, "middle": PP_ALIGN.CENTER, "end": PP_ALIGN.RIGHT}.get(anchor, PP_ALIGN.LEFT)
            for p in tf.paragraphs:
                p.alignment = align
                for run in p.runs:
                    run.font.name = style.get("font-family", "Microsoft YaHei")
                    run.font.size = m.pt(font_size)
                    run.font.bold = parse_num(style.get("font-weight"), 400) >= 600
                    run.font.color.rgb = parse_color(style.get("fill", "#000000"))
        elif tag == "image":
            href = el.get("href") or el.get(XLINK)
            if not href:
                continue
            img_path = decode_data_uri(href, tmp) or Path(href)
            if img_path.exists():
                slide.shapes.add_picture(str(img_path), m.x(el.get("x")), m.y(el.get("y")), m.w(el.get("width")), m.h(el.get("height")))

    args.out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.out)
    print(str(args.out))


if __name__ == "__main__":
    main()
