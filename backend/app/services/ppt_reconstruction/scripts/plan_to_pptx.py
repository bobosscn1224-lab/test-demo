#!/usr/bin/env python3
"""Generate an editable PPTX from layout_plan.json.

This is the preferred converter because the plan preserves semantics better than
arbitrary SVG. It maps supported plan elements to native PowerPoint objects and
uses image assets for complex regions.

Element types supported:
  rect, line, circle, arrow, polygon, text, image, table

Layer-based assembly ensures correct z-ordering:
  Layer 0 (bg): background rect/gradient/image
  Layer 1 (images): logos, photos, cropped assets
  Layer 2 (shapes): rects, circles, lines, arrows, polygons
  Layer 3 (text): titles, body text, labels, tables

Also callable from Python as: build_pptx_from_plan(plan, assets_dir, output_path)
"""
from __future__ import annotations

import argparse
import json
import os
from pathlib import Path
from typing import Any, Callable, Dict, Tuple

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def hex_to_rgb(value: str, default: Tuple[int, int, int] = (0, 0, 0)) -> RGBColor:
    if not value or value == "none":
        return RGBColor(*default)
    value = str(value).strip()
    if value.startswith("#"):
        value = value[1:]
    if len(value) == 3:
        value = "".join(ch * 2 for ch in value)
    try:
        return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))
    except Exception:
        return RGBColor(*default)


def resolve_font(font_family: str) -> str:
    """Map font family to best available equivalent on Windows."""
    font_map = {
        "Microsoft YaHei": "Microsoft YaHei",
        "SimHei": "SimHei",
        "SimSun": "SimSun",
        "PingFang SC": "Microsoft YaHei",
        "PingFang HK": "Microsoft YaHei",
        "Noto Sans CJK SC": "Microsoft YaHei",
        "Noto Sans CJK": "Microsoft YaHei",
        "Source Han Sans SC": "Microsoft YaHei",
        "Source Han Sans": "Microsoft YaHei",
        "Hiragino Sans GB": "Microsoft YaHei",
        "STHeiti": "SimHei",
        "STSong": "SimSun",
        "KaiTi": "KaiTi",
        "FangSong": "FangSong",
    }
    return font_map.get(font_family, font_family)


def set_element_name(shape, name: str):
    """Give a shape a readable name for identification in PowerPoint."""
    try:
        shape.name = name[:255]
    except Exception:
        pass


class Mapper:
    def __init__(self, canvas_w: float, canvas_h: float, slide_w_in: float = 13.333333):
        self.canvas_w = canvas_w
        self.canvas_h = canvas_h
        self.slide_w = Inches(slide_w_in)
        self.slide_h = int(self.slide_w * canvas_h / canvas_w)
        self.pt_per_px = (slide_w_in * 72.0) / canvas_w

    def x(self, px: float):
        return int(self.slide_w * float(px) / self.canvas_w)

    def y(self, px: float):
        return int(self.slide_h * float(px) / self.canvas_h)

    def w(self, px: float):
        return int(self.slide_w * float(px) / self.canvas_w)

    def h(self, px: float):
        return int(self.slide_h * float(px) / self.canvas_h)

    def pt(self, px: float):
        # The factor makes image-pixel typography visually close on standard 16:9 decks.
        return Pt(max(1, float(px) * self.pt_per_px))


def apply_fill(shape, fill_value: str | None):
    if fill_value is None or fill_value == "none":
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill_value, (255, 255, 255))


def apply_line(shape, stroke: str | None, stroke_width_pt: float | None):
    if stroke is None or stroke == "none":
        shape.line.fill.background()
    else:
        shape.line.color.rgb = hex_to_rgb(stroke)
        if stroke_width_pt is not None:
            shape.line.width = Pt(max(0.25, stroke_width_pt))


def add_rect(slide, m: Mapper, el: Dict[str, Any]):
    x, y, w, h = m.x(el.get("x", 0)), m.y(el.get("y", 0)), m.w(el.get("w", 0)), m.h(el.get("h", 0))
    radius = float(el.get("rx", 0) or 0)
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius > 0 else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    apply_fill(shape, el.get("fill", "#FFFFFF"))
    apply_line(shape, el.get("stroke", "none"), float(el.get("stroke_width", 0) or 0) * m.pt_per_px)
    return shape


def add_line(slide, m: Mapper, el: Dict[str, Any]):
    shape = slide.shapes.add_connector(
        1,
        m.x(el.get("x1", 0)),
        m.y(el.get("y1", 0)),
        m.x(el.get("x2", 0)),
        m.y(el.get("y2", 0)),
    )
    apply_line(shape, el.get("stroke", "#000000"), float(el.get("stroke_width", 1) or 1) * m.pt_per_px)
    return shape


def add_circle(slide, m: Mapper, el: Dict[str, Any]):
    cx = float(el.get("cx", el.get("x", 0)))
    cy = float(el.get("cy", el.get("y", 0)))
    r = float(el.get("r", 0))
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, m.x(cx - r), m.y(cy - r), m.w(2 * r), m.h(2 * r))
    apply_fill(shape, el.get("fill", "#FFFFFF"))
    apply_line(shape, el.get("stroke", "none"), float(el.get("stroke_width", 0) or 0) * m.pt_per_px)
    return shape


def add_text(slide, m: Mapper, el: Dict[str, Any]):
    box = slide.shapes.add_textbox(m.x(el.get("x", 0)), m.y(el.get("y", 0)), m.w(el.get("w", 1)), m.h(el.get("h", 1)))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    valign = str(el.get("valign", "top"))
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}.get(valign, MSO_ANCHOR.TOP)
    lines = str(el.get("text", "")).split("\n")
    p = tf.paragraphs[0]
    p.text = lines[0] if lines else ""
    for extra in lines[1:]:
        p = tf.add_paragraph()
        p.text = extra
    align = str(el.get("align", "left"))
    align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}
    for p in tf.paragraphs:
        p.alignment = align_map.get(align, PP_ALIGN.LEFT)
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        for run in p.runs:
            run.font.name = resolve_font(el.get("font_family", "Microsoft YaHei"))
            run.font.size = m.pt(float(el.get("font_size", 24)))
            run.font.bold = int(el.get("font_weight", 400) or 400) >= 600
            run.font.color.rgb = hex_to_rgb(el.get("color", "#000000"))
    return box


def add_image(slide, m: Mapper, el: Dict[str, Any], assets_dir: Path, asset_map: Dict[str, Any]):
    asset = asset_map.get(el.get("asset_id"), {})
    file = asset.get("file") or el.get("file")
    if not file:
        return None
    path = assets_dir / file
    if not path.exists():
        # Also allow absolute or plan-relative file paths.
        path = Path(file)
    if not path.exists():
        return None
    return slide.shapes.add_picture(str(path), m.x(el.get("x", 0)), m.y(el.get("y", 0)), m.w(el.get("w", 0)), m.h(el.get("h", 0)))


def add_table(slide, m: Mapper, el: Dict[str, Any]):
    rows = int(el.get("rows", 1))
    cols = int(el.get("cols", 1))
    shape = slide.shapes.add_table(rows, cols, m.x(el.get("x", 0)), m.y(el.get("y", 0)), m.w(el.get("w", 1)), m.h(el.get("h", 1)))
    table = shape.table
    cell_text = el.get("cell_text", [])
    for r in range(rows):
        for c in range(cols):
            text = ""
            if r < len(cell_text) and c < len(cell_text[r]):
                text = str(cell_text[r][c])
            cell = table.cell(r, c)
            cell.text = text
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = resolve_font(el.get("font_family", "Microsoft YaHei"))
                    run.font.size = m.pt(float(el.get("font_size", 18)))
                    run.font.color.rgb = hex_to_rgb(el.get("color", "#111111"))
    return shape


# ── Arrow element ──────────────────────────────────────────────────────────

def add_arrow(slide, m: Mapper, el: Dict[str, Any]):
    """Add an arrow connector, optionally with arrowheads."""
    shape = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR.STRAIGHT
        m.x(el.get("x1", 0)), m.y(el.get("y1", 0)),
        m.x(el.get("x2", 0)), m.y(el.get("y2", 0)),
    )
    apply_line(shape, el.get("stroke", "#333333"),
               float(el.get("stroke_width", 1.5) or 1.5) * m.pt_per_px)

    # Arrowheads via OOXML
    head_end = el.get("arrow_head", "end")
    if head_end in ("end", "both"):
        _add_arrowhead_oxml(shape, is_start=False)
    if head_end in ("start", "both"):
        _add_arrowhead_oxml(shape, is_start=True)
    return shape


def _add_arrowhead_oxml(shape, is_start: bool = False):
    """Add triangle arrowhead to a connector via OOXML."""
    from lxml import etree
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    spPr = shape._element.find(qn("p:spPr"))
    if spPr is None:
        return
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = etree.SubElement(spPr, qn("a:ln"))
    tag = qn("a:tailEnd") if is_start else qn("a:headEnd")
    # Remove existing head/tail
    for existing in ln.findall(tag):
        ln.remove(existing)
    head = etree.SubElement(ln, tag)
    head.set("type", "triangle")
    head.set("w", "med")
    head.set("len", "med")


# ── Polygon element ───────────────────────────────────────────────────────

def add_polygon(slide, m: Mapper, el: Dict[str, Any]):
    """Add a freeform polygon shape from [x,y] point array."""
    points = el.get("points", [])
    if not points or len(points) < 3:
        return None

    # Map points to EMU offsets relative to shape position
    min_x = min(p[0] for p in points)
    min_y = min(p[1] for p in points)
    max_x = max(p[0] for p in points)
    max_y = max(p[1] for p in points)

    # Use FREEFORM shape, then adjust path via OOXML
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,  # placeholder; we'll replace the geometry
        m.x(min_x), m.y(min_y),
        m.w(max_x - min_x), m.h(max_y - min_y),
    )
    apply_fill(shape, el.get("fill", "#FFFFFF"))
    apply_line(shape, el.get("stroke", "#333333"),
               float(el.get("stroke_width", 1.0) or 1.0) * m.pt_per_px)

    # Build custom path geometry
    _replace_shape_path(shape, points, min_x, min_y, max_x - min_x, max_y - min_y, m)
    return shape


def _replace_shape_path(shape, points: list, min_x: float, min_y: float,
                          w: float, h: float, m: Mapper):
    """Replace shape geometry with a custom freeform path."""
    from lxml import etree

    spPr = shape._element.find(qn("p:spPr"))
    if spPr is None:
        return

    # Remove existing geometry
    for child in list(spPr):
        if child.tag in (qn("a:prstGeom"), qn("a:custGeom")):
            spPr.remove(child)

    # Build custom geometry
    custGeom = etree.SubElement(spPr, qn("a:custGeom"))
    avLst = etree.SubElement(custGeom, qn("a:avLst"))

    # Path list
    pathLst = etree.SubElement(custGeom, qn("a:pathLst"))
    path = etree.SubElement(pathLst, qn("a:path"))
    path.set("w", str(int(m.w(w))))
    path.set("h", str(int(m.h(h))))

    # Build moveto + lineto commands
    cmds = []
    for i, pt in enumerate(points):
        rel_x = int(m.x(pt[0] - min_x))
        rel_y = int(m.y(pt[1] - min_y))
        if i == 0:
            cmd = etree.SubElement(path, qn("a:moveTo"))
        else:
            cmd = etree.SubElement(path, qn("a:lnTo"))
        pt_elem = etree.SubElement(cmd, qn("a:pt"))
        pt_elem.set("x", str(rel_x))
        pt_elem.set("y", str(rel_y))

    # Close path
    close = etree.SubElement(path, qn("a:close"))


# ── Gradient fill ────────────────────────────────────────────────────────

def apply_gradient_fill(shape, gradient_info: dict):
    """Apply a native gradient fill to a shape via OOXML."""
    from lxml import etree

    spPr = shape._element.find(qn("p:spPr"))
    if spPr is None:
        return

    # Remove existing fills
    for child in list(spPr):
        tag = etree.QName(child).localname if hasattr(child, 'tag') else ''
        if child.tag in (qn("a:solidFill"), qn("a:gradFill"), qn("a:noFill")):
            spPr.remove(child)

    gradFill = etree.SubElement(spPr, qn("a:gradFill"))
    gradFill.set("rotWithShape", "1")

    # Color stops
    gsLst = etree.SubElement(gradFill, qn("a:gsLst"))
    stops = gradient_info.get("stops", [])
    for stop in stops:
        pos_str = str(stop.get("pos", "0%")).replace("%", "")
        pos_val = int(float(pos_str) * 1000)
        gs = etree.SubElement(gsLst, qn("a:gs"))
        gs.set("pos", str(pos_val))
        srgbClr = etree.SubElement(gs, qn("a:srgbClr"))
        hex_val = stop.get("color", "#FFFFFF").lstrip("#")
        srgbClr.set("val", hex_val)

    # Linear gradient
    angle = gradient_info.get("angle", 90)
    lin = etree.SubElement(gradFill, qn("a:lin"))
    lin.set("ang", str(int(angle * 60000)))
    lin.set("scaled", "1")


# ── Main entry points (CLI + Python API) ──────────────────────────────────

def build_pptx_from_plan(
    plan: dict,
    assets_dir: Path | str = Path("work/assets"),
    output_path: Path | str | None = None,
    slide_width: float = 13.333333,
    bg_handler: dict | None = None,
) -> Presentation:
    """Python API: Build an element-by-element PPTX from a layout plan.

    Args:
        plan: The layout_plan.json data
        assets_dir: Path to cropped image assets
        output_path: If provided, save to this path
        slide_width: Slide width in inches (default 16:9)
        bg_handler: Optional background handler result dict with apply_to_slide callable

    Returns:
        python-pptx Presentation object
    """
    assets_dir = Path(assets_dir)
    canvas = plan["canvas"]
    m = Mapper(float(canvas["width"]), float(canvas["height"]), slide_width)

    prs = Presentation()
    prs.slide_width = m.slide_w
    prs.slide_height = m.slide_h
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ── Layer 0: Background ──
    if bg_handler and callable(bg_handler.get("apply_to_slide")):
        bg_handler["apply_to_slide"](slide)
    else:
        bg_color = canvas.get("background", "#FFFFFF")
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        apply_fill(bg, bg_color)
        apply_line(bg, "none", None)
        set_element_name(bg, "Background")

    asset_map = {a.get("id"): a for a in plan.get("assets", [])}
    elements = plan.get("elements", [])

    # ── Layer-based sorting ──
    def _layer_for(el: dict) -> int:
        t = el.get("type")
        if t == "image" or el.get("editability") == "asset":
            return 1
        if t in ("rect", "line", "circle", "arrow", "polygon"):
            return 2
        if t in ("text", "table"):
            return 3
        return 2

    # Sort by layer first, then by z within layer
    elements_sorted = sorted(elements, key=lambda e: (_layer_for(e), e.get("z", 0)))

    counts = {"rect": 0, "line": 0, "circle": 0, "arrow": 0, "polygon": 0,
              "text": 0, "image": 0, "table": 0}

    for el in elements_sorted:
        el_id = el.get("id", "element")
        t = el.get("type")
        try:
            if t == "rect":
                shape = add_rect(slide, m, el)
                if shape:
                    set_element_name(shape, f"Rect_{el_id}")
                    counts["rect"] += 1
            elif t == "line":
                shape = add_line(slide, m, el)
                if shape:
                    set_element_name(shape, f"Line_{el_id}")
                    counts["line"] += 1
            elif t == "circle":
                shape = add_circle(slide, m, el)
                if shape:
                    set_element_name(shape, f"Circle_{el_id}")
                    counts["circle"] += 1
            elif t == "arrow":
                shape = add_arrow(slide, m, el)
                if shape:
                    set_element_name(shape, f"Arrow_{el_id}")
                    counts["arrow"] += 1
            elif t == "polygon":
                shape = add_polygon(slide, m, el)
                if shape:
                    set_element_name(shape, f"Polygon_{el_id}")
                    counts["polygon"] += 1
            elif t == "text":
                shape = add_text(slide, m, el)
                if shape:
                    set_element_name(shape, f"Text_{el_id}")
                    counts["text"] += 1
            elif t == "image":
                shape = add_image(slide, m, el, assets_dir, asset_map)
                if shape:
                    set_element_name(shape, f"Image_{el_id}")
                    counts["image"] += 1
            elif t == "table":
                shape = add_table(slide, m, el)
                if shape:
                    set_element_name(shape, f"Table_{el_id}")
                    counts["table"] += 1
        except Exception as exc:
            logger.warning("Failed to place element %s: %s", el_id, exc)

    total = sum(counts.values())
    logger.info("PPTX built: %d elements (%s)", total,
                ", ".join(f"{k}={v}" for k, v in counts.items() if v > 0))

    if output_path:
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        logger.info("PPTX saved: %s", output_path)

    return prs


import logging
logger = logging.getLogger(__name__)


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("plan", type=Path)
    parser.add_argument("--assets", type=Path, default=Path("work/assets"))
    parser.add_argument("--out", type=Path, default=Path("work/reconstructed.pptx"))
    parser.add_argument("--slide-width", type=float, default=13.333333)
    args = parser.parse_args()

    plan = json.loads(args.plan.read_text(encoding="utf-8"))
    build_pptx_from_plan(plan, assets_dir=args.assets, output_path=args.out,
                         slide_width=args.slide_width)
    print(str(args.out))


if __name__ == "__main__":
    main()
