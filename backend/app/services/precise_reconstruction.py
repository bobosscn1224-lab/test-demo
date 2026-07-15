"""High-fidelity image-to-editable-PPTX reconstruction for PPT Maker step 4.

The service deliberately keeps the existing Claude/DeepSeek backend:

* OpenCV + OCR measure the source image locally.
* DeepSeek receives OCR/CV metadata only and may correct/group text.  It is
  never treated as a vision model and failure is non-fatal.
* Complex visual regions are preserved as multiple text-cleaned image crops.
  A full-slide bitmap is never used in the reconstructed slide.
* Main text is rebuilt as native PowerPoint text boxes.
* High-confidence basic geometry is added as native PowerPoint shapes.
* The exported PPTX is rendered by local Microsoft PowerPoint, compared with
  the source, and the best typography calibration is retained.
"""
from __future__ import annotations

import asyncio
import difflib
import json
import logging
import math
import os
import re
import shutil
import subprocess
import sys
import tempfile
import threading
import time
import uuid
from dataclasses import dataclass
from pathlib import Path
from typing import Any

import cv2
import numpy as np
from PIL import Image

logger = logging.getLogger(__name__)
PIPELINE_VERSION = "precise_v3_native_ppt_objects"

SLIDE_W_IN = 13.333333
SLIDE_H_IN = 7.5
SLIDE_W_PT = 960.0
SLIDE_H_PT = 540.0
DEFAULT_RENDER_W = 1600
DEFAULT_RENDER_H = 900
TESSERACT_PATHS = (
    r"C:\Program Files\Tesseract-OCR\tesseract.exe",
    r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
)
POWERPOINT_PATHS = (
    r"C:\Program Files\Microsoft Office\Office15\POWERPNT.EXE",
    r"C:\Program Files\Microsoft Office\root\Office16\POWERPNT.EXE",
    r"C:\Program Files (x86)\Microsoft Office\Office15\POWERPNT.EXE",
    r"C:\Program Files (x86)\Microsoft Office\root\Office16\POWERPNT.EXE",
)
_CJK_RE = re.compile(r"[\u3400-\u9fff]")
_CONTENT_RE = re.compile(r"[\u3400-\u9fffA-Za-z0-9]")
_PUNCT_ONLY_RE = re.compile(r"^[\s|_—\-~`'\".,，。:：;；()（）<>《》@©]+$")
_KNOWN_SHORT_LATIN = {
    "MO", "BG", "BU", "BP", "QBC", "UI", "UX", "AI", "IT", "OKR", "KPI",
    "ROI", "B2B", "B2C", "SLA", "CRM",
}
_PADDLE_OCR = None
_PADDLE_LOCK = threading.Lock()
DEFAULT_CJK_FONT = os.getenv("PPT_PRECISE_CJK_FONT", "Microsoft YaHei UI")
DEFAULT_LATIN_FONT = os.getenv("PPT_PRECISE_LATIN_FONT", "Arial")
DEFAULT_NUMBER_FONT = os.getenv("PPT_PRECISE_NUMBER_FONT", "Arial")

from ._paths import OUTPUTS_DIR as _DATA_OUTPUT_DIR


def _default_output_dir() -> Path:
    _DATA_OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    return _DATA_OUTPUT_DIR


@dataclass
class CoordinateMapper:
    """Map source pixels to a 16:9 PowerPoint canvas.

    Near-16:9 images are stretched by the tiny residual amount, matching how
    visual mockups are normally fitted to a widescreen slide.  Materially
    different aspect ratios are letterboxed without distortion.
    """

    image_width: int
    image_height: int
    aspect_tolerance: float = 0.035

    def __post_init__(self) -> None:
        source_aspect = self.image_width / max(self.image_height, 1)
        target_aspect = SLIDE_W_IN / SLIDE_H_IN
        relative_delta = abs(source_aspect - target_aspect) / target_aspect
        if relative_delta <= self.aspect_tolerance:
            self.mode = "stretch"
            self.sx = SLIDE_W_IN / self.image_width
            self.sy = SLIDE_H_IN / self.image_height
            self.ox = 0.0
            self.oy = 0.0
        else:
            self.mode = "contain"
            scale = min(SLIDE_W_IN / self.image_width, SLIDE_H_IN / self.image_height)
            self.sx = self.sy = scale
            self.ox = (SLIDE_W_IN - self.image_width * scale) / 2
            self.oy = (SLIDE_H_IN - self.image_height * scale) / 2

    def x(self, value: float) -> float:
        return self.ox + value * self.sx

    def y(self, value: float) -> float:
        return self.oy + value * self.sy

    def w(self, value: float) -> float:
        return value * self.sx

    def h(self, value: float) -> float:
        return value * self.sy

    def rect(self, x: float, y: float, w: float, h: float) -> tuple[float, float, float, float]:
        return self.x(x), self.y(y), self.w(w), self.h(h)

    def font_points(self, pixel_height: float, font_scale: float = 1.0) -> float:
        # At 120-DPI PowerPoint export, the mapped source glyph height in
        # points is already a close match for Microsoft YaHei.
        mapped_points = pixel_height * self.sy * 72.0
        return max(6.5, min(48.0, mapped_points * font_scale))

    def render_rect(
        self,
        x: float,
        y: float,
        w: float,
        h: float,
        render_w: int = DEFAULT_RENDER_W,
        render_h: int = DEFAULT_RENDER_H,
    ) -> tuple[int, int, int, int]:
        x_in, y_in, w_in, h_in = self.rect(x, y, w, h)
        return (
            round(x_in / SLIDE_W_IN * render_w),
            round(y_in / SLIDE_H_IN * render_h),
            max(1, round(w_in / SLIDE_W_IN * render_w)),
            max(1, round(h_in / SLIDE_H_IN * render_h)),
        )


def _read_cv_image(image_path: str) -> np.ndarray:
    data = np.fromfile(image_path, dtype=np.uint8)
    image = cv2.imdecode(data, cv2.IMREAD_COLOR)
    if image is None:
        raise ValueError(f"Cannot read image: {image_path}")
    return image


def _find_tesseract() -> str | None:
    configured = os.getenv("TESSERACT_CMD", "").strip()
    if configured and os.path.exists(configured):
        return configured
    for candidate in TESSERACT_PATHS:
        if os.path.exists(candidate):
            return candidate
    return shutil.which("tesseract")


def _contains_cjk(text: str) -> bool:
    return bool(_CJK_RE.search(text))


def _normalize_text_for_similarity(text: str) -> str:
    return re.sub(r"[^\u3400-\u9fffA-Za-z0-9]", "", text).lower()


def _text_similarity(a: str, b: str) -> float:
    aa = _normalize_text_for_similarity(a)
    bb = _normalize_text_for_similarity(b)
    if not aa or not bb:
        return 0.0
    return difflib.SequenceMatcher(None, aa, bb).ratio()


def _iou(a: dict[str, Any], b: dict[str, Any]) -> float:
    x1 = max(a["x"], b["x"])
    y1 = max(a["y"], b["y"])
    x2 = min(a["x"] + a["w"], b["x"] + b["w"])
    y2 = min(a["y"] + a["h"], b["y"] + b["h"])
    if x2 <= x1 or y2 <= y1:
        return 0.0
    inter = (x2 - x1) * (y2 - y1)
    union = a["w"] * a["h"] + b["w"] * b["h"] - inter
    return inter / max(union, 1)


def _valid_ocr_token(text: str, confidence: float) -> bool:
    text = text.strip()
    if not text or not _CONTENT_RE.search(text):
        return False
    if _PUNCT_ONLY_RE.match(text):
        return False
    if len(text) == 1 and text in {"I", "l", "O", "Q"} and confidence < 85:
        return False
    if _contains_cjk(text) and len(text) == 1 and confidence < 52:
        return False
    ascii_core = re.sub(r"[^A-Za-z0-9]", "", text)
    if (
        ascii_core
        and any(char.isalpha() for char in ascii_core)
        and not _contains_cjk(text)
        and ascii_core.upper() not in _KNOWN_SHORT_LATIN
        and confidence < 76
    ):
        return False
    if text.isascii() and text.isalpha() and len(text) <= 3:
        if text.upper() not in _KNOWN_SHORT_LATIN and confidence < 82:
            return False
    return confidence >= 27


def _sanitize_ocr_token(text: str) -> str:
    stripped = text.strip()
    ascii_core = re.sub(r"[^A-Za-z0-9]", "", stripped)
    if ascii_core.upper() in _KNOWN_SHORT_LATIN:
        return ascii_core.upper()
    return stripped


def _sample_background_bgr(cv_img: np.ndarray, x: int, y: int, w: int, h: int) -> np.ndarray:
    ih, iw = cv_img.shape[:2]
    margin = max(2, min(8, round(min(w, h) * 0.16)))
    samples: list[np.ndarray] = []
    if y >= margin:
        samples.append(cv_img[y - margin:y, max(0, x):min(iw, x + w)].reshape(-1, 3))
    if y + h + margin <= ih:
        samples.append(cv_img[y + h:y + h + margin, max(0, x):min(iw, x + w)].reshape(-1, 3))
    if x >= margin:
        samples.append(cv_img[max(0, y):min(ih, y + h), x - margin:x].reshape(-1, 3))
    if x + w + margin <= iw:
        samples.append(cv_img[max(0, y):min(ih, y + h), x + w:x + w + margin].reshape(-1, 3))
    samples = [sample for sample in samples if sample.size]
    if not samples:
        return np.array([255, 255, 255], dtype=np.float32)
    pixels = np.concatenate(samples, axis=0)
    return np.median(pixels, axis=0)


def _sample_foreground_hex(cv_img: np.ndarray, x: int, y: int, w: int, h: int) -> str:
    ih, iw = cv_img.shape[:2]
    x1, y1 = max(0, x), max(0, y)
    x2, y2 = min(iw, x + w), min(ih, y + h)
    roi = cv_img[y1:y2, x1:x2]
    if roi.size == 0:
        return "#111111"
    bg = _sample_background_bgr(cv_img, x1, y1, x2 - x1, y2 - y1)
    distance = np.linalg.norm(roi.astype(np.float32) - bg.reshape(1, 1, 3), axis=2)
    gray = cv2.cvtColor(roi, cv2.COLOR_BGR2GRAY)
    threshold = max(25.0, float(np.percentile(distance, 72)))
    mask = (distance >= threshold) & (gray < 235)
    if mask.sum() < 4:
        mask = gray <= np.percentile(gray, 22)
    pixels = roi[mask]
    if pixels.size == 0:
        return "#111111"
    bgr = np.median(pixels, axis=0).astype(int)
    return f"#{bgr[2]:02X}{bgr[1]:02X}{bgr[0]:02X}"


def _bgr_to_hex(bgr: np.ndarray | tuple[int, int, int] | list[int]) -> str:
    arr = np.asarray(bgr).astype(int)
    return f"#{arr[2]:02X}{arr[1]:02X}{arr[0]:02X}"


def _sample_fill_hex(cv_img: np.ndarray, x: int, y: int, w: int, h: int) -> str:
    ih, iw = cv_img.shape[:2]
    inset_x = max(2, min(18, round(w * 0.18)))
    inset_y = max(2, min(14, round(h * 0.18)))
    x1 = max(0, min(iw, x + inset_x))
    y1 = max(0, min(ih, y + inset_y))
    x2 = max(x1 + 1, min(iw, x + w - inset_x))
    y2 = max(y1 + 1, min(ih, y + h - inset_y))
    roi = cv_img[y1:y2, x1:x2]
    if roi.size == 0:
        return "#FFFFFF"
    return _bgr_to_hex(np.median(roi.reshape(-1, 3), axis=0))


def _sample_stroke_hex(cv_img: np.ndarray, x: int, y: int, w: int, h: int) -> str:
    ih, iw = cv_img.shape[:2]
    thickness = max(2, min(6, round(min(w, h) * 0.10)))
    x1, y1 = max(0, x), max(0, y)
    x2, y2 = min(iw, x + w), min(ih, y + h)
    if x2 <= x1 or y2 <= y1:
        return "#A9C4C7"
    bands = [
        cv_img[y1:min(y2, y1 + thickness), x1:x2],
        cv_img[max(y1, y2 - thickness):y2, x1:x2],
        cv_img[y1:y2, x1:min(x2, x1 + thickness)],
        cv_img[y1:y2, max(x1, x2 - thickness):x2],
    ]
    pixels = np.concatenate([band.reshape(-1, 3) for band in bands if band.size], axis=0)
    if pixels.size == 0:
        return "#A9C4C7"
    gray = 0.114 * pixels[:, 0] + 0.587 * pixels[:, 1] + 0.299 * pixels[:, 2]
    saturation = np.max(pixels, axis=1) - np.min(pixels, axis=1)
    mask = (gray < 246) & ((saturation > 7) | (gray < 226))
    if np.count_nonzero(mask) >= 6:
        pixels = pixels[mask]
    return _bgr_to_hex(np.median(pixels, axis=0))


def _is_nearly_background(fill_hex: str, background_hex: str, tolerance: int = 22) -> bool:
    fill = np.asarray(_hex_rgb(fill_hex, (255, 255, 255)), dtype=np.int16)
    bg = np.asarray(_hex_rgb(background_hex, (255, 255, 255)), dtype=np.int16)
    return int(np.max(np.abs(fill - bg))) <= tolerance


def _detect_words_tesseract(image_path: str, cv_img: np.ndarray) -> list[dict[str, Any]]:
    import pytesseract

    executable = _find_tesseract()
    if not executable:
        logger.warning("Tesseract is unavailable; editable text detection cannot run")
        return []
    pytesseract.pytesseract.tesseract_cmd = executable
    pil_img = Image.open(image_path).convert("RGB")
    all_words: list[dict[str, Any]] = []

    # Sparse-text mode is the primary pass.  Automatic page segmentation adds
    # occasional words missed by sparse mode, but its large merged rows are not
    # used directly.
    for psm in (11, 6):
        try:
            data = pytesseract.image_to_data(
                pil_img,
                lang="chi_sim+eng",
                config=f"--psm {psm}",
                output_type=pytesseract.Output.DICT,
            )
        except Exception as exc:
            logger.warning("Tesseract PSM %s failed: %s", psm, exc)
            continue
        for idx, raw_text in enumerate(data["text"]):
            text = (raw_text or "").strip()
            try:
                confidence = float(data["conf"][idx])
            except (TypeError, ValueError):
                confidence = -1
            if not _valid_ocr_token(text, confidence):
                continue
            text = _sanitize_ocr_token(text)
            word = {
                "text": text,
                "confidence": confidence,
                "x": int(data["left"][idx]),
                "y": int(data["top"][idx]),
                "w": int(data["width"][idx]),
                "h": int(data["height"][idx]),
                "psm": psm,
                "block_num": int(data["block_num"][idx]),
                "line_num": int(data["line_num"][idx]),
            }
            if word["w"] <= 0 or word["h"] <= 0:
                continue
            if word["text"].isdigit() and word["confidence"] < 65:
                continue
            if (
                _contains_cjk(word["text"])
                and len(word["text"]) == 1
                and word["w"] / max(word["h"], 1) > 2.2
            ):
                continue
            if (
                re.search(r"[^A-Za-z0-9]", (raw_text or "").strip())
                and re.sub(r"[^A-Za-z0-9]", "", (raw_text or "").strip()).upper()
                in _KNOWN_SHORT_LATIN
                and confidence < 80
            ):
                continue
            word["color"] = _sample_foreground_hex(
                cv_img, word["x"], word["y"], word["w"], word["h"]
            )
            all_words.append(word)

    # Deduplicate the two OCR passes.  Prefer PSM 11 for precise sparse layout,
    # unless the other pass has materially better confidence/text.
    primary = sorted(
        [item for item in all_words if item["psm"] == 11],
        key=lambda item: -item["confidence"],
    )
    secondary = sorted(
        [item for item in all_words if item["psm"] != 11],
        key=lambda item: -item["confidence"],
    )
    kept: list[dict[str, Any]] = []
    for word in primary:
        duplicate = any(
            _iou(word, existing) > 0.42
            and _text_similarity(word["text"], existing["text"]) > 0.45
            for existing in kept
        )
        if not duplicate:
            kept.append(word)
    for word in secondary:
        cx = word["x"] + word["w"] / 2
        cy = word["y"] + word["h"] / 2
        covered = False
        for existing in kept:
            expanded = (
                existing["x"] - existing["h"] * 0.35 <= cx
                <= existing["x"] + existing["w"] + existing["h"] * 0.35
                and existing["y"] - existing["h"] * 0.45 <= cy
                <= existing["y"] + existing["h"] * 1.45
            )
            same_visual = _iou(word, existing) > 0.12
            if expanded or same_visual:
                covered = True
                break
        if not covered:
            kept.append(word)
    return kept


_RAPID_OCR = None


def _get_rapid_ocr():
    """Lazy-load RapidOCR singleton.  ~30s for a full PPT slide vs 11min for PaddleOCR."""
    global _RAPID_OCR
    if _RAPID_OCR is None:
        from rapidocr_onnxruntime import RapidOCR
        _RAPID_OCR = RapidOCR()
        logger.info("RapidOCR initialized")
    return _RAPID_OCR


def _detect_regions_paddle(
    image_path: str,
    cv_img: np.ndarray,
) -> list[dict[str, Any]]:
    """Use RapidOCR as the primary layout reader (22x faster than PaddleOCR)."""
    if os.getenv("PPT_PRECISE_PADDLEOCR", "1").lower() in {"0", "false", "off"}:
        return []
    ih, iw = cv_img.shape[:2]
    try:
        ocr = _get_rapid_ocr()
        # RapidOCR accepts numpy array (BGR or RGB)
        ocr_result, _elapse = ocr(cv_img)
    except Exception as exc:
        logger.warning("RapidOCR unavailable, using Tesseract fallback: %s", exc)
        return []

    regions: list[dict[str, Any]] = []
    sequence = 1
    if ocr_result:
        for box, text, score in ocr_result:
            text = str(text or "").strip()
            confidence = float(score) * 100
            if not text or confidence < 45 or not _CONTENT_RE.search(text):
                continue
            # RapidOCR box: [[x1,y1],[x2,y1],[x2,y2],[x1,y2]]
            if len(box) < 4:
                continue
            xs = [int(p[0]) for p in box]
            ys = [int(p[1]) for p in box]
            x1, y1 = max(0, min(xs)), max(0, min(ys))
            x2, y2 = min(iw, max(xs)), min(ih, max(ys))
            if x2 <= x1 or y2 <= y1:
                continue
            width, height = x2 - x1, y2 - y1
            role = "body"
            if y1 < ih * 0.16 and height > ih * 0.022:
                role = "title"
            elif y1 > ih * 0.92:
                role = "page_number" if len(text) <= 4 else "footer"
            elif text.replace(" ", "").isdigit():
                role = "number"
            elif height > ih * 0.029:
                role = "heading"
            color = _sample_foreground_hex(cv_img, x1, y1, width, height)
            regions.append({
                "id": f"text_{sequence:03d}",
                "text": text,
                "x": x1, "y": y1, "w": width, "h": height,
                "words": [],
                "confidence": round(confidence, 1),
                "median_word_h": float(height),
                "color": color,
                "role": role,
                "align": "left",
                "bold": role in {"title", "heading", "number"},
                "editable": True,
                "source": "rapidocr",
                "image_width": iw, "image_height": ih,
            })
            sequence += 1
    return regions


def _vertical_overlap(a: dict[str, Any], b: dict[str, Any]) -> float:
    overlap = max(0, min(a["y"] + a["h"], b["y"] + b["h"]) - max(a["y"], b["y"]))
    return overlap / max(1, min(a["h"], b["h"]))


def _join_word_text(previous: dict[str, Any] | None, current: dict[str, Any]) -> str:
    if previous is None:
        return current["text"]
    left = previous["text"]
    right = current["text"]
    gap = current["x"] - (previous["x"] + previous["w"])
    char_width = max(
        4.0,
        (previous["w"] / max(len(left), 1) + current["w"] / max(len(right), 1)) / 2,
    )
    if _contains_cjk(left[-1:]) or _contains_cjk(right[:1]):
        separator = ""
    elif left[-1:].isalnum() and right[:1].isalnum() and gap > char_width * 0.35:
        separator = " "
    else:
        separator = ""
    return separator + right


def _make_text_region(
    words: list[dict[str, Any]],
    region_id: str,
    image_width: int,
    image_height: int,
) -> dict[str, Any]:
    words = sorted(words, key=lambda item: item["x"])
    x1 = min(item["x"] for item in words)
    y1 = min(item["y"] for item in words)
    x2 = max(item["x"] + item["w"] for item in words)
    y2 = max(item["y"] + item["h"] for item in words)
    text = ""
    previous = None
    for word in words:
        text += _join_word_text(previous, word)
        previous = word
    heights = [item["h"] for item in words]
    effective_heights = []
    for item in words:
        compact = _normalize_text_for_similarity(item["text"])
        character_count = max(len(compact), 1)
        width_based_height = item["w"] / character_count * 1.28
        effective_heights.append(min(item["h"], max(7.0, width_based_height)))
    confidence = float(np.average(
        [item["confidence"] for item in words],
        weights=[max(1, len(item["text"])) for item in words],
    ))
    role = "body"
    if y1 < image_height * 0.16 and np.median(heights) > image_height * 0.022:
        role = "title"
    elif y1 > image_height * 0.92:
        role = "page_number" if len(text) <= 4 else "footer"
    elif text.replace(" ", "").isdigit():
        role = "number"
    elif np.median(heights) > image_height * 0.027:
        role = "heading"

    colors = [word["color"] for word in words]
    dominant_color = max(set(colors), key=colors.count)
    return {
        "id": region_id,
        "text": text.strip(),
        "x": x1,
        "y": y1,
        "w": x2 - x1,
        "h": y2 - y1,
        "words": words,
        "confidence": round(confidence, 1),
        "median_word_h": float(np.median(effective_heights)),
        "color": dominant_color,
        "role": role,
        "align": "left",
        "bold": role in {"title", "heading", "number"},
        "editable": True,
        "source": "tesseract",
        "image_width": image_width,
    }


def _cluster_words_to_regions(
    words: list[dict[str, Any]],
    image_width: int,
    image_height: int,
) -> list[dict[str, Any]]:
    if not words:
        return []
    words = sorted(words, key=lambda item: (item["y"] + item["h"] / 2, item["x"]))
    rows: list[list[dict[str, Any]]] = []
    for word in words:
        best_row = None
        best_score = -1.0
        for row in rows:
            row_box = {
                "y": min(item["y"] for item in row),
                "h": max(item["y"] + item["h"] for item in row) - min(item["y"] for item in row),
            }
            overlap = _vertical_overlap(word, {"y": row_box["y"], "h": row_box["h"]})
            center_delta = abs(
                (word["y"] + word["h"] / 2)
                - np.median([item["y"] + item["h"] / 2 for item in row])
            )
            tolerance = max(word["h"], np.median([item["h"] for item in row])) * 0.58
            score = overlap if center_delta <= tolerance else -1
            if score > best_score:
                best_score = score
                best_row = row
        if best_row is not None and best_score >= 0.25:
            best_row.append(word)
        else:
            rows.append([word])

    regions: list[dict[str, Any]] = []
    sequence = 1
    for row in rows:
        row.sort(key=lambda item: item["x"])
        median_h = float(np.median([item["h"] for item in row]))
        pieces: list[list[dict[str, Any]]] = [[]]
        for word in row:
            if not pieces[-1]:
                pieces[-1].append(word)
                continue
            previous = pieces[-1][-1]
            gap = word["x"] - (previous["x"] + previous["w"])
            average_char = (
                previous["w"] / max(len(previous["text"]), 1)
                + word["w"] / max(len(word["text"]), 1)
            ) / 2
            # Split independent card/table columns but keep normal sentence gaps.
            split_threshold = max(median_h * 2.5, average_char * 4.2, image_width * 0.025)
            different_block = word.get("block_num") != previous.get("block_num")
            if gap > split_threshold or (different_block and gap > median_h * 0.45):
                pieces.append([word])
            else:
                pieces[-1].append(word)
        for piece in pieces:
            region = _make_text_region(
                piece, f"text_{sequence:03d}", image_width, image_height
            )
            sequence += 1
            normalized = _normalize_text_for_similarity(region["text"])
            if not normalized:
                continue
            # Latin-only low-confidence fragments inside icons are usually noise.
            if (
                not _contains_cjk(region["text"])
                and not any(char.isdigit() for char in region["text"])
                and region["text"].upper() not in _KNOWN_SHORT_LATIN
                and region["confidence"] < 72
            ):
                continue
            regions.append(region)

    # Remove duplicate rows introduced by the secondary OCR pass.
    final: list[dict[str, Any]] = []
    for region in sorted(regions, key=lambda item: (-item["confidence"], item["y"], item["x"])):
        duplicate = False
        for existing in final:
            if _iou(region, existing) > 0.58 and _text_similarity(
                region["text"], existing["text"]
            ) > 0.45:
                duplicate = True
                break
        if not duplicate:
            final.append(region)
    final.sort(key=lambda item: (item["y"], item["x"]))
    for idx, region in enumerate(final, start=1):
        region["id"] = f"text_{idx:03d}"
    return final


def _detect_badge_sequence(
    cv_img: np.ndarray,
    text_regions: list[dict[str, Any]],
) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    """Detect repeated numbered circular badges and infer a visible sequence.

    This is deterministic visual inference, not LLM invention.  It only fires
    for 3-9 similarly sized, approximately evenly spaced circles in the upper
    60% of a slide.
    """
    ih, iw = cv_img.shape[:2]
    gray = cv2.cvtColor(cv_img, cv2.COLOR_BGR2GRAY)
    blurred = cv2.GaussianBlur(gray, (7, 7), 1.4)
    min_radius = max(10, round(min(iw, ih) * 0.018))
    max_radius = max(min_radius + 2, round(min(iw, ih) * 0.052))
    circles = cv2.HoughCircles(
        blurred,
        cv2.HOUGH_GRADIENT,
        dp=1.2,
        minDist=min_radius * 2.2,
        param1=100,
        param2=30,
        minRadius=min_radius,
        maxRadius=max_radius,
    )
    if circles is None:
        return [], []
    candidates = []
    for cx, cy, radius in np.round(circles[0]).astype(int):
        if cy > ih * 0.62:
            continue
        x1, y1 = max(0, cx - radius), max(0, cy - radius)
        x2, y2 = min(iw, cx + radius), min(ih, cy + radius)
        roi = cv_img[y1:y2, x1:x2]
        if roi.size == 0:
            continue
        hsv = cv2.cvtColor(roi, cv2.COLOR_BGR2HSV)
        saturation = float(np.mean(hsv[:, :, 1]))
        if saturation < 35:
            continue
        candidates.append({"cx": cx, "cy": cy, "r": radius})

    groups: list[list[dict[str, Any]]] = []
    for circle in sorted(candidates, key=lambda item: item["cy"]):
        target = None
        for group in groups:
            if (
                abs(circle["cy"] - np.median([item["cy"] for item in group]))
                <= max(circle["r"], np.median([item["r"] for item in group])) * 0.65
                and abs(circle["r"] - np.median([item["r"] for item in group]))
                <= max(4, circle["r"] * 0.32)
            ):
                target = group
                break
        if target is None:
            groups.append([circle])
        else:
            target.append(circle)

    inferred_regions: list[dict[str, Any]] = []
    badge_shapes: list[dict[str, Any]] = []
    for group in groups:
        group = sorted(group, key=lambda item: item["cx"])
        if not 3 <= len(group) <= 9:
            continue
        gaps = np.diff([item["cx"] for item in group])
        if len(gaps) and np.std(gaps) / max(np.mean(gaps), 1) > 0.28:
            continue
        if group[-1]["cx"] - group[0]["cx"] < iw * 0.35:
            continue
        for index, circle in enumerate(group, start=1):
            overlap = any(
                region["x"] <= circle["cx"] <= region["x"] + region["w"]
                and region["y"] <= circle["cy"] <= region["y"] + region["h"]
                and region["text"].strip().isdigit()
                for region in text_regions
            )
            if not overlap:
                inferred_regions.append({
                    "id": f"badge_{index:02d}",
                    "text": str(index),
                    "x": circle["cx"] - circle["r"] * 0.48,
                    "y": circle["cy"] - circle["r"] * 0.66,
                    "w": circle["r"] * 0.96,
                    "h": circle["r"] * 1.32,
                    "words": [],
                    "confidence": 82.0,
                    "median_word_h": circle["r"] * 1.12,
                    "color": "#FFFFFF",
                    "role": "number",
                    "align": "center",
                    "bold": True,
                    "editable": True,
                    "source": "inferred_badge_sequence",
                    "image_width": iw,
                })
        break
    # Gradient/texture badges remain in the raster visual layer.  Replacing
    # them with default PowerPoint circles would reduce visual quality.
    return inferred_regions, []


def _detect_native_shapes(cv_img: np.ndarray) -> list[dict[str, Any]]:
    """Detect conservative, high-confidence editable PowerPoint geometry.

    v3 treats structure as first-class PPT objects instead of background
    chunks: rounded cards, pill boxes, grid dividers and long rules are rebuilt
    natively when their geometry is stable enough.  Icons/logos remain raster
    crops because redrawing them would reduce fidelity.
    """
    ih, iw = cv_img.shape[:2]
    gray = cv2.cvtColor(cv_img, cv2.COLOR_BGR2GRAY)
    edges = cv2.Canny(gray, 55, 150)
    shapes: list[dict[str, Any]] = []
    bg_samples = np.concatenate([
        cv_img[0:8, :, :].reshape(-1, 3),
        cv_img[-8:, :, :].reshape(-1, 3),
        cv_img[:, 0:8, :].reshape(-1, 3),
        cv_img[:, -8:, :].reshape(-1, 3),
    ])
    background_hex = _bgr_to_hex(np.median(bg_samples, axis=0))

    # Closed rounded rectangles / cards / pill bands.
    rect_edges = cv2.dilate(edges, np.ones((3, 3), np.uint8), iterations=1)
    contours, _hierarchy = cv2.findContours(rect_edges, cv2.RETR_LIST, cv2.CHAIN_APPROX_SIMPLE)
    rects: list[dict[str, Any]] = []
    for contour in contours:
        x, y, w, h = cv2.boundingRect(contour)
        if w < max(42, iw * 0.025) or h < max(18, ih * 0.020):
            continue
        if w > iw * 0.94 or h > ih * 0.42:
            continue
        box_area = w * h
        if box_area < iw * ih * 0.0008 or box_area > iw * ih * 0.18:
            continue
        aspect = w / max(h, 1)
        if aspect < 0.55 or aspect > 42:
            continue
        roi_edges = edges[y:y + h, x:x + w]
        if roi_edges.size == 0:
            continue
        band = max(2, min(7, round(min(w, h) * 0.11)))
        border_mask = np.zeros((h, w), dtype=np.uint8)
        border_mask[:band, :] = 1
        border_mask[-band:, :] = 1
        border_mask[:, :band] = 1
        border_mask[:, -band:] = 1
        border_density = float(np.count_nonzero((roi_edges > 0) & (border_mask > 0))) / max(
            1, int(np.count_nonzero(border_mask))
        )
        interior_mask = border_mask == 0
        interior_density = float(np.count_nonzero((roi_edges > 0) & interior_mask)) / max(
            1, int(np.count_nonzero(interior_mask))
        )
        if border_density < 0.012:
            continue
        # Reject text/icon clusters: they tend to have more interior edges than
        # a structural frame.  Soft cards with icons/text inside get a relaxed
        # allowance because OCR text is later masked/overlaid.
        if interior_density > max(0.115, border_density * 3.8) and box_area < iw * ih * 0.035:
            continue
        fill_hex = _sample_fill_hex(cv_img, x, y, w, h)
        stroke_hex = _sample_stroke_hex(cv_img, x, y, w, h)
        # Very small white-on-white contours are usually glyph artifacts.
        if (
            box_area < iw * ih * 0.006
            and _is_nearly_background(fill_hex, background_hex)
            and _is_nearly_background(stroke_hex, background_hex, tolerance=28)
        ):
            continue
        rect_type = "round_rect" if min(w, h) >= 20 else "rect"
        rects.append({
            "type": rect_type,
            "x": int(x),
            "y": int(y),
            "w": int(w),
            "h": int(h),
            "fill": fill_hex,
            "stroke": stroke_hex,
            "width_pt": 0.55 if min(w, h) > 34 else 0.35,
            "radius_ratio": round(min(0.45, max(0.08, min(w, h) / max(w, h) * 0.65)), 3),
            "confidence": 0.82,
        })

    rects.sort(key=lambda item: (item["w"] * item["h"], item["w"]), reverse=True)
    deduped_rects: list[dict[str, Any]] = []
    for rect in rects:
        duplicate = False
        for existing in deduped_rects:
            if _iou(rect, existing) > 0.72:
                duplicate = True
                break
            same_center = (
                abs((rect["x"] + rect["w"] / 2) - (existing["x"] + existing["w"] / 2)) < 8
                and abs((rect["y"] + rect["h"] / 2) - (existing["y"] + existing["h"] / 2)) < 8
            )
            if same_center and abs(rect["w"] - existing["w"]) < 18 and abs(rect["h"] - existing["h"]) < 18:
                duplicate = True
                break
        if not duplicate:
            deduped_rects.append(rect)
        if len(deduped_rects) >= 45:
            break
    shapes.extend(sorted(deduped_rects, key=lambda item: (item["y"], item["x"])))

    lines = cv2.HoughLinesP(
        edges,
        1,
        np.pi / 180,
        threshold=max(28, round(iw * 0.022)),
        minLineLength=max(36, round(iw * 0.035)),
        maxLineGap=max(6, round(iw * 0.006)),
    )
    line_shapes: list[dict[str, Any]] = []
    if lines is not None:
        for raw in lines[:, 0]:
            x1, y1, x2, y2 = map(int, raw)
            if x2 < x1 or y2 < y1:
                x1, y1, x2, y2 = x2, y2, x1, y1
            dx, dy = x2 - x1, y2 - y1
            length = math.hypot(dx, dy)
            angle = abs(math.degrees(math.atan2(dy, dx)))
            if not (angle <= 2.5 or angle >= 177.5 or 87.5 <= angle <= 92.5):
                continue
            horizontal = angle <= 2.5 or angle >= 177.5
            minimum_length = max(48, iw * 0.055) if horizontal else max(38, ih * 0.055)
            if length < minimum_length:
                continue
            mx, my = (x1 + x2) // 2, (y1 + y2) // 2
            sample = cv_img[max(0, my - 1):my + 2, max(0, mx - 1):mx + 2]
            bgr = np.median(sample.reshape(-1, 3), axis=0).astype(int)
            luminance = 0.114 * bgr[0] + 0.587 * bgr[1] + 0.299 * bgr[2]
            saturation = int(max(bgr) - min(bgr))
            if luminance < 105 and saturation < 30:
                continue
            # Rectangle borders are already represented by native rectangles.
            on_rect_border = False
            for rect in deduped_rects:
                near_horizontal_border = horizontal and (
                    abs(y1 - rect["y"]) < 5 or abs(y1 - (rect["y"] + rect["h"])) < 5
                ) and min(x2, rect["x"] + rect["w"]) - max(x1, rect["x"]) > length * 0.55
                near_vertical_border = (not horizontal) and (
                    abs(x1 - rect["x"]) < 5 or abs(x1 - (rect["x"] + rect["w"])) < 5
                ) and min(y2, rect["y"] + rect["h"]) - max(y1, rect["y"]) > length * 0.55
                if near_horizontal_border or near_vertical_border:
                    on_rect_border = True
                    break
            if on_rect_border:
                continue
            line_shapes.append({
                "type": "line",
                "x1": x1, "y1": y1, "x2": x2, "y2": y2,
                "stroke": f"#{bgr[2]:02X}{bgr[1]:02X}{bgr[0]:02X}",
                "width_pt": 0.35 if luminance > 135 else 0.55,
                "confidence": 0.78,
            })

    # Deduplicate collinear Hough lines.
    deduped: list[dict[str, Any]] = []
    for shape in sorted(
        line_shapes,
        key=lambda item: -math.hypot(item["x2"] - item["x1"], item["y2"] - item["y1"]),
    ):
        duplicate = False
        for existing in deduped:
            horizontal = abs(shape["y2"] - shape["y1"]) < abs(shape["x2"] - shape["x1"])
            if horizontal:
                if abs(shape["y1"] - existing["y1"]) < 4 and min(
                    shape["x2"], existing["x2"]
                ) - max(shape["x1"], existing["x1"]) > 20:
                    duplicate = True
            else:
                if abs(shape["x1"] - existing["x1"]) < 4 and min(
                    shape["y2"], existing["y2"]
                ) - max(shape["y1"], existing["y1"]) > 20:
                    duplicate = True
            if duplicate:
                break
        if not duplicate:
            deduped.append(shape)
    shapes.extend(deduped[:80])
    return shapes[:120]


def _region_local_complexity(cv_img: np.ndarray, region: dict[str, Any]) -> float:
    ih, iw = cv_img.shape[:2]
    x1 = max(0, int(region["x"]))
    y1 = max(0, int(region["y"]))
    x2 = min(iw, int(region["x"] + region["w"]))
    y2 = min(ih, int(region["y"] + region["h"]))
    roi = cv_img[y1:y2, x1:x2]
    if roi.size == 0:
        return 0.0
    gray = cv2.cvtColor(roi, cv2.COLOR_BGR2GRAY)
    edge_density = float(np.mean(cv2.Canny(gray, 60, 150) > 0))
    return min(1.0, edge_density * 5 + float(np.std(gray)) / 255.0)


def analyze_layers(image_path: str) -> dict[str, Any]:
    """Measure the slide and separate editable information from visual assets."""
    cv_img = _read_cv_image(image_path)
    ih, iw = cv_img.shape[:2]
    paddle_regions = _detect_regions_paddle(image_path, cv_img)
    words: list[dict[str, Any]] = []
    if paddle_regions:
        # Tesseract is still useful as a lightweight word/color segmenter.
        # Paddle remains the source of truth for text and geometry.
        words = _detect_words_tesseract(image_path, cv_img)
        for region in paddle_regions:
            matched = []
            for word in words:
                if word["confidence"] < 70:
                    continue
                center_x = word["x"] + word["w"] / 2
                center_y = word["y"] + word["h"] / 2
                if (
                    region["x"] - region["h"] * 0.35
                    <= center_x
                    <= region["x"] + region["w"] + region["h"] * 0.35
                    and region["y"] - region["h"] * 0.45
                    <= center_y
                    <= region["y"] + region["h"] * 1.35
                ):
                    matched.append(word)
            region["words"] = sorted(matched, key=lambda item: item["x"])
        text_regions = paddle_regions
    else:
        words = _detect_words_tesseract(image_path, cv_img)
        text_regions = _cluster_words_to_regions(words, iw, ih)
    inferred_badges, badge_shapes = _detect_badge_sequence(cv_img, text_regions)
    text_regions.extend(inferred_badges)
    text_regions.sort(key=lambda item: (item["y"], item["x"]))

    raster_text: list[dict[str, Any]] = []
    editable_text: list[dict[str, Any]] = []
    for region in text_regions:
        complexity = _region_local_complexity(cv_img, region)
        region["complexity"] = round(complexity, 3)
        # Tiny, uncertain labels embedded inside complex graphics are safer as
        # raster.  Titles/body/KPIs/page numbers are never routed here.
        raster_only = (
            region["median_word_h"] < max(8, ih * 0.009)
            and region["confidence"] < 48
            and complexity > 0.52
            and region["role"] not in {"title", "heading", "number", "page_number"}
        )
        if raster_only:
            region["editable"] = False
            raster_text.append(region)
        else:
            editable_text.append(region)

    bg_samples = np.concatenate([
        cv_img[0:8, :, :].reshape(-1, 3),
        cv_img[-8:, :, :].reshape(-1, 3),
        cv_img[:, 0:8, :].reshape(-1, 3),
        cv_img[:, -8:, :].reshape(-1, 3),
    ])
    bg_bgr = np.median(bg_samples, axis=0).astype(int)
    background = f"#{bg_bgr[2]:02X}{bg_bgr[1]:02X}{bg_bgr[0]:02X}"
    return {
        "image_width": iw,
        "image_height": ih,
        "background": background,
        "editable_text": editable_text,
        "raster_text": raster_text,
        "native_shapes": badge_shapes + _detect_native_shapes(cv_img),
        "ocr_word_count": len(words) if words else len(paddle_regions),
        "ocr_engine": "rapidocr+tesseract_color" if paddle_regions else "tesseract",
    }


async def _refine_text_with_deepseek(
    regions: list[dict[str, Any]],
    image_width: int,
    image_height: int,
) -> tuple[list[dict[str, Any]], dict[str, Any]]:
    """Let the existing DeepSeek backend correct OCR metadata.

    No image or alternate model is used.  Validation rejects corrections that
    depart too far from OCR evidence.
    """
    if os.getenv("PPT_PRECISE_LLM_REFINEMENT", "1").lower() in {"0", "false", "off"}:
        return regions, {"used": False, "reason": "disabled"}
    try:
        from app.config import settings
        from app.services.llm_service import llm_service
        if not settings.anthropic_api_key:
            return regions, {"used": False, "reason": "missing_api_key"}
    except Exception as exc:
        return regions, {"used": False, "reason": f"service_unavailable:{exc}"}

    payload = []
    for region in regions:
        payload.append({
            "id": region["id"],
            "ocr_text": region["text"],
            "bbox": [round(region["x"]), round(region["y"]), round(region["w"]), round(region["h"])],
            "confidence": region["confidence"],
            "role_hint": region["role"],
            "source": region["source"],
            "word_evidence": [
                {
                    "text": word["text"],
                    "confidence": round(word["confidence"]),
                    "x": word["x"],
                    "color": word["color"],
                }
                for word in region.get("words", [])
            ],
        })

    system_prompt = """你是PPT视觉稿OCR校对器。你看不到图片，只能使用OCR文字、坐标、置信度、
颜色和重复版式证据。禁止臆造图片中没有的业务内容。

输出纯JSON：
{"items":[{"id":"text_001","keep":true,"text":"校正文字","role":"title|heading|body|label|number|footer|page_number","align":"left|center|right","bold":true}]}

规则：
1. 修正常见中文OCR错字、断字、乱码和中英文空格。
2. 同一行的碎片已由程序合并，不要合并不同id，也不要改变id。
3. 图标误识别出的无意义拉丁字母/符号可设置keep=false。
4. 高置信度中文、数字、专有缩写必须保留；不能确定时保持原文。
5. inferred_badge_sequence是程序根据重复圆形卡点推断出的可见序号，可以保留。
6. 不要输出解释或Markdown。"""
    user_prompt = (
        f"画布尺寸：{image_width}x{image_height}\n"
        "请校对以下OCR区域：\n"
        + json.dumps(payload, ensure_ascii=False)
    )
    system_prompt = (
        "You are an OCR proofreading assistant for PPT visual mockups. "
        "You cannot see the image. Use only OCR text, bounding boxes, confidence, "
        "sampled colors, and repeated layout evidence. Do not invent business content. "
        "Return JSON only in this exact schema: "
        "{\"items\":[{\"id\":\"text_001\",\"keep\":true,\"text\":\"corrected text\","
        "\"role\":\"title|heading|body|label|number|footer|page_number\","
        "\"align\":\"left|center|right\",\"bold\":true}]}. "
        "Rules: fix obvious Chinese OCR errors, broken words, mojibake, and unnatural "
        "spacing between Chinese, English, and numbers; never merge different ids; "
        "set keep=false only for meaningless icon/decorative fragments; preserve "
        "high-confidence Chinese, numbers, and acronyms; if unsure, keep the original; "
        "output no Markdown and no explanation."
    )
    user_prompt = (
        f"Canvas size: {image_width}x{image_height}\n"
        "Please proofread these OCR regions:\n"
        + json.dumps(payload, ensure_ascii=False)
    )
    try:
        response = await llm_service.chat(
            system_prompt=system_prompt,
            messages=[{"role": "user", "content": user_prompt}],
            max_tokens=8192,
            temperature=0.1,
            timeout=120,
            thinking={"type": "disabled"},
        )
        raw = ""
        for block in getattr(response, "content", []) or []:
            raw += getattr(block, "text", "")
        match = re.search(r"\{.*\}", raw, re.DOTALL)
        data = json.loads(match.group(0) if match else raw)
        updates = {
            item["id"]: item
            for item in data.get("items", [])
            if isinstance(item, dict) and item.get("id")
        }
    except Exception as exc:
        logger.warning("DeepSeek OCR refinement failed: %s", exc)
        return regions, {"used": False, "reason": f"request_failed:{exc.__class__.__name__}"}

    refined: list[dict[str, Any]] = []
    accepted = 0
    rejected = 0
    for region in regions:
        update = updates.get(region["id"])
        if not update:
            refined.append(region)
            continue
        if update.get("keep") is False:
            if (
                region["confidence"] < 68
                and not _contains_cjk(region["text"])
                and region["source"] != "inferred_badge_sequence"
            ):
                accepted += 1
                continue
            rejected += 1
            refined.append(region)
            continue
        candidate = str(update.get("text", "")).strip()
        similarity = _text_similarity(region["text"], candidate)
        can_accept = (
            candidate
            and (
                similarity >= 0.52
                or region["confidence"] < 52
                or region["source"] == "inferred_badge_sequence"
            )
        )
        if can_accept:
            region = dict(region)
            region["text"] = candidate
            if region["confidence"] < 75 and update.get("role") in {
                "title", "heading", "body", "label", "number", "footer", "page_number"
            }:
                region["role"] = update["role"]
            if update.get("align") in {"left", "center", "right"}:
                region["align"] = update["align"]
            if isinstance(update.get("bold"), bool):
                if not region["bold"] or update["bold"]:
                    region["bold"] = update["bold"]
            region["llm_refined"] = candidate != region.get("ocr_text", region["text"])
            accepted += 1
        else:
            rejected += 1
        refined.append(region)
    return refined, {"used": True, "accepted": accepted, "rejected": rejected}


def _build_text_mask(
    cv_img: np.ndarray,
    editable_regions: list[dict[str, Any]],
) -> np.ndarray:
    ih, iw = cv_img.shape[:2]
    mask = np.zeros((ih, iw), dtype=np.uint8)
    for region in editable_regions:
        words = region.get("words", [])
        boxes = words if words else [region]
        for item in boxes:
            padding_x = max(1, round(item["h"] * 0.08))
            padding_y = max(1, round(item["h"] * 0.09))
            x1 = max(0, round(item["x"] - padding_x))
            y1 = max(0, round(item["y"] - padding_y))
            x2 = min(iw, round(item["x"] + item["w"] + padding_x))
            y2 = min(ih, round(item["y"] + item["h"] + padding_y))
            mask[y1:y2, x1:x2] = 255
    kernel = np.ones((3, 3), np.uint8)
    return cv2.dilate(mask, kernel, iterations=1)


def _make_clean_visual_base(
    cv_img: np.ndarray,
    editable_regions: list[dict[str, Any]],
) -> tuple[np.ndarray, np.ndarray]:
    mask = _build_text_mask(cv_img, editable_regions)
    if not np.any(mask):
        return cv_img.copy(), mask
    cleaned = cv_img.copy()
    ih, iw = cv_img.shape[:2]
    inpaint_mask = np.zeros((ih, iw), dtype=np.uint8)
    for region in editable_regions:
        padding_x = max(2, round(region["median_word_h"] * 0.16))
        padding_y = max(2, round(region["median_word_h"] * 0.18))
        x1 = max(0, round(region["x"] - padding_x))
        y1 = max(0, round(region["y"] - padding_y))
        x2 = min(iw, round(region["x"] + region["w"] + padding_x))
        y2 = min(ih, round(region["y"] + region["h"] + padding_y))
        if x2 <= x1 or y2 <= y1:
            continue
        if (
            region.get("source") == "inferred_badge_sequence"
            or (
                region.get("role") == "number"
                and len(str(region.get("text", ""))) <= 2
                and region["y"] < ih * 0.55
            )
        ):
            # Keep the original textured badge and place an editable number
            # exactly over it.  Inpainting small white digits creates a more
            # visible scar than the coincident editable overlay.
            continue
        background = _sample_background_bgr(
            cv_img, x1, y1, x2 - x1, y2 - y1
        ).astype(np.uint8)
        border_roi = cv_img[
            max(0, y1 - 4):min(ih, y2 + 4),
            max(0, x1 - 4):min(iw, x2 + 4),
        ]
        variation = float(np.mean(np.std(border_roi.reshape(-1, 3), axis=0)))
        if variation < 45:
            cleaned[y1:y2, x1:x2] = background
        else:
            inpaint_mask[y1:y2, x1:x2] = 255
    if np.any(inpaint_mask):
        cleaned = cv2.inpaint(
            cleaned, inpaint_mask, inpaintRadius=5, flags=cv2.INPAINT_TELEA
        )
    return cleaned, mask


def _extract_visual_object_assets(
    cleaned: np.ndarray,
    text_mask: np.ndarray,
    work_dir: Path,
    native_shapes: list[dict[str, Any]] | None = None,
    max_assets: int | None = None,
) -> tuple[list[dict[str, Any]], np.ndarray]:
    """Extract non-text visual strokes as independent movable image objects.

    This is intentionally conservative.  It does not try to redraw official
    icons, logos, UI screenshots or complex infographics as fake PPT vectors.
    Instead it cuts their source pixels into transparent PNG objects so they
    remain visually faithful while still being selectable/movable in PPT.
    """
    ih, iw = cleaned.shape[:2]
    max_assets = max_assets or int(os.getenv("PPT_PRECISE_MAX_OBJECT_ASSETS", "72"))
    object_dir = work_dir / "visual_objects"
    object_dir.mkdir(parents=True, exist_ok=True)

    gray = cv2.cvtColor(cleaned, cv2.COLOR_BGR2GRAY)
    hsv = cv2.cvtColor(cleaned, cv2.COLOR_BGR2HSV)
    saturation = hsv[:, :, 1]
    value = hsv[:, :, 2]

    # Corporate slides often use teal strokes/icons plus black/gray outline
    # icons.  Combine color and edge evidence, then remove OCR text regions.
    dark_or_colored = ((gray < 170) | ((saturation > 28) & (value < 248))).astype(np.uint8) * 255
    edges = cv2.Canny(gray, 55, 150)
    candidate = cv2.bitwise_or(dark_or_colored, cv2.dilate(edges, np.ones((3, 3), np.uint8)))
    candidate[text_mask > 0] = 0
    if native_shapes:
        native_mask = np.zeros_like(candidate)
        for shape in native_shapes:
            try:
                shape_type = shape.get("type")
                if shape_type == "line":
                    cv2.line(
                        native_mask,
                        (int(shape["x1"]), int(shape["y1"])),
                        (int(shape["x2"]), int(shape["y2"])),
                        255,
                        max(5, round(min(iw, ih) * 0.006)),
                    )
                elif shape_type in {"rect", "round_rect", "ellipse"}:
                    x = int(shape["x"])
                    y = int(shape["y"])
                    w = int(shape["w"])
                    h = int(shape["h"])
                    thickness = max(5, min(12, round(min(w, h) * 0.18)))
                    # Mask only the border.  The interior may contain icons
                    # that should remain individually selectable.
                    cv2.rectangle(native_mask, (x, y), (x + w, y + h), 255, thickness)
            except Exception:
                continue
        native_mask = cv2.dilate(native_mask, np.ones((3, 3), np.uint8), iterations=1)
        candidate[native_mask > 0] = 0

    # Remove slide margins and tiny salt noise.
    margin = max(2, round(min(iw, ih) * 0.004))
    candidate[:margin, :] = 0
    candidate[-margin:, :] = 0
    candidate[:, :margin] = 0
    candidate[:, -margin:] = 0
    candidate = cv2.morphologyEx(candidate, cv2.MORPH_OPEN, np.ones((2, 2), np.uint8))
    # A smaller close kernel prevents neighboring icons, arrows and card
    # borders from being fused into one uneditable blob.
    candidate = cv2.morphologyEx(candidate, cv2.MORPH_CLOSE, np.ones((3, 3), np.uint8))

    num_labels, labels, stats, _centroids = cv2.connectedComponentsWithStats(candidate, 8)
    min_area = max(28, int(iw * ih * 0.000025))
    max_box_area = iw * ih * 0.60  # was 0.07 - allow large visual objects
    components: list[dict[str, Any]] = []
    for label in range(1, num_labels):
        x, y, w, h, area = [int(v) for v in stats[label]]
        if area < min_area or w < 6 or h < 6:
            continue
        box_area = w * h
        if box_area > max_box_area:
            continue
        aspect = w / max(h, 1)
        # Long separators/grid lines are already preserved by the tile layer
        # and optionally by native line detection.
        if aspect > 18 or aspect < 1 / 18:
            continue
        if w > iw * 0.34 and h > ih * 0.22:
            continue
        component_crop = (labels[y:y + h, x:x + w] == label)
        text_overlap = np.count_nonzero(component_crop & (text_mask[y:y + h, x:x + w] > 0))
        if text_overlap / max(1, int(component_crop.sum())) > 0.08:
            continue
        stroke_density = area / max(1, box_area)
        if stroke_density < 0.004:
            continue
        components.append({
            "x": x,
            "y": y,
            "w": w,
            "h": h,
            "area": area,
            "box_area": box_area,
            "label": label,
        })

    # Prefer meaningful medium/large objects; dedupe near-identical boxes.
    components.sort(key=lambda item: (item["box_area"], item["area"]), reverse=True)
    kept: list[dict[str, Any]] = []
    for component in components:
        box = {
            "x": component["x"],
            "y": component["y"],
            "w": component["w"],
            "h": component["h"],
        }
        duplicate = False
        for existing in kept:
            existing_box = {
                "x": existing["x"],
                "y": existing["y"],
                "w": existing["w"],
                "h": existing["h"],
            }
            if _iou(box, existing_box) > 0.78:
                duplicate = True
                break
        if not duplicate:
            kept.append(component)
        if len(kept) >= max_assets:
            break

    assets: list[dict[str, Any]] = []
    object_mask = np.zeros((ih, iw), dtype=np.uint8)
    for index, component in enumerate(sorted(kept, key=lambda item: (item["y"], item["x"])), start=1):
        pad = max(2, round(min(component["w"], component["h"]) * 0.08))
        x1 = max(0, component["x"] - pad)
        y1 = max(0, component["y"] - pad)
        x2 = min(iw, component["x"] + component["w"] + pad)
        y2 = min(ih, component["y"] + component["h"] + pad)
        if x2 <= x1 or y2 <= y1:
            continue

        crop_bgr = cleaned[y1:y2, x1:x2]
        component_alpha = (labels[y1:y2, x1:x2] == component["label"]).astype(np.uint8) * 255
        component_alpha = cv2.dilate(component_alpha, np.ones((3, 3), np.uint8), iterations=1)
        component_alpha = cv2.GaussianBlur(component_alpha, (3, 3), 0)
        if np.count_nonzero(component_alpha > 12) < min_area:
            continue

        bgra = cv2.cvtColor(crop_bgr, cv2.COLOR_BGR2BGRA)
        bgra[:, :, 3] = component_alpha
        filename = f"object_{index:02d}.png"
        path = object_dir / filename
        cv2.imencode(".png", bgra)[1].tofile(str(path))
        object_mask[y1:y2, x1:x2] = cv2.max(object_mask[y1:y2, x1:x2], component_alpha)
        assets.append({
            "path": str(path),
            "x": x1,
            "y": y1,
            "w": x2 - x1,
            "h": y2 - y1,
            "kind": "object",
            "reason": "transparent crop of non-text visual object/icon/logo",
        })

    if assets:
        object_mask = cv2.dilate(object_mask, np.ones((3, 3), np.uint8), iterations=1)
    return assets, object_mask


def _remove_object_assets_from_base(cleaned: np.ndarray, object_mask: np.ndarray) -> np.ndarray:
    if object_mask is None or not np.any(object_mask):
        return cleaned
    # Inpaint only the extracted strokes/icons so the top-layer object crops
    # become the visible source of truth at their original position.
    return cv2.inpaint(cleaned, object_mask, inpaintRadius=3, flags=cv2.INPAINT_TELEA)


def _tile_visual_base(
    cleaned: np.ndarray,
    work_dir: Path,
    columns: int | None = None,
    rows: int | None = None,
) -> list[dict[str, Any]]:
    """Optionally save local background crops; default is native PPT background.

    The user-facing editable mode should not leave a tiled bitmap substrate
    behind.  For simple business slides we instead rely on the slide background
    plus native shapes and independently movable icon crops.  A tile fallback is
    retained only for deliberately complex/photo backgrounds via env opt-in.
    """
    ih, iw = cleaned.shape[:2]
    mode = os.getenv("PPT_PRECISE_BACKGROUND_MODE", "native").strip().lower()
    if mode not in {"tile", "tiles", "tiled"}:
        return []
    if columns is None:
        columns = max(2, int(os.getenv("PPT_PRECISE_TILE_COLUMNS", "6")))
    if rows is None:
        rows = max(2, int(os.getenv("PPT_PRECISE_TILE_ROWS", "4")))
    tile_dir = work_dir / "visual_tiles"
    tile_dir.mkdir(parents=True, exist_ok=True)

    assets: list[dict[str, Any]] = []
    for row in range(rows):
        y1 = round(ih * row / rows)
        y2 = round(ih * (row + 1) / rows)
        for column in range(columns):
            x1 = round(iw * column / columns)
            x2 = round(iw * (column + 1) / columns)
            crop = cleaned[y1:y2, x1:x2]
            if crop.size == 0:
                continue
            filename = f"tile_r{row + 1}_c{column + 1}.png"
            path = tile_dir / filename
            cv2.imencode(".png", crop)[1].tofile(str(path))
            assets.append({
                "path": str(path),
                "x": x1,
                "y": y1,
                "w": x2 - x1,
                "h": y2 - y1,
                "kind": "tile",
                "reason": "text/object-cleaned local background crop",
            })
    return assets


def _hex_rgb(value: str, default: tuple[int, int, int] = (0, 0, 0)) -> tuple[int, int, int]:
    value = (value or "").strip().lstrip("#")
    if len(value) != 6:
        return default
    try:
        return int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16)
    except ValueError:
        return default


def _set_slide_background(slide, color: str) -> None:
    from lxml import etree
    from pptx.oxml.ns import qn

    r, g, b = _hex_rgb(color, (255, 255, 255))
    bg = slide.background._element
    bg_pr = bg.find(qn("p:bgPr"))
    if bg_pr is None:
        bg_pr = etree.SubElement(bg, qn("p:bgPr"))
    for child in list(bg_pr):
        bg_pr.remove(child)
    solid = etree.SubElement(bg_pr, qn("a:solidFill"))
    rgb = etree.SubElement(solid, qn("a:srgbClr"))
    rgb.set("val", f"{r:02X}{g:02X}{b:02X}")
    etree.SubElement(bg_pr, qn("a:effectLst"))


def _set_east_asian_font(run, font_name: str) -> None:
    from lxml import etree
    from pptx.oxml.ns import qn

    run.font.name = font_name
    r_pr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        existing = r_pr.find(qn(tag))
        if existing is None:
            existing = etree.SubElement(r_pr, qn(tag))
        existing.set("typeface", font_name)


def _paragraph_alignment(value: str):
    from pptx.enum.text import PP_ALIGN
    return {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get(value, PP_ALIGN.LEFT)


def _infer_alignment(region: dict[str, Any], image_width: int) -> str:
    if region.get("align") in {"center", "right"}:
        return region["align"]
    if region["role"] == "number" and len(region["text"]) <= 4:
        return "center"
    # Tight OCR boxes are safest as left aligned.  Center only short labels
    # located close to one of six/eight common grid centers.
    if len(region["text"]) <= 12 and region["w"] < image_width * 0.14:
        center = region["x"] + region["w"] / 2
        for columns in (6, 5, 4, 3):
            grid_centers = [(i + 0.5) * image_width / columns for i in range(columns)]
            if min(abs(center - item) for item in grid_centers) < image_width * 0.018:
                return "center"
    return "left"


def _hex_from_bgr(bgr: np.ndarray | list[int] | tuple[int, int, int]) -> str:
    b, g, r = [int(max(0, min(255, round(float(value))))) for value in bgr]
    return f"#{r:02X}{g:02X}{b:02X}"


def _rgb_distance(a: str, b: str) -> float:
    ar, ag, ab = _hex_rgb(a, (0, 0, 0))
    br, bg, bb = _hex_rgb(b, (0, 0, 0))
    return math.sqrt((ar - br) ** 2 + (ag - bg) ** 2 + (ab - bb) ** 2)


def _is_latin_dominant(text: str) -> bool:
    cjk = len(_CJK_RE.findall(text))
    latin_digits = len(re.findall(r"[A-Za-z0-9]", text))
    return latin_digits > cjk * 1.2


def _font_name_for_region(region: dict[str, Any]) -> str:
    role = region.get("role", "body")
    text = str(region.get("text", ""))
    if role in {"number", "page_number"} and not _contains_cjk(text):
        return DEFAULT_NUMBER_FONT
    if _is_latin_dominant(text) and not _contains_cjk(text):
        return DEFAULT_LATIN_FONT
    return DEFAULT_CJK_FONT


def _font_multiplier_for_region(region: dict[str, Any]) -> float:
    role = region.get("role", "body")
    text = str(region.get("text", ""))
    if role == "title":
        return 0.99
    if role == "heading":
        return 0.98
    if role == "number":
        return 1.03
    if role == "page_number":
        return 0.94
    if len(text) > 36:
        return 0.94
    return 0.97 if role == "body" else 1.0


def _effective_bold(region: dict[str, Any]) -> bool:
    if bool(region.get("bold")):
        return True
    role = region.get("role", "body")
    if role in {"title", "heading", "number", "page_number"}:
        return True
    image_height = max(1, int(region.get("image_height") or 900))
    if region.get("median_word_h", 0) >= image_height * 0.022 and len(str(region.get("text", ""))) <= 18:
        return True
    # Teal labels in corporate slides are usually semibold; PowerPoint exposes
    # only a bold switch through python-pptx, so use it conservatively.
    color = str(region.get("color", "#111111"))
    r, g, b = _hex_rgb(color, (17, 17, 17))
    if g > r * 1.15 and g >= b * 0.85 and len(str(region.get("text", ""))) <= 24:
        return True
    return False


def _visual_color_runs(
    cv_img: np.ndarray,
    region: dict[str, Any],
    text: str,
) -> list[tuple[str, str]] | None:
    """Approximate per-character colors from the source pixels.

    OCR engines often return a full title as one text region.  This helper
    recovers mixed black/teal emphasis such as "6 key checkpoints" by sampling
    foreground pixels by character column.
    """
    text = text or ""
    if len(text) < 2 or len(text) > 120:
        return None
    ih, iw = cv_img.shape[:2]
    pad = max(2, round(float(region.get("median_word_h", region.get("h", 16))) * 0.12))
    x1 = max(0, round(region["x"] - pad))
    y1 = max(0, round(region["y"] - pad))
    x2 = min(iw, round(region["x"] + region["w"] + pad))
    y2 = min(ih, round(region["y"] + region["h"] + pad))
    if x2 <= x1 or y2 <= y1:
        return None
    crop = cv_img[y1:y2, x1:x2]
    if crop.size == 0:
        return None

    bg = _sample_background_bgr(cv_img, x1, y1, x2 - x1, y2 - y1)
    distance = np.linalg.norm(crop.astype(np.float32) - bg.reshape(1, 1, 3), axis=2)
    gray = cv2.cvtColor(crop, cv2.COLOR_BGR2GRAY)
    threshold = max(22.0, float(np.percentile(distance, 70)))
    mask = (distance >= threshold) & (gray < 248)
    if int(mask.sum()) < max(10, len(text) * 2):
        return None

    default_color = region.get("color") or _sample_foreground_hex(
        cv_img, int(region["x"]), int(region["y"]), int(region["w"]), int(region["h"])
    )
    width = crop.shape[1]
    colors: list[str] = []
    for index, _character in enumerate(text):
        start = int(round(width * index / len(text)))
        end = int(round(width * (index + 1) / len(text)))
        end = max(end, start + 1)
        band_mask = mask[:, start:end]
        if int(band_mask.sum()) < 3:
            colors.append(default_color)
            continue
        pixels = crop[:, start:end][band_mask]
        color = _hex_from_bgr(np.median(pixels, axis=0))
        if _rgb_distance(color, default_color) < 48:
            color = default_color
        colors.append(color)

    # Smooth isolated one-character color glitches caused by anti-aliasing.
    for index in range(1, len(colors) - 1):
        if colors[index - 1] == colors[index + 1] and _rgb_distance(colors[index], colors[index - 1]) < 80:
            colors[index] = colors[index - 1]

    runs: list[tuple[str, str]] = []
    for character, color in zip(text, colors):
        if runs and _rgb_distance(runs[-1][1], color) < 38:
            runs[-1] = (runs[-1][0] + character, runs[-1][1])
        else:
            runs.append((character, color))
    distinct = {color for _, color in runs}
    if len(distinct) <= 1:
        return None
    if len(runs) > 10:
        return None
    return runs


def _augment_text_styles(cv_img: np.ndarray, regions: list[dict[str, Any]]) -> None:
    for region in regions:
        text = str(region.get("text", "")).strip()
        if not text:
            continue
        region.setdefault("image_width", cv_img.shape[1])
        region.setdefault("image_height", cv_img.shape[0])
        region["font_name"] = _font_name_for_region(region)
        region["bold"] = _effective_bold(region)
        visual_runs = _visual_color_runs(cv_img, region, text)
        if visual_runs:
            region["visual_color_runs"] = [
                {"text": run_text, "color": run_color}
                for run_text, run_color in visual_runs
            ]
            # Use the longest run as the default fallback color.
            longest = max(visual_runs, key=lambda item: len(item[0]))
            region["color"] = longest[1]


def _colored_text_runs(
    region: dict[str, Any],
    target_text: str,
) -> list[tuple[str, str]] | None:
    visual_runs = region.get("visual_color_runs") or []
    if visual_runs:
        joined = "".join(str(item.get("text", "")) for item in visual_runs)
        if joined == target_text:
            return [
                (str(item.get("text", "")), str(item.get("color", region.get("color", "#111111"))))
                for item in visual_runs
                if item.get("text")
            ]
    words = region.get("words", [])
    if not words:
        return None
    source_text = ""
    source_colors: list[str] = []
    previous = None
    for word in words:
        segment = _join_word_text(previous, word)
        source_text += segment
        source_colors.extend([word["color"]] * len(segment))
        previous = word
    if not source_text or difflib.SequenceMatcher(None, source_text, target_text).ratio() < 0.62:
        return None

    default_color = region.get("color", "#111111")
    target_colors = [default_color] * len(target_text)
    matcher = difflib.SequenceMatcher(None, source_text, target_text)
    for tag, i1, i2, j1, j2 in matcher.get_opcodes():
        if tag != "equal":
            continue
        for offset in range(min(i2 - i1, j2 - j1)):
            if i1 + offset < len(source_colors):
                target_colors[j1 + offset] = source_colors[i1 + offset]

    runs: list[tuple[str, str]] = []
    for character, color in zip(target_text, target_colors):
        if runs and runs[-1][1] == color:
            runs[-1] = (runs[-1][0] + character, color)
        else:
            runs.append((character, color))
    return runs


def _build_pptx(
    output_path: str,
    mapper: CoordinateMapper,
    background: str,
    assets: list[dict[str, Any]],
    native_shapes: list[dict[str, Any]],
    text_regions: list[dict[str, Any]],
    font_scale: float,
) -> dict[str, int]:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR
    from pptx.util import Inches, Pt

    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_W_IN)
    presentation.slide_height = Inches(SLIDE_H_IN)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _set_slide_background(slide, background)

    counts = {"visual_assets": 0, "object_assets": 0, "editable_shapes": 0, "editable_text": 0}

    # Optional visual substrate.  In v3 this is normally empty; tile fallback is
    # opt-in for genuinely complex backgrounds.
    tile_assets = [asset for asset in assets if asset.get("kind") != "object"]
    object_assets = [asset for asset in assets if asset.get("kind") == "object"]
    for index, asset in enumerate(tile_assets, start=1):
        x, y, w, h = mapper.rect(asset["x"], asset["y"], asset["w"], asset["h"])
        picture = slide.shapes.add_picture(
            asset["path"], Inches(x), Inches(y), Inches(w), Inches(h)
        )
        picture.name = f"VisualTile_{index:02d}"
        counts["visual_assets"] += 1

    # High-confidence editable geometry: cards, pills, row/column separators
    # and rules sit below icons and editable text.
    for index, shape_data in enumerate(native_shapes, start=1):
        try:
            shape_type = shape_data["type"]
            stroke_rgb = RGBColor(*_hex_rgb(shape_data.get("stroke", "#A8BFC2"), (168, 191, 194)))
            if shape_type == "line":
                x1, y1 = mapper.x(shape_data["x1"]), mapper.y(shape_data["y1"])
                x2, y2 = mapper.x(shape_data["x2"]), mapper.y(shape_data["y2"])
                shape = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    Inches(x1), Inches(y1), Inches(x2), Inches(y2),
                )
                shape.line.color.rgb = stroke_rgb
                shape.line.width = Pt(shape_data.get("width_pt", 0.8))
            elif shape_type == "ellipse":
                x, y, w, h = mapper.rect(
                    shape_data["x"], shape_data["y"], shape_data["w"], shape_data["h"]
                )
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h)
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(
                    *_hex_rgb(shape_data.get("fill", "#08766F"), (8, 118, 111))
                )
                shape.line.color.rgb = stroke_rgb
                shape.line.width = Pt(0.5)
            elif shape_type in {"rect", "round_rect"}:
                x, y, w, h = mapper.rect(
                    shape_data["x"], shape_data["y"], shape_data["w"], shape_data["h"]
                )
                mso_shape = MSO_SHAPE.ROUNDED_RECTANGLE if shape_type == "round_rect" else MSO_SHAPE.RECTANGLE
                shape = slide.shapes.add_shape(
                    mso_shape, Inches(x), Inches(y), Inches(w), Inches(h)
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(
                    *_hex_rgb(shape_data.get("fill", "#FFFFFF"), (255, 255, 255))
                )
                shape.line.color.rgb = stroke_rgb
                shape.line.width = Pt(shape_data.get("width_pt", 0.55))
                try:
                    if shape_type == "round_rect" and shape.adjustments:
                        shape.adjustments[0] = float(shape_data.get("radius_ratio", 0.18))
                except Exception:
                    pass
            else:
                continue
            shape.name = f"EditableShape_{index:02d}"
            counts["editable_shapes"] += 1
        except Exception as exc:
            logger.debug("Skipping native shape %s: %s", index, exc)

    # Movable faithful visual objects: icons, logos, complex badges/marks.
    for index, asset in enumerate(object_assets, start=1):
        x, y, w, h = mapper.rect(asset["x"], asset["y"], asset["w"], asset["h"])
        picture = slide.shapes.add_picture(
            asset["path"], Inches(x), Inches(y), Inches(w), Inches(h)
        )
        picture.name = f"MovableVisualObject_{index:02d}"
        counts["object_assets"] += 1

    image_width = mapper.image_width
    for index, region in enumerate(text_regions, start=1):
        text = str(region.get("text", "")).strip()
        if not text:
            continue
        x, y, w, h = mapper.rect(region["x"], region["y"], region["w"], region["h"])
        font_name = str(region.get("font_name") or _font_name_for_region(region))
        font_pt = mapper.font_points(
            region["median_word_h"],
            font_scale * _font_multiplier_for_region(region),
        )
        # Expand vertically for PowerPoint font metrics while preserving the
        # measured baseline.  Width padding prevents accidental wraps.
        y_adjust = min(h * 0.22, 0.06)
        box_y = max(0, y - y_adjust)
        box_h = min(SLIDE_H_IN - box_y, max(h * 1.48, font_pt / 72 * 1.35))
        box_w = min(SLIDE_W_IN - x, max(w + 0.04, w * 1.035))
        if region["role"] == "number" and len(text) <= 4:
            box_x = max(0, x - w * 0.22)
            box_w = min(SLIDE_W_IN - box_x, w * 1.44)
        else:
            box_x = max(0, x - 0.01)

        textbox = slide.shapes.add_textbox(
            Inches(box_x), Inches(box_y), Inches(box_w), Inches(box_h)
        )
        textbox.name = f"EditableText_{index:03d}_{region['role']}"
        frame = textbox.text_frame
        frame.clear()
        frame.word_wrap = False
        frame.margin_left = 0
        frame.margin_right = 0
        frame.margin_top = 0
        frame.margin_bottom = 0
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = frame.paragraphs[0]
        paragraph.alignment = _paragraph_alignment(
            _infer_alignment(region, image_width)
        )
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = 1.0

        # Preserve visible word-level color changes when OCR evidence is still
        # compatible with the corrected text.  Otherwise use the line color.
        color_runs = _colored_text_runs(region, text)
        if color_runs:
            for run_text, run_color in color_runs:
                run = paragraph.add_run()
                run.text = run_text
                run.font.size = Pt(font_pt)
                run.font.bold = _effective_bold(region)
                run.font.color.rgb = RGBColor(*_hex_rgb(run_color, (17, 17, 17)))
                _set_east_asian_font(run, font_name)
        else:
            run = paragraph.add_run()
            run.text = text
            run.font.size = Pt(font_pt)
            run.font.bold = _effective_bold(region)
            run.font.color.rgb = RGBColor(*_hex_rgb(region["color"], (17, 17, 17)))
            _set_east_asian_font(run, font_name)
        counts["editable_text"] += 1

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)
    return counts


def _powerpoint_installed() -> bool:
    if any(os.path.exists(path) for path in POWERPOINT_PATHS):
        return True
    return bool(shutil.which("powerpnt"))


def _powerpoint_process_ids() -> set[int]:
    try:
        import psutil
        return {
            process.pid
            for process in psutil.process_iter(["name"])
            if (process.info.get("name") or "").lower() == "powerpnt.exe"
        }
    except Exception:
        return set()


def _terminate_process_tree(process_id: int) -> None:
    try:
        import psutil
        process = psutil.Process(process_id)
        children = process.children(recursive=True)
        for child in children:
            child.terminate()
        process.terminate()
        _, alive = psutil.wait_procs(children + [process], timeout=3)
        for item in alive:
            item.kill()
    except Exception:
        return


def _render_powerpoint_once(
    pptx_path: str,
    output_path: str,
    width: int,
    height: int,
) -> None:
    """Worker entry point isolated in a subprocess.

    Some older Office installations successfully export the PNG and then hang
    during COM shutdown.  The parent process observes the completed file and
    safely cleans up only the PowerPoint process created for this render.
    """
    import pythoncom
    import win32com.client

    pythoncom.CoInitialize()
    app = None
    deck = None
    try:
        app = win32com.client.DispatchEx("PowerPoint.Application")
        app.DisplayAlerts = 0
        deck = app.Presentations.Open(
            os.path.abspath(pptx_path),
            ReadOnly=True,
            Untitled=False,
            WithWindow=False,
        )
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        deck.Slides(1).Export(os.path.abspath(output_path), "PNG", width, height)
    finally:
        if deck is not None:
            deck.Close()
        if app is not None:
            app.Quit()
        pythoncom.CoUninitialize()


def _render_libreoffice_once(
    pptx_path: str,
    output_path: str,
    width: int,
    height: int,
) -> tuple[bool, str | None]:
    """Render slide 1 through LibreOffice + Poppler when PowerPoint COM is blocked."""
    soffice = shutil.which("soffice") or shutil.which("soffice.exe")
    pdftoppm = shutil.which("pdftoppm") or shutil.which("pdftoppm.exe")
    if not soffice or not pdftoppm:
        missing = []
        if not soffice:
            missing.append("soffice")
        if not pdftoppm:
            missing.append("pdftoppm")
        return False, "LibreOffice fallback unavailable: missing " + ", ".join(missing)

    output = Path(output_path)
    output.unlink(missing_ok=True)
    with tempfile.TemporaryDirectory(prefix="ppt_precise_render_") as tmp:
        tmp_dir = Path(tmp)
        src = tmp_dir / "deck.pptx"
        shutil.copy2(pptx_path, src)
        try:
            subprocess.run(
                [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(tmp_dir), str(src)],
                check=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
                encoding="utf-8",
                errors="replace",
                timeout=90,
            )
            pdf_path = tmp_dir / "deck.pdf"
            if not pdf_path.exists():
                return False, "LibreOffice fallback failed: PDF not produced"
            preview_prefix = tmp_dir / "page"
            subprocess.run(
                [pdftoppm, "-png", "-singlefile", "-r", "96", str(pdf_path), str(preview_prefix)],
                check=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
                encoding="utf-8",
                errors="replace",
                timeout=90,
            )
            rendered = tmp_dir / "page.png"
            if not rendered.exists():
                return False, "LibreOffice fallback failed: PNG not produced"
            img = cv2.imread(str(rendered), cv2.IMREAD_COLOR)
            if img is None:
                return False, "LibreOffice fallback failed: rendered PNG unreadable"
            img = cv2.resize(img, (int(width), int(height)), interpolation=cv2.INTER_LANCZOS4)
            output.parent.mkdir(parents=True, exist_ok=True)
            cv2.imencode(".png", img)[1].tofile(str(output))
            return True, None
        except Exception as exc:
            return False, f"LibreOffice fallback failed: {exc.__class__.__name__}: {exc}"


def render_pptx_to_png(
    pptx_path: str,
    output_path: str,
    width: int = DEFAULT_RENDER_W,
    height: int = DEFAULT_RENDER_H,
) -> tuple[bool, str | None]:
    """Render slide 1 using PowerPoint COM, falling back to LibreOffice if available."""
    if not _powerpoint_installed():
        return _render_libreoffice_once(pptx_path, output_path, width, height)
    existing_powerpoint = _powerpoint_process_ids()
    output = Path(output_path)
    output.unlink(missing_ok=True)
    code = (
        "from app.services.precise_reconstruction import _render_powerpoint_once;"
        f"_render_powerpoint_once({pptx_path!r},{output_path!r},{int(width)},{int(height)})"
    )
    creation_flags = getattr(subprocess, "CREATE_NO_WINDOW", 0)
    try:
        worker = subprocess.Popen(
            [sys.executable, "-c", code],
            cwd=str(Path(__file__).resolve().parents[2]),
            stdout=subprocess.DEVNULL,
            stderr=subprocess.PIPE,
            text=True,
            encoding="utf-8",
            errors="replace",
            creationflags=creation_flags,
        )
        deadline = time.monotonic() + 75
        last_size = -1
        stable_since = None
        while time.monotonic() < deadline:
            return_code = worker.poll()
            if output.exists() and output.stat().st_size > 1000:
                size = output.stat().st_size
                if size == last_size:
                    stable_since = stable_since or time.monotonic()
                    if time.monotonic() - stable_since >= 1.0:
                        break
                else:
                    last_size = size
                    stable_since = None
            if return_code is not None:
                break
            time.sleep(0.25)

        if worker.poll() is None:
            _terminate_process_tree(worker.pid)
        stderr = ""
        try:
            _, stderr = worker.communicate(timeout=3)
        except subprocess.TimeoutExpired:
            _terminate_process_tree(worker.pid)

        new_powerpoint = _powerpoint_process_ids() - existing_powerpoint
        for process_id in new_powerpoint:
            _terminate_process_tree(process_id)

        if output.exists() and output.stat().st_size > 1000:
            return True, None
        powerpoint_error = stderr.strip() or "PowerPoint render timed out"
    except Exception as exc:
        logger.warning("PowerPoint rendering failed: %s", exc)
        powerpoint_error = f"{exc.__class__.__name__}: {exc}"

    fallback_ok, fallback_error = _render_libreoffice_once(pptx_path, output_path, width, height)
    if fallback_ok:
        return True, None
    return False, f"{powerpoint_error}\n{fallback_error}"


def _source_canvas(
    source_path: str,
    mapper: CoordinateMapper,
    width: int,
    height: int,
) -> np.ndarray:
    source = _read_cv_image(source_path)
    canvas = np.full((height, width, 3), 255, dtype=np.uint8)
    if mapper.mode == "stretch":
        return cv2.resize(source, (width, height), interpolation=cv2.INTER_LANCZOS4)
    x = round(mapper.ox / SLIDE_W_IN * width)
    y = round(mapper.oy / SLIDE_H_IN * height)
    w = round(mapper.image_width * mapper.sx / SLIDE_W_IN * width)
    h = round(mapper.image_height * mapper.sy / SLIDE_H_IN * height)
    resized = cv2.resize(source, (w, h), interpolation=cv2.INTER_LANCZOS4)
    canvas[y:y + h, x:x + w] = resized
    return canvas


def compare_images(
    source_path: str,
    rendered_path: str,
    mapper: CoordinateMapper,
    text_regions: list[dict[str, Any]] | None = None,
    heatmap_path: str | None = None,
    comparison_path: str | None = None,
) -> dict[str, Any]:
    from skimage.metrics import structural_similarity

    try:
        rendered = _read_cv_image(rendered_path)
    except (ValueError, OSError):
        return {"error": "Cannot read rendered preview"}
    height, width = rendered.shape[:2]
    source = _source_canvas(source_path, mapper, width, height)
    if source.shape != rendered.shape:
        rendered = cv2.resize(rendered, (source.shape[1], source.shape[0]))

    gray_source = cv2.cvtColor(source, cv2.COLOR_BGR2GRAY)
    gray_rendered = cv2.cvtColor(rendered, cv2.COLOR_BGR2GRAY)
    full_ssim, ssim_map = structural_similarity(
        gray_source, gray_rendered, full=True, data_range=255
    )
    edge_source = cv2.Canny(gray_source, 60, 150)
    edge_rendered = cv2.Canny(gray_rendered, 60, 150)
    edge_ssim = structural_similarity(
        edge_source, edge_rendered, data_range=255
    )
    diff = cv2.absdiff(source, rendered)
    mean_diff = float(np.mean(cv2.cvtColor(diff, cv2.COLOR_BGR2GRAY)))

    text_scores: list[float] = []
    for region in text_regions or []:
        x, y, w, h = mapper.render_rect(
            region["x"], region["y"], region["w"], region["h"], width, height
        )
        padding = max(3, round(h * 0.35))
        x1, y1 = max(0, x - padding), max(0, y - padding)
        x2, y2 = min(width, x + w + padding), min(height, y + h + padding)
        a = gray_source[y1:y2, x1:x2]
        b = gray_rendered[y1:y2, x1:x2]
        if min(a.shape[:2]) < 7:
            continue
        try:
            text_scores.append(float(structural_similarity(a, b, data_range=255)))
        except ValueError:
            continue
    text_ssim = float(np.mean(text_scores)) if text_scores else float(full_ssim)
    composite = (
        float(full_ssim) * 0.42
        + float(edge_ssim) * 0.22
        + text_ssim * 0.30
        + max(0.0, 1.0 - mean_diff / 255.0) * 0.06
    )

    if heatmap_path:
        error = np.clip((1.0 - ssim_map) * 255 * 2.2, 0, 255).astype(np.uint8)
        heatmap = cv2.applyColorMap(error, cv2.COLORMAP_JET)
        overlay = cv2.addWeighted(rendered, 0.58, heatmap, 0.42, 0)
        cv2.imencode(".png", overlay)[1].tofile(heatmap_path)

    if comparison_path:
        label_h = 44
        comparison = np.full((height + label_h, width * 2 + 16, 3), 238, dtype=np.uint8)
        comparison[label_h:, :width] = source
        comparison[label_h:, width + 16:] = rendered
        cv2.putText(
            comparison, "SOURCE", (16, 30), cv2.FONT_HERSHEY_SIMPLEX,
            0.8, (8, 118, 111), 2, cv2.LINE_AA,
        )
        cv2.putText(
            comparison, "PPTX RENDER", (width + 32, 30), cv2.FONT_HERSHEY_SIMPLEX,
            0.8, (8, 118, 111), 2, cv2.LINE_AA,
        )
        cv2.imencode(".png", comparison)[1].tofile(comparison_path)

    return {
        "ssim": round(float(full_ssim), 4),
        "edge_ssim": round(float(edge_ssim), 4),
        "text_ssim": round(text_ssim, 4),
        "mean_diff": round(mean_diff, 2),
        "score": round(float(composite), 4),
        "quality": (
            "strong" if composite >= 0.86
            else "good" if composite >= 0.76
            else "fair" if composite >= 0.64
            else "needs_review"
        ),
    }


def _safe_url(filename: str) -> str:
    return f"/api/skills/download/{filename}"


async def reconstruct(
    image_path: str,
    session_id: str = "",
    output_dir: str | None = None,
    refine_with_llm: bool = True,
    render_preview: bool = True,
    max_calibration_passes: int = 2,
) -> dict[str, Any]:
    """Reconstruct one visual mockup as a single-slide editable 16:9 PPTX."""
    source = Path(image_path)
    if not source.exists():
        return {"error": f"Image not found: {source.name}"}

    out_dir = Path(output_dir) if output_dir else _default_output_dir()
    out_dir = out_dir.resolve()
    out_dir.mkdir(parents=True, exist_ok=True)
    sid = (session_id or uuid.uuid4().hex)[:8]
    run_id = uuid.uuid4().hex[:6]
    stem = f"ppt_precise_{sid}_{run_id}"
    work_root = out_dir.parent / "precise_work" / stem
    work_root.mkdir(parents=True, exist_ok=True)

    try:
        layers = await asyncio.to_thread(analyze_layers, str(source))
    except Exception as exc:
        logger.exception("Precise reconstruction analysis failed")
        return {"error": f"Analysis failed: {exc}"}

    mapper = CoordinateMapper(layers["image_width"], layers["image_height"])
    editable_text = layers["editable_text"]
    llm_report = {"used": False, "reason": "not_requested"}
    if refine_with_llm and editable_text:
        editable_text, llm_report = await _refine_text_with_deepseek(
            editable_text, layers["image_width"], layers["image_height"]
        )
    # Reassign stable ordering after optional filtering.
    editable_text.sort(key=lambda item: (item["y"], item["x"]))
    for index, region in enumerate(editable_text, start=1):
        region["id"] = f"text_{index:03d}"

    cv_img = _read_cv_image(str(source))
    _augment_text_styles(cv_img, editable_text)
    cleaned, text_mask = await asyncio.to_thread(
        _make_clean_visual_base, cv_img, editable_text
    )
    clean_path = work_root / "text_cleaned_visual.png"
    mask_path = work_root / "editable_text_mask.png"
    cv2.imencode(".png", cleaned)[1].tofile(str(clean_path))
    cv2.imencode(".png", text_mask)[1].tofile(str(mask_path))
    object_assets, object_mask = await asyncio.to_thread(
        _extract_visual_object_assets, cleaned, text_mask, work_root, layers["native_shapes"]
    )
    object_mask_path = work_root / "visual_object_mask.png"
    cv2.imencode(".png", object_mask)[1].tofile(str(object_mask_path))
    tile_base = await asyncio.to_thread(
        _remove_object_assets_from_base, cleaned, object_mask
    )
    cv2.imencode(".png", tile_base)[1].tofile(str(work_root / "tile_visual_base.png"))
    tile_assets = await asyncio.to_thread(_tile_visual_base, tile_base, work_root)
    assets = tile_assets + object_assets

    final_pptx = out_dir / f"{stem}.pptx"
    final_preview = out_dir / f"{stem}_preview.png"
    final_comparison = out_dir / f"{stem}_comparison.png"
    final_heatmap = out_dir / f"{stem}_diff.png"
    final_report = out_dir / f"{stem}_report.json"

    # A small font-scale search is more reliable than assuming OCR glyph height
    # maps perfectly across every installed font/rendering engine.
    scale_pool = [1.0, 0.94, 1.06, 0.90, 1.10, 0.97, 1.03]
    scales = scale_pool[: max(1, max_calibration_passes if render_preview else 1)]
    candidates: list[dict[str, Any]] = []
    build_counts = {}
    for candidate_index, scale in enumerate(scales, start=1):
        candidate_pptx = work_root / f"candidate_{candidate_index}.pptx"
        build_counts = await asyncio.to_thread(
            _build_pptx,
            str(candidate_pptx),
            mapper,
            layers["background"],
            assets,
            layers["native_shapes"],
            editable_text,
            scale,
        )
        candidate = {
            "scale": scale,
            "pptx": str(candidate_pptx),
            "score": -1.0,
            "metrics": {},
        }
        if render_preview:
            candidate_png = work_root / f"candidate_{candidate_index}.png"
            rendered, render_error = await asyncio.to_thread(
                render_pptx_to_png, str(candidate_pptx), str(candidate_png)
            )
            candidate["render_error"] = render_error
            if rendered:
                metrics = await asyncio.to_thread(
                    compare_images,
                    str(source),
                    str(candidate_png),
                    mapper,
                    editable_text,
                    None,
                    None,
                )
                candidate["metrics"] = metrics
                candidate["score"] = metrics.get("score", -1.0)
                candidate["preview"] = str(candidate_png)
        candidates.append(candidate)

    rendered_candidates = [item for item in candidates if item["score"] >= 0]
    best = max(rendered_candidates, key=lambda item: item["score"]) if rendered_candidates else candidates[0]
    shutil.copy2(best["pptx"], final_pptx)

    render_error = best.get("render_error")
    metrics: dict[str, Any] = best.get("metrics", {})
    if best.get("preview") and os.path.exists(best["preview"]):
        shutil.copy2(best["preview"], final_preview)
        metrics = await asyncio.to_thread(
            compare_images,
            str(source),
            str(final_preview),
            mapper,
            editable_text,
            str(final_heatmap),
            str(final_comparison),
        )
    elif render_preview:
        rendered, render_error = await asyncio.to_thread(
            render_pptx_to_png, str(final_pptx), str(final_preview)
        )
        if rendered:
            metrics = await asyncio.to_thread(
                compare_images,
                str(source),
                str(final_preview),
                mapper,
                editable_text,
                str(final_heatmap),
                str(final_comparison),
            )

    report = {
        "source": source.name,
        "pipeline_version": PIPELINE_VERSION,
        "slide_size": "16:9",
        "coordinate_mode": mapper.mode,
        "model_backend": "existing Claude/DeepSeek endpoint (text metadata refinement only)",
        "llm_refinement": llm_report,
        "editable": {
            "text_boxes": build_counts.get("editable_text", 0),
            "native_shapes": build_counts.get("editable_shapes", 0),
            "text_roles": sorted(set(region["role"] for region in editable_text)),
        },
        "image_assets": {
            "visual_tiles": build_counts.get("visual_assets", 0),
            "movable_objects": build_counts.get("object_assets", 0),
            "strategy": (
                "native PowerPoint background and geometry first; movable transparent crops only for icons/logos/complex visual marks; no full-slide bitmap"
            ),
            "background_mode": os.getenv("PPT_PRECISE_BACKGROUND_MODE", "native").strip().lower(),
        },
        "raster_text": [
            {
                "text": region["text"],
                "reason": "tiny low-confidence text embedded in a complex visual region",
            }
            for region in layers["raster_text"]
        ],
        "tradeoffs": [
            "Complex visuals, icons, photos, textures, subtle gradients and unconfident marks remain raster crops to preserve design quality.",
            "Recognized cards, pills, dividers and long rules are rebuilt as native editable PowerPoint shapes.",
            "Recognized non-text icons/logos are transparent movable image objects instead of low-quality hand-drawn shapes.",
            "Main text is rebuilt as native PowerPoint text boxes; font, boldness and mixed colors are inferred from local source pixels.",
            "Chart values are not invented; unreadable chart detail remains in the visual asset layer.",
        ],
        "ocr": {
            "engine": layers.get("ocr_engine", "unknown"),
            "words": layers["ocr_word_count"],
            "editable_regions": len(editable_text),
            "raster_regions": len(layers["raster_text"]),
        },
        "visual_qa": {
            **metrics,
            "font_scale_candidates": [
                {
                    "scale": item["scale"],
                    "score": item["score"],
                    "render_error": item.get("render_error"),
                }
                for item in candidates
            ],
            "selected_font_scale": best["scale"],
            "render_error": render_error,
        },
    }
    final_report.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")

    result = {
        "filename": final_pptx.name,
        "pipeline_version": PIPELINE_VERSION,
        "path": str(final_pptx),
        "url": _safe_url(final_pptx.name),
        "preview_filename": final_preview.name if final_preview.exists() else None,
        "preview_path": str(final_preview) if final_preview.exists() else None,
        "preview_url": _safe_url(final_preview.name) if final_preview.exists() else None,
        "comparison_filename": final_comparison.name if final_comparison.exists() else None,
        "comparison_path": str(final_comparison) if final_comparison.exists() else None,
        "comparison_url": _safe_url(final_comparison.name) if final_comparison.exists() else None,
        "diff_filename": final_heatmap.name if final_heatmap.exists() else None,
        "diff_url": _safe_url(final_heatmap.name) if final_heatmap.exists() else None,
        "report_filename": final_report.name,
        "report_path": str(final_report),
        "report_url": _safe_url(final_report.name),
        "element_count": sum(build_counts.values()),
        "breakdown": build_counts,
        "report": report,
    }
    return result
