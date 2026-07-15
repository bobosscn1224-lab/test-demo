"""Async batch PPTX conversion with state-machine progress tracking.

Jobs are queued in-memory and processed page-by-page in the background.
The frontend polls GET /api/skills/batch-status/{job_id} for progress.
"""
from __future__ import annotations

import asyncio
import logging
import os
import uuid
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

# Global job store (survives between requests but not server restarts)
_jobs: dict[str, dict[str, Any]] = {}


def create_job(image_paths: list[str], output_dir: str) -> str:
    """Create a new batch job and start background processing."""
    job_id = uuid.uuid4().hex[:12]
    now = datetime.now(timezone.utc).isoformat()
    _jobs[job_id] = {
        "job_id": job_id,
        "status": "queued",
        "total": len(image_paths),
        "current": 0,
        "current_file": "",
        "progress": 0,
        "results": [],
        "errors": [],
        "combined_pptx_path": None,
        "combined_pptx_url": None,
        "combined_pptx_filename": None,
        "created_at": now,
    }
    # Launch background task
    asyncio.create_task(_process_job(job_id, image_paths, output_dir))
    return job_id


def get_job(job_id: str) -> dict[str, Any] | None:
    return _jobs.get(job_id)


async def _process_job(job_id: str, image_paths: list[str], output_dir: str) -> None:
    """Process all images sequentially, updating job state after each."""
    job = _jobs[job_id]
    job["status"] = "processing"
    output_dir_path = Path(output_dir)
    output_dir_path.mkdir(parents=True, exist_ok=True)

    from app.services.deckweaver_service import convert_image_to_pptx

    individual_pptx: list[str] = []
    valid_paths = [p for p in image_paths if os.path.exists(p)]

    if not valid_paths:
        job["status"] = "failed"
        job["errors"].append({"page": 0, "file": "", "error": "No valid images found"})
        return

    total = len(valid_paths)
    for idx, img_path in enumerate(valid_paths, start=1):
        job["current"] = idx
        job["current_file"] = os.path.basename(img_path)
        job["progress"] = int((idx - 1) / total * 90)  # 0-90% for conversion phase

        try:
            r = await convert_image_to_pptx(
                img_path,
                session_id=f"{job_id}_p{idx}",
                output_dir=str(output_dir_path),
            )
        except Exception as exc:
            job["errors"].append({"page": idx, "file": os.path.basename(img_path), "error": str(exc)})
            continue

        if r.get("error"):
            job["errors"].append({"page": idx, "file": os.path.basename(img_path), "error": r["error"]})
            continue

        individual_pptx.append(r["path"])
        job["results"].append({
            "page": idx,
            "file": os.path.basename(img_path),
            "text_items": r.get("report", {}).get("text_items", 0),
            "status": "done",
        })

    if not individual_pptx:
        job["status"] = "failed"
        return

    # Merge phase
    job["status"] = "merging"
    job["progress"] = 92
    job["current_file"] = "合并中..."

    combined_name = f"ppt_batch_{job_id}.pptx"
    combined_path = output_dir_path / combined_name
    try:
        await asyncio.to_thread(_merge_pptx_files, individual_pptx, str(combined_path))
        job["combined_pptx_path"] = str(combined_path)
        job["combined_pptx_url"] = f"/api/skills/download/{combined_name}"
        job["combined_pptx_filename"] = combined_name
    except Exception as exc:
        job["errors"].append({"page": 0, "file": "", "error": f"Merge failed: {exc}"})
        job["status"] = "failed"
        return

    job["status"] = "completed"
    job["progress"] = 100
    job["current"] = total
    job["current_file"] = "完成"


def _merge_pptx_files(source_paths: list[str], output_path: str) -> None:
    """Merge single-slide PPTX files into one combined PPTX. Preserves images."""
    from pptx import Presentation
    from pptx.util import Inches
    from copy import deepcopy
    import io

    merged = Presentation()
    merged.slide_width = Inches(13.333)
    merged.slide_height = Inches(7.5)

    for src_path in source_paths:
        src = Presentation(src_path)
        for slide in src.slides:
            blank_layout = merged.slide_layouts[6]
            new_slide = merged.slides.add_slide(blank_layout)
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture
                    try:
                        blob = shape.image.blob
                        buf = io.BytesIO(blob)
                        new_slide.shapes.add_picture(buf, shape.left, shape.top, shape.width, shape.height)
                    except Exception:
                        try:
                            new_slide.shapes._spTree.append(deepcopy(shape._element))
                        except Exception:
                            pass
                else:
                    try:
                        new_slide.shapes._spTree.append(deepcopy(shape._element))
                    except Exception:
                        pass

    merged.save(output_path)
