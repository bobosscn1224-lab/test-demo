"""Image-to-Editable-PPTX — PPT Object-Principled Element Reconstruction.

Strategy: Each PPT element is a complete, independent, selectable object.

Element types and their reconstruction:
  - ICONS: complete image crops (connected-component detection, NOT entropy)
  - TEXT: editable text boxes with real OCR text (PaddleOCR + Tesseract)
  - SHAPES: native PPTX lines/rectangles (OpenCV Hough + contour)
  - CARDS: native rectangles with fill (contour detection)
  - BACKGROUND: inpainted clean image or solid fill (cv2.inpaint)

All primary work is local tools — no LLM for recognition.
"""
from __future__ import annotations

import logging
import os
import uuid
from collections import defaultdict

import cv2
import numpy as np
from pptx.oxml.ns import qn

logger = logging.getLogger(__name__)

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
TESSERACT_PATH = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
MIN_CONFIDENCE = 30
MIN_ELEMENT_AREA = 0.001  # fraction of image area


async def convert_image_to_pptx(
    image_path: str,
    output_dir: str | None = None,
    session_id: str = "",
) -> dict | None:
    import asyncio
    return await asyncio.to_thread(_convert_sync, image_path, output_dir, session_id)


# ═══════════════════════════════════════════════════════════════════════════
# Main conversion
# ═══════════════════════════════════════════════════════════════════════════

def _convert_sync(image_path: str, output_dir: str | None, session_id: str) -> dict | None:
    """PPT-object-principled reconstruction.

    Every element (icon, text box, shape, card) is an independent PPTX object.
    Nothing is split across boundaries — icons are whole, text is editable.
    """
    try:
        import pytesseract
        from PIL import Image
        from pptx import Presentation
        from pptx.util import Inches

        if os.path.exists(TESSERACT_PATH):
            pytesseract.pytesseract.tesseract_cmd = TESSERACT_PATH

        # ── Load image ──
        img = Image.open(image_path)
        iw, ih = img.size
        buf = np.fromfile(image_path, dtype=np.uint8)
        cv_img = cv2.imdecode(buf, cv2.IMREAD_COLOR)
        if cv_img is None:
            cv_img = np.array(img.convert("RGB"))[:, :, ::-1]

        sx = SLIDE_W_IN / iw
        sy = SLIDE_H_IN / ih

        logger.info("PPTX: %s (%dx%d)", os.path.basename(image_path), iw, ih)

        # ── Step 1: Text regions (PaddleOCR + Tesseract) ──
        text_regions = _detect_text(img, iw, ih)

        # ── Step 2: Shapes (OpenCV) ──
        shapes = _detect_shapes(cv_img, iw, ih, sx, sy)

        # ── Step 3: Visual elements via connected components ──
        #   Build mask: content = everything NOT background
        #   Subtract text + shapes → remaining = visual elements
        #   Connected components → each is ONE image crop
        visual_elements = _detect_visual_elements(cv_img, text_regions, shapes, iw, ih, sx, sy)

        # ── Step 4: Clean background via inpainting ──
        bg_img, bg_type, bg_color = _make_clean_background(cv_img, text_regions, visual_elements)

        # ── Build PPTX ──
        out_dir = output_dir or os.path.abspath(os.path.join("data", "outputs"))
        os.makedirs(out_dir, exist_ok=True)
        sid = session_id[:8] if session_id else uuid.uuid4().hex[:8]

        from pptx import Presentation
        prs = Presentation()
        prs.slide_width = Inches(SLIDE_W_IN)
        prs.slide_height = Inches(SLIDE_H_IN)
        blank = prs.slide_layouts[6]

        # Page 1: Reference
        s1 = prs.slides.add_slide(blank)
        s1.shapes.add_picture(image_path, Inches(0), Inches(0), Inches(SLIDE_W_IN), Inches(SLIDE_H_IN))

        # Page 2: Reconstruction
        s2 = prs.slides.add_slide(blank)
        asset_dir = os.path.join(out_dir, f"assets_{sid}")
        os.makedirs(asset_dir, exist_ok=True)

        # Layer 0: Clean background
        if bg_type == "solid":
            _set_slide_bg(s2, *bg_color)
        elif bg_img is not None:
            bg_path = os.path.join(asset_dir, "bg.png")
            cv2.imwrite(bg_path, bg_img)
            s2.shapes.add_picture(bg_path, Inches(0), Inches(0), Inches(SLIDE_W_IN), Inches(SLIDE_H_IN))
        else:
            _set_slide_bg(s2, 255, 255, 255)

        # Layer 1: Visual elements (icons, photos — each as ONE complete image crop)
        placed_visuals = _place_visual_elements(s2, cv_img, visual_elements, asset_dir)

        # Layer 2: Native shapes
        _place_shapes(s2, shapes)

        # Layer 3: Editable text boxes
        _place_text_boxes(s2, text_regions, iw, ih, sx, sy)

        # ── Save ──
        fname = f"ppt_hybrid_{sid}.pptx"
        fpath = os.path.join(out_dir, fname)
        prs.save(fpath)

        total_chars = sum(len(t["text"]) for t in text_regions)
        return {
            "filename": fname, "path": fpath,
            "url": f"/api/skills/download/{fname}",
            "pages": 2,
            "chars_extracted": total_chars,
            "text_regions": len(text_regions),
            "shapes_detected": len(shapes),
            "visual_elements": len(visual_elements),
            "editable_elements": len(text_regions) + len(shapes) + len(visual_elements),
            "report": {
                "method": "PPT-object-principled: connected-component icons + OCR text + native shapes + inpainted bg",
                "background_type": bg_type,
                "editable_text_boxes": len(text_regions),
                "editable_shapes": len(shapes),
                "image_elements": len(visual_elements),
                "text_sample": [t["text"][:60] for t in text_regions[:5]],
                "shapes_sample": [f"{s['type']}:{s.get('subtype','')}" for s in shapes[:5]],
                "visual_sample": [f"{v['category']} ({v['w']:.1f}x{v['h']:.1f}in)" for v in visual_elements[:5]],
            },
        }
    except Exception as exc:
        logger.error("PPTX failed: %s", exc)
        import traceback
        traceback.print_exc()
        return None


# ═══════════════════════════════════════════════════════════════════════════
# STEP 1: Text Detection (PaddleOCR + Tesseract)
# ═══════════════════════════════════════════════════════════════════════════

_PADDLE = None


def _get_paddle():
    global _PADDLE
    if _PADDLE is None:
        try:
            from paddleocr import PaddleOCR
            _PADDLE = PaddleOCR(lang="ch", ocr_version="PP-OCRv4", use_doc_orientation_classify=False, use_doc_unwarping=False)
            logger.info("PaddleOCR ready")
        except Exception as exc:
            logger.warning("PaddleOCR unavailable: %s", exc)
            _PADDLE = False
    return _PADDLE if _PADDLE is not False else None


def _detect_text(img, iw: int, ih: int) -> list[dict]:
    """PaddleOCR (Chinese primary) + Tesseract (English supplement)."""
    paddle = _get_paddle()
    raw_lines: list[dict] = []

    # PaddleOCR
    if paddle:
        try:
            arr = np.array(img.convert("RGB"))
            for res in paddle.predict(arr):
                polys = res.get("dt_polys", []) if isinstance(res, dict) else getattr(res, "dt_polys", [])
                texts = res.get("rec_texts", []) if isinstance(res, dict) else getattr(res, "rec_texts", [])
                scores = res.get("rec_scores", []) if isinstance(res, dict) else getattr(res, "rec_scores", [])
                for poly, text, score in zip(polys or [], texts or [], scores or []):
                    if float(score) < 0.5 or not text or not text.strip():
                        continue
                    poly_arr = np.asarray(poly)
                    if poly_arr.ndim != 2 or poly_arr.shape[1] != 2:
                        continue
                    xs, ys = poly_arr[:, 0], poly_arr[:, 1]
                    raw_lines.append({
                        "text": text.strip(), "x": int(min(xs)), "y": int(min(ys)),
                        "w": int(max(xs) - min(xs)), "h": int(max(ys) - min(ys)),
                        "source": "paddleocr",
                    })
        except Exception as exc:
            logger.warning("PaddleOCR error: %s", exc)

    # Tesseract
    import pytesseract
    try:
        data = pytesseract.image_to_data(img, lang="chi_sim+eng", output_type=pytesseract.Output.DICT, config="--psm 6")
        for i in range(len(data["text"])):
            t = (data["text"][i] or "").strip()
            if not t:
                continue
            conf = int(data["conf"][i]) if data["conf"][i] != "-1" else 100
            if conf < MIN_CONFIDENCE:
                continue
            raw_lines.append({
                "text": t, "x": int(data["left"][i]), "y": int(data["top"][i]),
                "w": int(data["width"][i]), "h": int(data["height"][i]),
                "source": "tesseract",
            })
    except Exception as exc:
        logger.warning("Tesseract error: %s", exc)

    if not raw_lines:
        return []

    # ── Merge strategy: PaddleOCR is the primary engine ──
    # PaddleOCR handles both Chinese AND English text well.
    # Tesseract is ONLY used as complete fallback (PaddleOCR found nothing).
    # We NEVER merge Tesseract with PaddleOCR because Tesseract's Chinese
    # output is character-level and creates fragmentation.
    paddle_lines = [l for l in raw_lines if l.get("source") == "paddleocr"]

    if paddle_lines:
        # PaddleOCR succeeded — use it exclusively
        merged = paddle_lines
    else:
        # PaddleOCR failed completely — fall back to Tesseract only
        merged = [l for l in raw_lines if l.get("source") == "tesseract"]

    # Group into blocks
    return _group_into_blocks(merged, iw, ih)


def _group_into_blocks(lines: list[dict], iw: int, ih: int) -> list[dict]:
    """Group line-level detections into semantic text blocks."""
    if not lines:
        return []
    lines.sort(key=lambda r: r["y"])

    blocks = []
    current = [lines[0]]
    for i in range(1, len(lines)):
        prev, cur = lines[i - 1], lines[i]
        gap = cur["y"] - (prev["y"] + prev["h"])
        avg_h = max((prev["h"] + cur["h"]) / 2, 1)
        x_overlap = max(0, min(prev["x"] + prev["w"], cur["x"] + cur["w"]) - max(prev["x"], cur["x"]))
        overlap_frac = x_overlap / max(prev["w"], cur["w"], 1)

        # Lines belong to same block if vertically close AND horizontally overlapping
        if gap < avg_h * 2.5 and (overlap_frac > 0.15 or abs(cur["x"] - prev["x"]) < 80):
            current.append(cur)
        else:
            blocks.append(_merge_block(current, iw, ih))
            current = [cur]
    blocks.append(_merge_block(current, iw, ih))
    return blocks


def _merge_block(lines: list[dict], iw: int, ih: int) -> dict:
    text = "\n".join(l["text"] for l in lines)
    min_x = min(l["x"] for l in lines)
    min_y = min(l["y"] for l in lines)
    max_x = max(l["x"] + l["w"] for l in lines)
    max_y = max(l["y"] + l["h"] for l in lines)
    avg_h = sum(l["h"] for l in lines) / len(lines)
    rel_y = min_y / max(ih, 1)
    font_est = avg_h * (SLIDE_H_IN / ih) * 72 * 0.75
    return {
        "text": text.strip(), "x": min_x, "y": min_y,
        "w": max_x - min_x, "h": max_y - min_y,
        "font_size_est": round(max(8, min(72, font_est)), 1),
        "is_title": rel_y < 0.25 or avg_h > 35,
        "is_footer": rel_y > 0.85,
        "num_lines": len(lines),
    }


# ═══════════════════════════════════════════════════════════════════════════
# STEP 2: Shape Detection (OpenCV)
# ═══════════════════════════════════════════════════════════════════════════

def _detect_shapes(cv_img, iw, ih, sx, sy) -> list[dict]:
    """Detect lines and rectangles via OpenCV."""
    h, w = cv_img.shape[:2]
    min_len = min(w, h) * 0.05
    results = []

    try:
        gray = cv2.cvtColor(cv_img, cv2.COLOR_BGR2GRAY)

        # Hough lines
        edges = cv2.Canny(gray, 50, 150, apertureSize=3)
        hough = cv2.HoughLinesP(edges, 1, np.pi / 180, threshold=120,
                                minLineLength=int(min_len * 1.5), maxLineGap=15)
        if hough is not None:
            for line in hough:
                x1, y1, x2, y2 = line[0]
                length = np.sqrt((x2 - x1) ** 2 + (y2 - y1) ** 2)
                if length < min_len:
                    continue
                angle = abs(np.arctan2(y2 - y1, x2 - x1) * 180 / np.pi)
                lt = "horizontal" if angle < 5 or angle > 175 else ("vertical" if 85 < angle < 95 else "diagonal")
                results.append({
                    "type": "line", "subtype": lt,
                    "x1": round(x1 * sx, 4), "y1": round(y1 * sy, 4),
                    "x2": round(x2 * sx, 4), "y2": round(y2 * sy, 4),
                    "width_pt": max(0.5, min(4, ((y2 - y1) if lt == "horizontal" else (x2 - x1)) * sy * 72 / 10)),
                })

        # Rectangles via contour approximation
        blurred = cv2.GaussianBlur(gray, (5, 5), 0)
        thresh = cv2.adaptiveThreshold(blurred, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
                                       cv2.THRESH_BINARY_INV, 11, 2)
        contours, _ = cv2.findContours(thresh, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
        for cnt in contours:
            area = cv2.contourArea(cnt)
            if area < w * h * MIN_ELEMENT_AREA:
                continue
            peri = cv2.arcLength(cnt, True)
            approx = cv2.approxPolyDP(cnt, 0.02 * peri, True)
            if len(approx) == 4:
                rx, ry, rw, rh = cv2.boundingRect(cnt)
                if cv2.contourArea(cnt) / max(rw * rh, 1) > 0.7 and rw > 30 and rh > 15:
                    roi = cv_img[ry:ry + rh, rx:rx + rw]
                    avg_bgr = np.mean(roi, axis=(0, 1)).astype(int)
                    edge_bgr = _edge_color(cv_img, rx, ry, rw, rh)
                    results.append({
                        "type": "rectangle", "x": round(rx * sx, 4), "y": round(ry * sy, 4),
                        "w": round(rw * sx, 4), "h": round(rh * sy, 4),
                        "fill": _b2r(avg_bgr), "stroke": _b2r(edge_bgr),
                    })
    except Exception as exc:
        logger.warning("Shape detection: %s", exc)

    return _dedup(results)[:30]


# ═══════════════════════════════════════════════════════════════════════════
# STEP 3: Visual Element Detection (connected components, NOT entropy)
# ═══════════════════════════════════════════════════════════════════════════

def _detect_visual_elements(cv_img, text_regions, shapes, iw, ih, sx, sy) -> list[dict]:
    """Find icons/photos/charts as complete connected regions.

    Approach (NOT entropy — connected components):
      1. Build a "content mask" via edge detection
      2. Subtract text areas and shape areas
      3. Find connected components → each is ONE visual element
      4. Filter by size and aspect ratio
    """
    h, w = cv_img.shape[:2]
    elements = []

    try:
        gray = cv2.cvtColor(cv_img, cv2.COLOR_BGR2GRAY)

        # Edge-based content detection
        edges = cv2.Canny(gray, 40, 120, apertureSize=3)

        # Dilate to connect nearby edges → solid regions
        kernel = np.ones((7, 7), np.uint8)
        content_mask = cv2.dilate(edges, kernel, iterations=2)
        content_mask = cv2.morphologyEx(content_mask, cv2.MORPH_CLOSE, np.ones((12, 12), np.uint8))

        # Subtract text regions
        for tr in text_regions:
            x1 = max(0, int(tr["x"]) - 8)
            y1 = max(0, int(tr["y"]) - 8)
            x2 = min(w, int(tr["x"] + tr["w"]) + 8)
            y2 = min(h, int(tr["y"] + tr["h"]) + 8)
            content_mask[y1:y2, x1:x2] = 0

        # Subtract shape regions (rectangles)
        for s in shapes:
            if s["type"] == "rectangle":
                rx = max(0, int(s["x"] / sx) - 4)
                ry = max(0, int(s["y"] / sy) - 4)
                rw_ = min(w - rx, int(s["w"] / sx) + 8)
                rh_ = min(h - ry, int(s["h"] / sy) + 8)
                content_mask[ry:ry + rh_, rx:rx + rw_] = 0

        # Also blank out thin horizontal/vertical line areas
        for s in shapes:
            if s["type"] == "line":
                lx1 = int(s["x1"] / sx)
                ly1 = int(s["y1"] / sy)
                lx2 = int(s["x2"] / sx)
                ly2 = int(s["y2"] / sy)
                cv2.line(content_mask, (lx1, ly1), (lx2, ly2), 0, thickness=6)

        # Find connected components
        num_labels, labels, stats, centroids = cv2.connectedComponentsWithStats(content_mask, connectivity=8)
        min_area = w * h * 0.002  # ~0.2% of image

        for i in range(1, num_labels):  # skip background (label 0)
            cx, cy, cw, ch, area = stats[i]
            if area < min_area:
                continue
            aspect = cw / max(ch, 1)
            if aspect > 15 or aspect < 0.07:
                continue
            # Skip very thin regions (likely lines that slipped through)
            if cw < 12 or ch < 12:
                continue

            elements.append({
                "x": round(cx * sx, 4), "y": round(cy * sy, 4),
                "w": round(cw * sx, 4), "h": round(ch * sy, 4),
                "pixel_x": int(cx), "pixel_y": int(cy),
                "pixel_w": int(cw), "pixel_h": int(ch),
                "category": "visual_element",
            })

        # Sort by area (largest first) for better visual hierarchy
        elements.sort(key=lambda e: e["pixel_w"] * e["pixel_h"], reverse=True)

    except Exception as exc:
        logger.warning("Visual element detection: %s", exc)

    return elements[:15]


# ═══════════════════════════════════════════════════════════════════════════
# STEP 4: Clean Background (cv2.inpaint)
# ═══════════════════════════════════════════════════════════════════════════

def _make_clean_background(cv_img, text_regions, visual_elements) -> tuple:
    """Inpaint all foreground content to create clean background."""
    h, w = cv_img.shape[:2]
    mask = np.zeros((h, w), dtype=np.uint8)

    # Mask text regions
    for tr in text_regions:
        px = max(0, int(tr["x"]) - 5)
        py = max(0, int(tr["y"]) - 5)
        pw = min(w - px, int(tr["w"]) + 10)
        ph = min(h - py, int(tr["h"]) + 10)
        mask[py:py + ph, px:px + pw] = 255

    # Mask visual elements
    for ve in visual_elements:
        px, py = ve["pixel_x"], ve["pixel_y"]
        pw, ph = ve["pixel_w"], ve["pixel_h"]
        mask[py:py + ph, px:px + pw] = 255

    if mask.sum() == 0:
        bg = _bg_color(cv_img)
        return (None, "solid", bg)

    kernel = np.ones((5, 5), np.uint8)
    mask = cv2.dilate(mask, kernel, iterations=3)

    inpainted = cv2.inpaint(cv_img, mask, inpaintRadius=7, flags=cv2.INPAINT_TELEA)

    # Classify background
    h2, w2 = inpainted.shape[:2]
    ss = max(8, min(w2, h2) // 30)
    samples = [
        inpainted[0:ss, 0:ss], inpainted[0:ss, -ss:],
        inpainted[-ss:, 0:ss], inpainted[-ss:, -ss:],
        inpainted[h2 // 2 - 4:h2 // 2 + 4, 0:ss],
        inpainted[h2 // 2 - 4:h2 // 2 + 4, -ss:],
    ]
    means = [np.mean(r.reshape(-1, 3), axis=0) for r in samples if r.size > 0]
    if means:
        means = np.array(means)
        if np.max(np.var(means, axis=0)) < 200:
            gm = np.mean(means, axis=0).astype(int)
            return (None, "solid", (int(gm[2]), int(gm[1]), int(gm[0])))

    return (inpainted, "complex", (255, 255, 255))


# ═══════════════════════════════════════════════════════════════════════════
# PPTX Construction Helpers
# ═══════════════════════════════════════════════════════════════════════════

def _place_visual_elements(slide, cv_img, elements, asset_dir) -> list[dict]:
    """Place each visual element as ONE complete image crop."""
    from pptx.util import Inches
    placed = []
    for i, el in enumerate(elements):
        try:
            px, py = el["pixel_x"], el["pixel_y"]
            pw, ph = el["pixel_w"], el["pixel_h"]
            crop = cv_img[py:py + ph, px:px + pw]
            path = os.path.join(asset_dir, f"visual_{i:02d}.png")
            cv2.imwrite(path, crop)
            slide.shapes.add_picture(path, Inches(el["x"]), Inches(el["y"]), Inches(el["w"]), Inches(el["h"]))
            placed.append({**el, "path": path})
        except Exception as exc:
            logger.warning("Visual placement: %s", exc)
    return placed


def _place_shapes(slide, shapes) -> int:
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    count = 0
    for s in shapes:
        try:
            if s["type"] == "line":
                c = slide.shapes.add_connector(1, Inches(s["x1"]), Inches(s["y1"]), Inches(s["x2"]), Inches(s["y2"]))
                c.line.width = Pt(s.get("width_pt", 1.5))
                c.line.color.rgb = RGBColor(0x66, 0x66, 0x66)
                c.name = f"Line_{s.get('subtype','')}_{count + 1}"
                count += 1
            elif s["type"] == "rectangle":
                shp = slide.shapes.add_shape(1, Inches(s["x"]), Inches(s["y"]), Inches(s["w"]), Inches(s["h"]))
                fr, fg, fb = s.get("fill", (240, 240, 240))
                sr, sg, sb = s.get("stroke", (200, 200, 200))
                shp.fill.solid(); shp.fill.fore_color.rgb = RGBColor(fr, fg, fb)
                shp.line.color.rgb = RGBColor(sr, sg, sb); shp.line.width = Pt(0.5)
                shp.name = f"Rect_{count + 1}"
                _fix_ooxml(shp)
                count += 1
        except Exception:
            pass
    return count


def _place_text_boxes(slide, text_regions, iw, ih, sx, sy) -> int:
    """Place transparent editable text boxes with real OCR text."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    count = 0
    for tr in text_regions:
        try:
            x_in = max(0, tr["x"] * sx)
            y_in = max(0, tr["y"] * sy)
            w_in = min(SLIDE_W_IN - x_in, max(tr["w"] * sx, 0.5))
            h_in = min(SLIDE_H_IN - y_in, max(tr["h"] * sy, 0.3))

            ocr_text = tr.get("text", "").strip()
            if not ocr_text:
                continue

            tb = slide.shapes.add_textbox(Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in))
            tf = tb.text_frame; tf.word_wrap = True

            is_title = tr.get("is_title", False)
            fs = tr.get("font_size_est", 14)

            for li, line in enumerate(ocr_text.split("\n")):
                line = line.strip()
                if not line:
                    continue
                p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                p.text = line
                p.font.size = Pt(max(9, min(48, fs)))
                if is_title or fs > 22:
                    p.font.bold = True

            # Dark text for light backgrounds (most common case)
            tc = (26, 26, 26) if is_title else (51, 51, 51)
            for p in tf.paragraphs:
                p.font.color.rgb = RGBColor(*tc)

            rel_x = (x_in + w_in / 2) / SLIDE_W_IN
            p = tf.paragraphs[0]
            if rel_x < 0.25:
                p.alignment = PP_ALIGN.LEFT
            elif rel_x > 0.75:
                p.alignment = PP_ALIGN.RIGHT
            elif 0.42 < rel_x < 0.58:
                p.alignment = PP_ALIGN.CENTER
            else:
                p.alignment = PP_ALIGN.LEFT

            _autofit(tf)
            role = "Title" if is_title else ("Footer" if tr.get("is_footer") else ("Sub" if fs > 20 else "Body"))
            tb.name = f"Text_{role}_{count + 1}"
            count += 1
        except Exception as exc:
            logger.warning("Text box: %s", exc)
    return count


# ═══════════════════════════════════════════════════════════════════════════
# Utilities
# ═══════════════════════════════════════════════════════════════════════════

def _iou(a, b) -> float:
    x1 = max(a["x"], b["x"]); y1 = max(a["y"], b["y"])
    x2 = min(a["x"] + a["w"], b["x"] + b["w"]); y2 = min(a["y"] + a["h"], b["y"] + b["h"])
    if x2 <= x1 or y2 <= y1: return 0.0
    inter = (x2 - x1) * (y2 - y1)
    return inter / max(a["w"] * a["h"] + b["w"] * b["h"] - inter, 1)


def _is_contained(small, big) -> bool:
    """Check if small region is fully contained within big region."""
    return (
        small["x"] >= big["x"] - 5
        and small["y"] >= big["y"] - 5
        and small["x"] + small["w"] <= big["x"] + big["w"] + 5
        and small["y"] + small["h"] <= big["y"] + big["h"] + 5
    )

def _edge_color(cv_img, x, y, w, h):
    hi, wi = cv_img.shape[:2]
    ep = []
    if y > 0: ep.append(cv_img[max(0,y-1):y+2, x:x+w].reshape(-1,3))
    if y+h < hi: ep.append(cv_img[y+h-2:y+h+1, x:x+w].reshape(-1,3))
    if x > 0: ep.append(cv_img[y:y+h, max(0,x-1):x+2].reshape(-1,3))
    if x+w < wi: ep.append(cv_img[y:y+h, x+w-2:x+w+1].reshape(-1,3))
    return tuple(int(c) for c in np.mean(np.concatenate(ep), axis=0)) if ep else (0,0,0)

def _b2r(bgr): return (int(bgr[2]), int(bgr[1]), int(bgr[0])) if isinstance(bgr, np.ndarray) else (bgr[2], bgr[1], bgr[0])

def _bg_color(cv_img):
    h, w = cv_img.shape[:2]
    cs = [cv_img[0:8,0:8], cv_img[0:8,-8:], cv_img[-8:,0:8], cv_img[-8:,-8:]]
    sm = [np.mean(s.reshape(-1,3), axis=0) for s in cs if s.size > 0]
    avg = np.mean(sm, axis=0).astype(int)
    return (int(avg[2]), int(avg[1]), int(avg[0]))

def _set_slide_bg(slide, r, g, b):
    from lxml import etree
    bgPr = slide.background._element.find("{http://schemas.openxmlformats.org/presentationml/2006/main}bgPr")
    if bgPr is None:
        bgPr = etree.SubElement(slide.background._element, qn("p:bgPr"))
    sf = etree.SubElement(bgPr, qn("a:solidFill"))
    sc = etree.SubElement(sf, qn("a:srgbClr")); sc.set("val", f"{r:02x}{g:02x}{b:02x}")

def _autofit(tf):
    from lxml import etree
    bp = tf._element.find(qn("a:bodyPr"))
    if bp is not None:
        for tag in ("noAutofit", "normAutofit", "spAutoFit"):
            for e in bp.findall(qn(f"a:{tag}")): bp.remove(e)
        na = etree.SubElement(bp, qn("a:normAutofit")); na.set("fontScale", "70000")

def _fix_ooxml(shape):
    from lxml import etree
    sp = shape._element.find(qn("p:spPr"))
    if sp is None: return
    children = list(sp)
    gi = next((i for i, c in enumerate(children) if c.tag == qn("a:prstGeom")), None)
    if gi is None: return
    ft = {qn("a:solidFill"), qn("a:gradFill"), qn("a:noFill"), qn("a:pattFill")}
    to_move = [i for i in range(gi+1, len(children)) if children[i].tag in ft]
    for i in reversed(to_move):
        f = children[i]; sp.remove(f); sp.insert(gi, f); gi += 1

def _dedup(shapes):
    if len(shapes) <= 1: return shapes
    shapes.sort(key=lambda s: s.get("w",0)*s.get("h",0) if s.get("type")=="rectangle" else 0, reverse=True)
    kept = []
    for s in shapes:
        dup = False
        for k in kept:
            if s["type"] != k["type"]: continue
            if s["type"] == "line":
                if np.sqrt((s["x1"]-k["x1"])**2+(s["y1"]-k["y1"])**2)<0.15 and np.sqrt((s["x2"]-k["x2"])**2+(s["y2"]-k["y2"])**2)<0.15:
                    dup = True; break
            elif s["type"] == "rectangle":
                if abs(s["x"]+s["w"]/2-k["x"]-k["w"]/2)<0.1 and abs(s["y"]+s["h"]/2-k["y"]-k["h"]/2)<0.1:
                    dup = True; break
        if not dup: kept.append(s)
    return kept
