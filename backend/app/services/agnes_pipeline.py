"""Agnes-enhanced PPTX reconstruction pipeline (Pipeline B v3).

Rich object extraction: every visual element becomes an independently
selectable PowerPoint object. Combines:
  - Side-agreement text erasure (DeckWeaver algorithm, cleans backgrounds)
  - Canny + connected-components → individual visual objects as transparent PNGs
  - OpenCV shape detection → native PPTX shapes (rectangles, circles, lines)
  - RapidOCR + Tesseract → text content + color + size sampling
  - Agnes Vision → layout understanding + color scheme + element classification
  - python-pptx → layered 16:9 slide (bg → shapes → objects → text)
  - PowerPoint COM → SSIM calibration loop

Pipeline v3: OCR → Agnes → erase text → extract objects → build_rich_pptx → calibrate
"""
from __future__ import annotations

import asyncio
import json
import logging
import os
import shutil
import uuid
from collections import Counter
from pathlib import Path
from typing import Any

import cv2
import numpy as np

logger = logging.getLogger(__name__)
PIPELINE_VERSION = "agnes_v3"

from .precise_reconstruction import (
    CoordinateMapper,
    SLIDE_W_IN, SLIDE_H_IN,
    DEFAULT_RENDER_W, DEFAULT_RENDER_H,
    _read_cv_image,
    _augment_text_styles,
    analyze_layers,
    render_pptx_to_png,
    compare_images,
)
from .agnes_layout import (
    analyze_layout as analyze_agnes_layout,
    fuse_agnes_with_ocr,
)
from ._paths import OUTPUTS_DIR

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn


def _hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    hex_color = hex_color.lstrip("#")
    if len(hex_color) == 3:
        hex_color = "".join(c * 2 for c in hex_color)
    return tuple(int(hex_color[i:i + 2], 16) for i in (0, 2, 4))


def _normalize_fused_elements(elements: list[dict[str, Any]]) -> None:
    for elem in elements:
        elem.setdefault("median_word_h", elem.get("h", 16))
        elem.setdefault("median_char_w", elem.get("median_word_h", 16) * 0.55)
        elem.setdefault("color", elem.get("color", "#111111"))
        elem.setdefault("confidence", 80.0)
        elem.setdefault("words", [])
        elem.setdefault("text", elem.get("text", ""))
        elem.setdefault("role", elem.get("agnes_role", "body"))


# ═══════════════════════════════════════════════════════════════════
#  Main entry point
# ═══════════════════════════════════════════════════════════════════

async def reconstruct(
    image_path: str,
    session_id: str = "",
    output_dir: str | None = None,
    agnes_api_key: str | None = None,
    agnes_model: str = "agnes-2.0-flash",
    render_preview: bool = True,
    max_calibration_passes: int = 2,
) -> dict[str, Any]:
    source = Path(image_path)
    if not source.exists():
        return {"error": f"Image not found: {source.name}"}

    out_dir = Path(output_dir) if output_dir else OUTPUTS_DIR
    out_dir = out_dir.resolve()
    out_dir.mkdir(parents=True, exist_ok=True)
    sid = (session_id or uuid.uuid4().hex)[:8]
    run_id = uuid.uuid4().hex[:6]
    stem = f"ppt_agnes_{sid}_{run_id}"
    work_root = out_dir.parent / "precise_work" / stem
    work_root.mkdir(parents=True, exist_ok=True)

    # ── 1. OCR analysis ──
    try:
        layers = await asyncio.to_thread(analyze_layers, str(source))
    except Exception as exc:
        logger.exception("OCR analysis failed")
        return {"error": f"Analysis failed: {exc}"}

    mapper = CoordinateMapper(layers["image_width"], layers["image_height"])
    iw, ih = layers["image_width"], layers["image_height"]
    cv_img = _read_cv_image(str(source))

    # ── 2. Agnes layout understanding ──
    logger.info("Agnes layout analysis...")
    agnes_data = await analyze_agnes_layout(str(source), api_key=agnes_api_key, model=agnes_model)
    agnes_used = not agnes_data.get("_empty", False)

    # ── 3. Fuse Agnes + OCR ──
    editable_text = layers["editable_text"]
    fused = fuse_agnes_with_ocr(agnes_data, editable_text, iw, ih)
    fused_elements = fused["fused_elements"]
    fused_elements.sort(key=lambda e: (e["y"], e["x"]))
    _normalize_fused_elements(fused_elements)
    for idx, elem in enumerate(fused_elements, start=1):
        elem["id"] = f"text_{idx:03d}"

    # ── 4. Augment text styles (OCR color sampling) ──
    _augment_text_styles(cv_img, fused_elements)

    # ── 5. Erase text from background (side-agreement, cleaner than inpaint) ──
    cleaned = await asyncio.to_thread(_erase_text_regions, cv_img, fused_elements)
    clean_path = work_root / "clean_base.png"
    cv2.imencode(".png", cleaned)[1].tofile(str(clean_path))

    # ── 6. Extract individual visual objects ──
    native_shapes = layers.get("native_shapes", [])
    individual_objects = await asyncio.to_thread(
        _extract_individual_objects, cleaned, cv_img, fused_elements, native_shapes, iw, ih, work_root
    )
    logger.info("Extracted %d individual visual objects", len(individual_objects))

    # ── 7. Build PPTX (rich layered reconstruction) ──
    final_pptx = out_dir / f"{stem}.pptx"
    final_preview = out_dir / f"{stem}_preview.png"
    final_comparison = out_dir / f"{stem}_comparison.png"
    final_heatmap = out_dir / f"{stem}_diff.png"
    final_report = out_dir / f"{stem}_report.json"

    scale_pool = [1.0, 0.94, 1.06, 0.90, 1.10, 0.97, 1.03]
    scales = scale_pool[:max(1, max_calibration_passes if render_preview else 1)]
    candidates: list[dict[str, Any]] = []
    build_counts = {}
    for candidate_index, scale in enumerate(scales, start=1):
        candidate_pptx = work_root / f"candidate_{candidate_index}.pptx"
        build_counts = await asyncio.to_thread(
            _build_pptx_rich,
            str(candidate_pptx), mapper, str(clean_path),
            fused_elements, native_shapes, individual_objects,
            agnes_data, scale,
        )
        candidate = {"scale": scale, "pptx": str(candidate_pptx), "score": -1.0, "metrics": {}}
        if render_preview:
            candidate_png = work_root / f"candidate_{candidate_index}.png"
            rendered, render_error = await asyncio.to_thread(
                render_pptx_to_png, str(candidate_pptx), str(candidate_png)
            )
            candidate["render_error"] = render_error
            if rendered:
                metrics = await asyncio.to_thread(
                    compare_images, str(source), str(candidate_png), mapper,
                    fused_elements, None, None,
                )
                candidate["metrics"] = metrics
                candidate["score"] = metrics.get("score", -1.0)
                candidate["preview"] = str(candidate_png)
        candidates.append(candidate)

    # ── 8. Select best candidate ──
    rendered_candidates = [c for c in candidates if c["score"] >= 0]
    best = max(rendered_candidates, key=lambda c: c["score"]) if rendered_candidates else candidates[0]
    shutil.copy2(best["pptx"], final_pptx)

    render_error = best.get("render_error")
    metrics = best.get("metrics", {})
    if best.get("preview") and os.path.exists(best["preview"]):
        shutil.copy2(best["preview"], final_preview)
        try:
            metrics = await asyncio.to_thread(
                compare_images, str(source), str(final_preview), mapper,
                fused_elements, str(final_heatmap), str(final_comparison),
            )
        except Exception:
            pass
    elif render_preview:
        rendered, render_error = await asyncio.to_thread(
            render_pptx_to_png, str(final_pptx), str(final_preview)
        )
        if rendered:
            try:
                metrics = await asyncio.to_thread(
                    compare_images, str(source), str(final_preview), mapper,
                    fused_elements, str(final_heatmap), str(final_comparison),
                )
            except Exception:
                pass

    # ── 9. Report ──
    report = _build_detailed_report(
        source.name, agnes_data, agnes_used, agnes_model,
        fused_elements, build_counts, individual_objects,
        metrics, best, candidates, mapper.mode, render_error,
    )
    final_report.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")

    result = {
        "filename": final_pptx.name,
        "pipeline_version": PIPELINE_VERSION,
        "path": str(final_pptx),
        "url": _safe_url(final_pptx.name),
        "preview_filename": final_preview.name if final_preview.exists() else None,
        "preview_path": str(final_preview) if final_preview.exists() else None,
        "preview_url": _safe_url(final_preview.name) if final_preview.exists() else None,
        "comparison_filename": final_comparison.name if final_comparison.exists() else None,
        "comparison_path": str(final_comparison) if final_comparison.exists() else None,
        "comparison_url": _safe_url(final_comparison.name) if final_comparison.exists() else None,
        "heatmap_filename": final_heatmap.name if final_heatmap.exists() else None,
        "heatmap_url": _safe_url(final_heatmap.name) if final_heatmap.exists() else None,
        "report": report,
        "report_filename": final_report.name,
        "report_url": _safe_url(final_report.name),
        "report_path": str(final_report),
    }
    return result


# ═══════════════════════════════════════════════════════════════════
#  Side-agreement text erasure (DeckWeaver algorithm)
# ═══════════════════════════════════════════════════════════════════

def _erase_text_regions(cv_img: np.ndarray, text_items: list[dict]) -> np.ndarray:
    """Erase text using side-agreement voting for clean background fill."""
    cleaned = cv_img.copy()
    h, w = cv_img.shape[:2]
    for t in text_items:
        x1, y1 = max(0, t["x"] - 2), max(0, t["y"] - 2)
        x2, y2 = min(w, t["x"] + t["w"] + 2), min(h, t["y"] + t["h"] + 2)
        if x2 <= x1 or y2 <= y1:
            continue
        bg_color = _side_agreement_bg(cv_img, x1, y1, x2, y2)
        cleaned[y1:y2, x1:x2] = bg_color
    return cleaned


def _side_agreement_bg(img: np.ndarray, x1: int, y1: int, x2: int, y2: int) -> np.ndarray:
    h, w = img.shape[:2]
    sides = {}
    margin = 3
    if y1 >= margin:
        m = _strip_bg_median(img[y1 - margin:y1, x1:x2])
        if m is not None: sides["top"] = m
    if y2 + margin <= h:
        m = _strip_bg_median(img[y2:y2 + margin, x1:x2])
        if m is not None: sides["bot"] = m
    if x1 >= margin:
        m = _strip_bg_median(img[y1:y2, x1 - margin:x1])
        if m is not None: sides["left"] = m
    if x2 + margin <= w:
        m = _strip_bg_median(img[y1:y2, x2:x2 + margin])
        if m is not None: sides["right"] = m
    if not sides:
        return np.array([255, 255, 255], dtype=np.uint8)
    quantized = [tuple(((c // 20) * 20).astype(int)) for c in sides.values()]
    votes = Counter(quantized)
    winner_q, _ = votes.most_common(1)[0]
    close = [c for c in sides.values() if max(abs(c - np.array(winner_q))) < 30]
    if close:
        return np.median(close, axis=0).astype(np.uint8)
    return np.array(winner_q, dtype=np.uint8)


def _strip_bg_median(strip: np.ndarray) -> np.ndarray | None:
    if strip.size == 0:
        return None
    flat = strip.reshape(-1, 3).astype(np.float32)
    if len(flat) < 3:
        return flat[0].astype(np.uint8)
    median = np.median(flat, axis=0)
    dists = np.linalg.norm(flat - median, axis=1)
    close = flat[dists < np.percentile(dists, 75)]
    if len(close) >= 2:
        return np.median(close, axis=0).astype(np.uint8)
    return median.astype(np.uint8)


# ═══════════════════════════════════════════════════════════════════
#  Rich visual object extraction
# ═══════════════════════════════════════════════════════════════════

def _extract_individual_objects(
    cleaned: np.ndarray,
    original: np.ndarray,
    text_items: list[dict],
    native_shapes: list[dict],
    iw: int, ih: int,
    work_root: Path,
) -> list[dict]:
    """Extract individual visual objects from the clean background.

    Uses Canny edge detection + connected components to find distinct
    visual blobs (icons, logos, decorations).  Each becomes an independent
    transparent PNG that can be selected in PowerPoint.
    """
    h, w = cleaned.shape[:2]

    # Build mask excluding text regions and native shapes
    exclude_mask = np.zeros((h, w), dtype=np.uint8)
    for t in text_items:
        pad = 4
        x1, y1 = max(0, t["x"] - pad), max(0, t["y"] - pad)
        x2, y2 = min(w, t["x"] + t["w"] + pad), min(h, t["y"] + t["h"] + pad)
        exclude_mask[y1:y2, x1:x2] = 255
    for s in native_shapes:
        pad = 6
        sx = s.get("x", s.get("x1", 0))
        sy = s.get("y", s.get("y1", 0))
        sw = s.get("w", abs(s.get("x2", 0) - s.get("x1", 0)))
        sh = s.get("h", abs(s.get("y2", 0) - s.get("y1", 0)))
        x1, y1 = max(0, sx - pad), max(0, sy - pad)
        x2, y2 = min(w, sx + sw + pad), min(h, sy + sh + pad)
        exclude_mask[y1:y2, x1:x2] = 255

    gray = cv2.cvtColor(cleaned, cv2.COLOR_BGR2GRAY)
    edges = cv2.Canny(gray, 40, 130)
    edges[exclude_mask > 0] = 0
    kernel = np.ones((3, 3), np.uint8)
    closed = cv2.morphologyEx(edges, cv2.MORPH_CLOSE, kernel)
    kernel2 = np.ones((5, 5), np.uint8)
    closed = cv2.dilate(closed, kernel2, iterations=1)

    num_labels, labels, stats, centroids = cv2.connectedComponentsWithStats(closed, 8)
    min_area = max(400, int(w * h * 0.0008))
    max_area = int(w * h * 0.45)

    objects_dir = work_root / "objects"
    objects_dir.mkdir(exist_ok=True)
    objects = []

    for i in range(1, num_labels):
        cx, cy, cw, ch, area = [int(v) for v in stats[i]]
        if area < min_area or area > max_area:
            continue
        aspect = cw / max(ch, 1)
        if aspect > 20 or aspect < 0.05:
            continue
        # Expand bounding box slightly for context
        pad_x = max(2, cw // 20)
        pad_y = max(2, ch // 20)
        x1 = max(0, cx - pad_x)
        y1 = max(0, cy - pad_y)
        x2 = min(w, cx + cw + pad_x)
        y2 = min(h, cy + ch + pad_y)
        # Create transparent PNG: original pixels inside contour, alpha elsewhere
        obj_mask = np.zeros((h, w), dtype=np.uint8)
        obj_mask[labels == i] = 255
        obj_mask = cv2.dilate(obj_mask, np.ones((3, 3), np.uint8), iterations=1)
        crop = original[y1:y2, x1:x2].copy()
        alpha = obj_mask[y1:y2, x1:x2]
        rgba = np.dstack([crop, alpha])
        obj_path = objects_dir / f"obj_{i:03d}.png"
        cv2.imencode(".png", rgba)[1].tofile(str(obj_path))
        objects.append({
            "x": x1, "y": y1, "w": x2 - x1, "h": y2 - y1,
            "path": str(obj_path), "area": area,
        })

    return sorted(objects, key=lambda o: o["area"], reverse=True)[:50]


# ═══════════════════════════════════════════════════════════════════
#  Rich PPTX Builder
# ═══════════════════════════════════════════════════════════════════

def _build_pptx_rich(
    pptx_path: str,
    mapper: CoordinateMapper,
    clean_base_path: str,
    text_elements: list[dict],
    native_shapes: list[dict],
    visual_objects: list[dict],
    agnes_data: dict,
    font_scale: float,
) -> dict[str, int]:
    """Build PPTX with every element as an independent PowerPoint object.

    Layer 0: Clean visual base (full-slide background)
    Layer 1: Native shapes (rectangles, circles, lines, arrows)
    Layer 2: Individual visual objects (icons, logos, decorations)
    Layer 3: Editable text boxes
    """
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    counts = {"editable_text": 0, "editable_shapes": 0, "visual_objects": 0}

    # ── Layer 0: Clean visual base ──
    if os.path.exists(clean_base_path):
        slide.shapes.add_picture(clean_base_path, Inches(0), Inches(0),
                                 Inches(SLIDE_W_IN), Inches(SLIDE_H_IN))

    color_scheme = agnes_data.get("color_scheme", {})

    # ── Layer 1: Native shapes ──
    for idx, s in enumerate(native_shapes):
        try:
            shape_type = s.get("type", "rect")
            if shape_type == "line":
                x1, y1 = mapper.x(s["x1"]), mapper.y(s["y1"])
                x2, y2 = mapper.x(s["x2"]), mapper.y(s["y2"])
                shp = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2),
                )
                shp.line.color.rgb = RGBColor(*_hex_to_rgb(s.get("stroke", "#A8BFC2")))
                shp.line.width = Pt(s.get("width_pt", 0.8))
            elif shape_type == "ellipse":
                x, y, w, h = mapper.rect(s["x"], s["y"], s["w"], s["h"])
                shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
                shp.fill.solid()
                shp.fill.fore_color.rgb = RGBColor(*_hex_to_rgb(s.get("fill", "#08766F")))
                shp.line.color.rgb = RGBColor(*_hex_to_rgb(s.get("stroke", "#A8BFC2")))
                shp.line.width = Pt(0.5)
            elif shape_type in ("rect", "round_rect"):
                x, y, w, h = mapper.rect(s["x"], s["y"], s["w"], s["h"])
                mso = MSO_SHAPE.ROUNDED_RECTANGLE if shape_type == "round_rect" else MSO_SHAPE.RECTANGLE
                shp = slide.shapes.add_shape(mso, Inches(x), Inches(y), Inches(w), Inches(h))
                fill_hex = s.get("fill", "#FFFFFF")
                shp.fill.solid()
                shp.fill.fore_color.rgb = RGBColor(*_hex_to_rgb(fill_hex))
                shp.line.color.rgb = RGBColor(*_hex_to_rgb(s.get("stroke", "#A8BFC2")))
                shp.line.width = Pt(s.get("width_pt", 0.55))
                try:
                    if shape_type == "round_rect" and shp.adjustments:
                        shp.adjustments[0] = float(s.get("radius_ratio", 0.18))
                except Exception:
                    pass
            else:
                continue
            shp.name = f"Shape_{idx:02d}_{shape_type}"
            counts["editable_shapes"] += 1
        except Exception as exc:
            logger.debug("Shape %s skip: %s", idx, exc)

    # ── Layer 2: Individual visual objects ──
    for idx, obj in enumerate(visual_objects):
        try:
            x, y, w, h = mapper.rect(obj["x"], obj["y"], obj["w"], obj["h"])
            if os.path.exists(obj["path"]):
                pic = slide.shapes.add_picture(obj["path"], Inches(x), Inches(y), Inches(w), Inches(h))
                pic.name = f"VisualObj_{idx:02d}"
                counts["visual_objects"] += 1
        except Exception as exc:
            logger.debug("Object %s skip: %s", idx, exc)

    # ── Layer 3: Editable text boxes ──
    text_roles = {
        "title", "subtitle", "body", "kpi-number", "kpi-label",
        "bullet-list", "page-number", "footer", "tag", "callout-box",
    }
    for elem in text_elements:
        agnes_type = elem.get("agnes_type", "body")
        if agnes_type not in text_roles:
            continue
        text = elem.get("text", "").strip()
        if not text:
            continue

        rect = mapper.rect(elem["x"], elem["y"], elem["w"], elem["h"])
        pixel_h = elem.get("median_word_h", 16)
        ocr_pt = max(8, pixel_h * 0.75 * font_scale)
        role = elem.get("agnes_role", "body")
        role_map = {"title": 1.0, "heading": 0.85, "body": 0.7, "number": 0.9}
        font_size = max(8, ocr_pt * role_map.get(role, 0.7))
        text_color = elem.get("color", "#111111")
        font_style = elem.get("agnes_font_style", {})
        align_str = font_style.get("alignment", "left")
        align_map = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
        weight = font_style.get("weight", "regular")

        text = elem.get("text", "")
        has_cjk = any('一' <= c <= '鿿' for c in text)
        font_name = "Microsoft YaHei" if has_cjk else "Arial"

        # Expand box slightly for PowerPoint metrics
        box_y = max(0, rect[1] - min(rect[3] * 0.15, 0.06))
        box_h = min(SLIDE_H_IN - box_y, rect[3] * 1.35)
        box_w = min(SLIDE_W_IN - rect[0], rect[2] * 1.04)

        box = slide.shapes.add_textbox(
            Inches(rect[0]), Inches(box_y), Inches(box_w), Inches(box_h)
        )
        box.name = f"Text_{elem['id']}_{role}"
        tf = box.text_frame
        tf.word_wrap = False
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0

        p = tf.paragraphs[0]
        p.alignment = align_map.get(align_str, PP_ALIGN.LEFT)
        p.text = text
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = Pt(font_size)
        run.font.name = font_name
        r, g, b = _hex_to_rgb(text_color)
        run.font.color.rgb = RGBColor(r, g, b)
        run.font.bold = weight in ("bold", "semibold") or role in ("title", "number")
        try:
            ea = run._r.get_or_add_rPr().makeelement(qn("a:ea"), {"typeface": font_name if has_cjk else "Microsoft YaHei"})
            run._r.get_or_add_rPr().append(ea)
        except Exception:
            pass
        counts["editable_text"] += 1

    prs.save(pptx_path)
    return counts


# ═══════════════════════════════════════════════════════════════════
#  Report
# ═══════════════════════════════════════════════════════════════════

def _build_detailed_report(
    source_name: str,
    agnes_data: dict, agnes_used: bool, agnes_model: str,
    fused_elements: list, build_counts: dict,
    visual_objects: list,
    metrics: dict, best: dict, candidates: list,
    coord_mode: str, render_error: str | None,
) -> dict[str, Any]:
    editable_text = []
    editable_shapes = []
    image_objects = []
    raster_embedded = []

    for elem in fused_elements:
        agnes_type = elem.get("agnes_type", "body")
        text = elem.get("text", "").strip()
        is_text = agnes_type in {
            "title", "subtitle", "body", "kpi-number", "kpi-label",
            "bullet-list", "page-number", "footer", "tag", "callout-box",
        }
        is_visual = agnes_type in {"icon", "logo", "photo", "illustration", "chart"}

        if is_visual:
            image_objects.append({"label": elem.get("id", "?"), "type": agnes_type,
                                  "reason": "独立图标/图片素材"})
        elif is_text and text:
            editable_text.append({"label": elem.get("id", "?"), "type": agnes_type,
                                  "text": text[:50], "color": elem.get("color", "#111"),
                                  "font_size_pt": round(elem.get("median_word_h", 16) * 0.75, 1)})
        elif agnes_type in ("card", "circle-badge", "divider-line", "arrow", "button",
                            "rect", "rounded_rect", "line", "ellipse"):
            editable_shapes.append({"label": elem.get("id", "?"), "type": agnes_type,
                                    "shape": elem.get("agnes_shape", {}).get("shape_type", "rect")})
        elif text:
            raster_embedded.append({"label": elem.get("id", "?"), "type": agnes_type,
                                     "text": text[:30], "reason": "嵌入复杂视觉区域"})

    return {
        "source": source_name,
        "pipeline_version": PIPELINE_VERSION,
        "slide_size": "16:9",
        "coordinate_mode": coord_mode,
        "agnes": {
            "used": agnes_used, "model": agnes_model,
            "layout_type": agnes_data.get("page_structure", {}).get("layout_type", "?"),
            "visual_style": agnes_data.get("page_structure", {}).get("visual_style", "?"),
            "element_count": len(agnes_data.get("elements", [])),
            "group_count": len(agnes_data.get("groups", [])),
        },
        "elements_summary": {
            "editable_text_boxes": build_counts.get("editable_text", 0),
            "editable_native_shapes": build_counts.get("editable_shapes", 0),
            "individual_visual_objects": build_counts.get("visual_objects", 0),
            "raster_embedded_text": len(raster_embedded),
        },
        "editable_text": editable_text,
        "editable_shapes": editable_shapes,
        "individual_objects": [
            {"label": f"obj_{i:03d}", "area_px": obj["area"], "size": f"{obj['w']}x{obj['h']}"}
            for i, obj in enumerate(visual_objects[:20])
        ],
        "raster_embedded_text": raster_embedded,
        "color_scheme": agnes_data.get("color_scheme", {}),
        "visual_qa": {**metrics,
            "font_scale_candidates": [{"scale": c["scale"], "score": c["score"]} for c in candidates],
            "selected_font_scale": best["scale"], "render_error": render_error},
    }


def _safe_url(filename: str) -> str:
    return f"/api/skills/download/{filename}"
